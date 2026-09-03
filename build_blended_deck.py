import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = "/Users/subham/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx"
OUTPUT_PATH = "/Users/subham/Downloads/ORCA_4.0_SIH2026_Official_Deck.pptx"

# ==============================================================================
# UNIFIED COLOR PALETTE - MODERN MARITIME MISSION CONTROL
# ==============================================================================
# Backgrounds & Cards
CLR_CARD_BG        = RGBColor(255, 255, 255)       # Clean crisp white
CLR_CARD_BG_MUTED  = RGBColor(241, 245, 249)       # Slate 100
CLR_CARD_BG_OCEAN  = RGBColor(240, 249, 255)       # Sky 50
CLR_CARD_BG_DARK   = RGBColor(10, 25, 47)          # Deep Navy 950
CLR_CARD_BG_NAVY   = RGBColor(15, 34, 64)          # Navy 900
CLR_HERO_BG        = RGBColor(14, 30, 54)          # Deep Maritime Hero

# Brand Primary & Accents
CLR_NAVY_TITLE     = RGBColor(13, 71, 161)         # Official SIH Navy Deep (#0d47a1)
CLR_OCEAN_BLUE     = RGBColor(2, 132, 199)         # Sky/Ocean Blue (#0284c7)
CLR_CYAN_ACCENT    = RGBColor(6, 182, 212)         # Vibrant Cyan (#06b6d4)
CLR_EMERALD        = RGBColor(16, 185, 129)        # Emerald Safe (#10b981)
CLR_AMBER          = RGBColor(217, 119, 6)         # Warning Amber (#d97706)
CLR_ROSE           = RGBColor(225, 29, 72)         # Danger Rose (#e11d48)
CLR_INDIGO         = RGBColor(67, 56, 202)         # Indigo (#4338ca)

# Text Colors
CLR_TEXT_MAIN      = RGBColor(15, 23, 42)          # Slate 900
CLR_TEXT_BODY      = RGBColor(51, 65, 85)          # Slate 700
CLR_TEXT_MUTED     = RGBColor(100, 116, 139)       # Slate 500
CLR_TEXT_WHITE     = RGBColor(255, 255, 255)       # White
CLR_TEXT_CYAN      = RGBColor(56, 189, 248)        # Sky 400

# Borders
CLR_BORDER_LIGHT   = RGBColor(226, 232, 240)       # Slate 200
CLR_BORDER_BLUE    = RGBColor(186, 230, 253)       # Sky 200
CLR_BORDER_PRIMARY = RGBColor(13, 71, 161)         # Primary Dark Blue
CLR_BORDER_CYAN    = RGBColor(6, 182, 212)         # Cyan Border

def set_shape_style(shape, fill_color, border_color=None, border_width=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

def add_card(slide, left, top, width, height, bg_color=CLR_CARD_BG, border_color=CLR_BORDER_LIGHT, border_width=1.0, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_style(shape, bg_color, border_color, border_width)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return shape, tf

def setup_slide_header(slide, title_text, badges):
    for s in slide.shapes:
        if s.name == "Title 1":
            s.left = Inches(1.85)
            s.top = Inches(0.12)
            s.width = Inches(8.75)
            s.height = Inches(0.85)
            tf = s.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = "Arial"
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = CLR_NAVY_TITLE
            p.alignment = PP_ALIGN.LEFT
        elif "Oval" in s.name:
            s.left = Inches(0.36)
            s.top = Inches(0.18)
            s.width = Inches(1.35)
            s.height = Inches(0.75)
            tf = s.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "ORCA 4.0"
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = CLR_TEXT_WHITE
            p.alignment = PP_ALIGN.CENTER
            s.fill.solid()
            s.fill.fore_color.rgb = CLR_NAVY_TITLE
            s.line.fill.background()
        elif "TextBox" in s.name:
            s.text_frame.clear()

    # Add Official Idea Pointer Badges Bar across top (top: 1.02 to 1.34 in)
    badge_w = Inches(2.95)
    badge_h = Inches(0.30)
    for idx, (b_text, b_col) in enumerate(badges):
        bx = Inches(0.6 + idx * 3.06)
        by = Inches(1.02)
        b_shape, b_tf = add_card(slide, bx, by, badge_w, badge_h, bg_color=b_col, border_color=None)
        p = b_tf.paragraphs[0]
        p.text = b_text
        p.font.name = "Arial"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = CLR_TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

def build_presentation():
    print(f"Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)

    # Delete Slide 7 if present
    if len(prs.slides) >= 7:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        print("Deleted instruction slide. Exactly 6 slides remain.")

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides[0]
    print("Formatting Slide 1: TITLE PAGE...")

    for shape in slide1.shapes:
        if shape.name == "TextBox 9":
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fields = [
                ("Problem Statement ID:", " SIH26176"),
                ("Problem Statement Title:", " AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform for Artisanal Fishermen & Coastal Fleets"),
                ("Theme:", " Space Technology / Blue Economy / Disaster Management"),
                ("PS Category:", " Software (with Edge Hardware compatibility)"),
                ("Team ID:", " [Team ID]"),
                ("Team Name:", " ORCA 4.0")
            ]
            
            for idx, (label, val) in enumerate(fields):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.space_after = Pt(7)
                
                r1 = p.add_run()
                r1.text = label
                r1.font.name = "Arial"
                r1.font.bold = True
                r1.font.size = Pt(12)
                r1.font.color.rgb = CLR_NAVY_TITLE
                
                r2 = p.add_run()
                r2.text = val
                r2.font.name = "Arial"
                r2.font.bold = False
                r2.font.size = Pt(12)
                r2.font.color.rgb = CLR_TEXT_MAIN
        elif shape.name == "Subtitle 3":
            shape.text_frame.text = "TITLE PAGE"

    # Right Hero Container on Slide 1
    _, tf_h1 = add_card(slide1, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5),
                        bg_color=CLR_HERO_BG, border_color=CLR_CYAN_ACCENT, border_width=1.5)
    
    p0 = tf_h1.paragraphs[0]
    p0.text = "SMART INDIA HACKATHON 2026 · PROJECT SHOWCASE"
    p0.font.name = "Arial"
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = CLR_CYAN_ACCENT
    p0.space_after = Pt(4)

    p1 = tf_h1.add_paragraph()
    p1.text = "ORCA 4.0"
    p1.font.name = "Arial"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = CLR_TEXT_WHITE
    p1.space_after = Pt(2)

    p2 = tf_h1.add_paragraph()
    p2.text = "“FROM OCEAN DATA TO SAFER DECISIONS”"
    p2.font.name = "Arial"
    p2.font.size = Pt(13.5)
    p2.font.bold = True
    p2.font.color.rgb = CLR_CYAN_ACCENT
    p2.space_after = Pt(6)

    p3 = tf_h1.add_paragraph()
    p3.text = "AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform"
    p3.font.name = "Arial"
    p3.font.size = Pt(10.5)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(10)

    p4 = tf_h1.add_paragraph()
    p4.text = "CORE DECISION CHAIN:"
    p4.font.name = "Arial"
    p4.font.size = Pt(9.5)
    p4.font.bold = True
    p4.font.color.rgb = CLR_EMERALD
    
    p5 = tf_h1.add_paragraph()
    p5.text = "OCEAN  ➔  WEATHER  ➔  VESSEL  ➔  ROUTE  ➔  ORCA  ➔  RISK  ➔  DECISION"
    p5.font.name = "Arial"
    p5.font.size = Pt(8.5)
    p5.font.bold = True
    p5.font.color.rgb = CLR_TEXT_WHITE
    p5.space_after = Pt(12)

    s1_cards = [
        ("🛡️ Deterministic Safety Shield", "Hard physics capsizing check (Hcrit = 0.6·L) with zero AI hallucination risk", CLR_EMERALD),
        ("🚢 Parametric Vessel Twin", "Dynamic wave encounter physics tailored to vessel hull length, draft & beam", CLR_CYAN_ACCENT),
        ("📻 Multi-Hop LoRa Fleet Mesh", "Overcomes 12 NM 4G blackout for zero-cost offshore telemetry up to 50 km", CLR_AMBER)
    ]
    for c_title, c_desc, c_col in s1_cards:
        p_c = tf_h1.add_paragraph()
        p_c.space_after = Pt(5)
        r_t = p_c.add_run()
        r_t.text = f"{c_title}: "
        r_t.font.name = "Arial"
        r_t.font.bold = True
        r_t.font.size = Pt(9.5)
        r_t.font.color.rgb = c_col
        r_d = p_c.add_run()
        r_d.text = c_desc
        r_d.font.name = "Arial"
        r_d.font.size = Pt(8.5)
        r_d.font.color.rgb = RGBColor(203, 213, 225)

    # =========================================================================
    # SLIDE 2: IDEA TITLE / PROPOSED SOLUTION
    # =========================================================================
    slide2 = prs.slides[1]
    print("Formatting Slide 2: IDEA TITLE / PROPOSED SOLUTION...")
    setup_slide_header(slide2, "IDEA TITLE: ORCA 4.0 — Maritime Decision Intelligence",
                       [("PROPOSED SOLUTION", CLR_OCEAN_BLUE),
                        ("DETAILED EXPLANATION", CLR_NAVY_TITLE),
                        ("HOW IT ADDRESSES PROBLEM", CLR_NAVY_TITLE),
                        ("INNOVATION & UNIQUENESS", CLR_EMERALD)])

    # Top Thesis Callout Box (top: 1.38, height: 0.52)
    _, tf_th = add_card(slide2, Inches(0.6), Inches(1.38), Inches(12.13), Inches(0.52),
                        bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_th = tf_th.paragraphs[0]
    r1 = p_th.add_run()
    r1.text = "CORE PARADIGM SHIFT:  FORECAST ≠ OPERATIONAL RISK.  "
    r1.font.name = "Arial"
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.color.rgb = CLR_AMBER
    r2 = p_th.add_run()
    r2.text = "Existing systems ask 'What is happening at sea?'. ORCA computes 'What does it mean for THIS vessel, on THIS route, at THIS time?'"
    r2.font.name = "Arial"
    r2.font.size = Pt(9.5)
    r2.font.color.rgb = CLR_TEXT_WHITE

    # Left Container: Data Transformation Pipeline (top: 1.95, width: 7.3, height: 4.25)
    _, tf_t_hdr = add_card(slide2, Inches(0.6), Inches(1.95), Inches(7.3), Inches(4.25),
                           bg_color=CLR_CARD_BG_MUTED, border_color=CLR_BORDER_BLUE, border_width=1.0)
    p_thd = tf_t_hdr.paragraphs[0]
    p_thd.text = "VISUAL TRANSFORMATION: FROM DATA CHAOS TO OPERATIONAL DECISION"
    p_thd.font.name = "Arial"
    p_thd.font.bold = True
    p_thd.font.size = Pt(10)
    p_thd.font.color.rgb = CLR_NAVY_TITLE

    # 8 Input Nodes in 2 columns
    inputs_1 = [
        ("🌊 WAVES (Hs, Swell Tp)", "Sea & swell wave spectra"),
        ("💨 WIND (Vel, Gusts)", "Surface boundary wind vectors"),
        ("🌀 CYCLONES (IMD Alerts)", "Storm center coordinates"),
        ("🌊 CURRENTS (Surface u,v)", "Ocean drift velocity vectors")
    ]
    inputs_2 = [
        ("🛰️ SATELLITES (SST / Chlor)", "Thermal front boundaries"),
        ("🚢 VESSEL (Lwl, Draft, Beam)", "Displacement hull profile"),
        ("📡 AIS FEEDS (MMSI Data)", "Real-time vessel position"),
        ("🗺️ ROUTE (Planned Waypoints)", "Departure & target docks")
    ]

    for idx, (title, sub) in enumerate(inputs_1):
        _, tf_in = add_card(slide2, Inches(0.72), Inches(2.32 + idx * 0.90), Inches(2.0), Inches(0.80),
                            bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
        p = tf_in.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = CLR_CYAN_ACCENT
        p2 = tf_in.add_paragraph()
        p2.text = sub
        p2.font.name = "Arial"
        p2.font.size = Pt(7)
        p2.font.color.rgb = RGBColor(203, 213, 225)

    for idx, (title, sub) in enumerate(inputs_2):
        _, tf_in = add_card(slide2, Inches(2.80), Inches(2.32 + idx * 0.90), Inches(2.0), Inches(0.80),
                            bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
        p = tf_in.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = CLR_CYAN_ACCENT
        p2 = tf_in.add_paragraph()
        p2.text = sub
        p2.font.name = "Arial"
        p2.font.size = Pt(7)
        p2.font.color.rgb = RGBColor(203, 213, 225)

    # Central Engine Box
    _, tf_ce = add_card(slide2, Inches(4.90), Inches(2.32), Inches(2.85), Inches(1.65),
                        bg_color=CLR_OCEAN_BLUE, border_color=CLR_CYAN_ACCENT, border_width=1.5)
    p_ce = tf_ce.paragraphs[0]
    p_ce.text = "⚡ ORCA 4.0 ENGINE"
    p_ce.font.name = "Arial"
    p_ce.font.bold = True
    p_ce.font.size = Pt(11)
    p_ce.font.color.rgb = CLR_TEXT_WHITE
    p_ce.alignment = PP_ALIGN.CENTER
    p_ce2 = tf_ce.add_paragraph()
    p_ce2.text = "• 0.083° Spatial Grid Normalization\n• Parametric Vessel Digital Twin\n• Multi-Species Habitat Index"
    p_ce2.font.name = "Arial"
    p_ce2.font.size = Pt(8)
    p_ce2.font.color.rgb = CLR_TEXT_WHITE

    # Downstream Decision Boxes
    out_cards = [
        ("🛡️ DETERMINISTIC RISK", "Hcrit capsizing check (<10ms)", CLR_EMERALD),
        ("🗺️ ROUTE RISK FIELD", "Dynamic 2D risk surface A*", CLR_CYAN_ACCENT),
        ("📢 ACTIONABLE VERDICT", "Clear Go/No-Go + Audio Voice", CLR_AMBER)
    ]
    for idx, (title, sub, col) in enumerate(out_cards):
        _, tf_oc = add_card(slide2, Inches(4.90), Inches(4.08 + idx * 0.70), Inches(2.85), Inches(0.62),
                            bg_color=CLR_CARD_BG_DARK, border_color=col, border_width=1.0)
        p = tf_oc.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = col
        r2 = p.add_run()
        r2.text = sub
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_WHITE

    # Right Box: Competitive Differentiation Matrix
    _, tf_diff = add_card(slide2, Inches(8.05), Inches(1.95), Inches(4.68), Inches(4.25),
                          bg_color=CLR_CARD_BG, border_color=CLR_NAVY_TITLE, border_width=1.0)
    p_dfh = tf_diff.paragraphs[0]
    p_dfh.text = "COMPETITIVE DIFFERENTIATION MATRIX"
    p_dfh.font.name = "Arial"
    p_dfh.font.bold = True
    p_dfh.font.size = Pt(10)
    p_dfh.font.color.rgb = CLR_NAVY_TITLE
    p_dfh.space_after = Pt(4)

    comp_rows = [
        ("Capability", "Existing Systems", "ORCA 4.0"),
        ("Weather Forecast", "● Regional Only", "✓ High-Res (0.083°)"),
        ("Cyclone Warnings", "● Generic SMS", "✓ Dynamic Geofenced"),
        ("Vessel Hull Context", "— Ignored", "✓ Vessel Twin (Lwl, Draft)"),
        ("Wave Physics Math", "— None", "✓ Hcrit & Resonance Math"),
        ("Safe Route Pathfinder", "— Static Path", "✓ Dynamic A* Avoidance"),
        ("Deterministic Safety", "— Unverified", "✓ Non-bypassable Rule"),
        ("Offshore Reach", "— Dead >12 NM", "✓ LoRa Mesh (50 km)")
    ]
    for idx, (c1, c2, c3) in enumerate(comp_rows):
        p_r = tf_diff.add_paragraph()
        p_r.space_after = Pt(2)
        r1 = p_r.add_run()
        r1.text = f"{c1.ljust(20)}: "
        r1.font.name = "Arial"
        r1.font.bold = True if idx == 0 else False
        r1.font.size = Pt(8)
        r1.font.color.rgb = CLR_NAVY_TITLE if idx == 0 else CLR_TEXT_MAIN
        
        r2 = p_r.add_run()
        r2.text = f"{c2}  ➔  "
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_ROSE if "—" in c2 else CLR_TEXT_MUTED
        
        r3 = p_r.add_run()
        r3.text = c3
        r3.font.name = "Arial"
        r3.font.bold = True
        r3.font.size = Pt(8)
        r3.font.color.rgb = CLR_EMERALD if "✓" in c3 else CLR_TEXT_MAIN

    # Bottom Innovation Strip
    _, tf_btm2 = add_card(slide2, Inches(0.6), Inches(6.25), Inches(12.13), Inches(0.45),
                          bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_b2 = tf_btm2.paragraphs[0]
    p_b2.text = "① MULTI-SOURCE HARMONIZATION   ➔   ② VESSEL-AWARE PHYSICS TWIN   ➔   ③ ROUTE RISK OPTIMIZATION   ➔   ④ EXPLAINABLE SAFETY SHIELD"
    p_b2.font.name = "Arial"
    p_b2.font.bold = True
    p_b2.font.size = Pt(9.5)
    p_b2.font.color.rgb = CLR_CYAN_ACCENT
    p_b2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides[2]
    print("Formatting Slide 3: TECHNICAL APPROACH...")
    setup_slide_header(slide3, "TECHNICAL APPROACH: End-to-End Architecture",
                       [("TECHNOLOGIES TO BE USED", CLR_OCEAN_BLUE),
                        ("METHODOLOGY & PROCESS", CLR_NAVY_TITLE),
                        ("SYSTEM FLOWCHARTS", CLR_NAVY_TITLE),
                        ("PROTOTYPE SIMULATION", CLR_EMERALD)])

    # Top Architecture Box (top: 1.38, height: 3.25)
    _, tf_a_bg = add_card(slide3, Inches(0.6), Inches(1.38), Inches(12.13), Inches(3.25),
                          bg_color=CLR_CARD_BG_DARK, border_color=CLR_BORDER_BLUE, border_width=1.0)

    # 6 Ingestion Nodes
    sources = [
        ("INCOIS ERDDAP", "Wave spectrum & currents"),
        ("IMD Alerts", "Real-time cyclone tracks"),
        ("Copernicus Marine", "0.083° global ocean grids"),
        ("ISRO MOSDAC", "INSAT-3DR SST & Oceansat-3"),
        ("AIS Streams", "Vessel transponder positions"),
        ("NMEA Sensors", "Onboard GPS / IMU / Pitch")
    ]
    for idx, (s_name, s_sub) in enumerate(sources):
        _, tf_s = add_card(slide3, Inches(0.72 + idx * 1.96), Inches(1.48), Inches(1.88), Inches(0.62),
                           bg_color=CLR_CARD_BG_NAVY, border_color=CLR_CYAN_ACCENT, border_width=1.0)
        p = tf_s.paragraphs[0]
        p.text = f"📡 {s_name}"
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8)
        p.font.color.rgb = CLR_CYAN_ACCENT
        p2 = tf_s.add_paragraph()
        p2.text = s_sub
        p2.font.name = "Arial"
        p2.font.size = Pt(7)
        p2.font.color.rgb = RGBColor(203, 213, 225)

    # Normalization Funnel Bar
    _, tf_fn = add_card(slide3, Inches(0.72), Inches(2.18), Inches(11.89), Inches(0.42),
                        bg_color=CLR_OCEAN_BLUE, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = "⚡ DATA ENGINE & DUAL DIGITAL STATE: Spatial H3 Indexing (0.083°) + Environmental State  ⟷  Parametric Vessel Digital Twin (Lwl, Draft, Beam)"
    p_fn.font.name = "Arial"
    p_fn.font.bold = True
    p_fn.font.size = Pt(8.5)
    p_fn.font.color.rgb = CLR_TEXT_WHITE
    p_fn.alignment = PP_ALIGN.CENTER

    # Tripartite Engines
    engines = [
        ("⚙️ DETERMINISTIC PHYSICS (<10ms)", 
         "• Critical Wave Height (Hcrit = 0.6 · Lhull)\n• Wave Steepness (S = Hs / λ) & Period\n• Beam-Sea Encounter Angle & Roll Resonance",
         CLR_EMERALD),
        ("🧠 PREDICTIVE ML & ANALYTICS", 
         "• Multi-Species Habitat Index (HSI XGBoost)\n• 2D Sobel Filter SST Thermal Fronts\n• 1,000-Particle Monte Carlo SAR Drift",
         CLR_CYAN_ACCENT),
        ("📍 DYNAMIC A* OPTIMIZATION", 
         "• 2D Dynamic Risk Cost Surface Mesh\n• A* Navigational Pathfinding & Avoidance\n• Optimal Fuel Burn & Dock Allocation",
         CLR_AMBER)
    ]
    for idx, (e_title, e_desc, e_col) in enumerate(engines):
        _, tf_e = add_card(slide3, Inches(0.72 + idx * 4.02), Inches(2.68), Inches(3.85), Inches(1.22),
                           bg_color=CLR_CARD_BG_NAVY, border_color=e_col, border_width=1.0)
        p = tf_e.paragraphs[0]
        p.text = e_title
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = e_col
        p2 = tf_e.add_paragraph()
        p2.text = e_desc
        p2.font.name = "Arial"
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = CLR_TEXT_WHITE

    # Hard Safety Circuit Breaker Banner
    _, tf_scb = add_card(slide3, Inches(0.72), Inches(3.98), Inches(11.89), Inches(0.52),
                         bg_color=CLR_CARD_BG_NAVY, border_color=CLR_EMERALD, border_width=1.0)
    p_scb = tf_scb.paragraphs[0]
    r1 = p_scb.add_run()
    r1.text = "🛡️ NON-BYPASSABLE SAFETY CIRCUIT BREAKER: "
    r1.font.name = "Arial"
    r1.font.bold = True
    r1.font.size = Pt(9)
    r1.font.color.rgb = CLR_EMERALD
    r2 = p_scb.add_run()
    r2.text = "Pure deterministic compiled logic evaluates physical capsizing risk (<10ms). LLM is strictly constrained downstream for vernacular audio synthesis only."
    r2.font.name = "Arial"
    r2.font.size = Pt(8)
    r2.font.color.rgb = CLR_TEXT_WHITE

    # Bottom Left: Live Prototype Simulation
    _, tf_dm = add_card(slide3, Inches(0.6), Inches(4.72), Inches(5.9), Inches(1.98),
                        bg_color=CLR_CARD_BG_OCEAN, border_color=CLR_NAVY_TITLE, border_width=1.0)
    p_dm_h = tf_dm.paragraphs[0]
    p_dm_h.text = "LIVE WORKING PROTOTYPE SIMULATION (KANYAKUMARI CASE)"
    p_dm_h.font.name = "Arial"
    p_dm_h.font.bold = True
    p_dm_h.font.size = Pt(9.5)
    p_dm_h.font.color.rgb = CLR_NAVY_TITLE
    p_dm_h.space_after = Pt(3)

    k_steps = [
        ("1. Input", "Kanyakumari Port (8.08°N, 77.55°E), 9m FRP Craft"),
        ("2. Live State", "Wave Hs = 2.8m, Swell Tp = 8.2s, Wind = 24 kts"),
        ("3. Physics Check", "Hcrit = 5.4m > 2.8m (Safe from Immediate Capsize)"),
        ("4. Route Engine", "Diverts 3.2 NM East to avoid shallow reef breaker"),
        ("5. Output Verdict", "PROCEED WITH CAUTION (Audio Advisory Generated)")
    ]
    for sn, sv in k_steps:
        p = tf_dm.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"{sn}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = CLR_NAVY_TITLE
        r2 = p.add_run()
        r2.text = sv
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_MAIN

    # Bottom Right: Production Tech Stack
    _, tf_tc = add_card(slide3, Inches(6.65), Inches(4.72), Inches(6.08), Inches(1.98),
                        bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_tc_h = tf_tc.paragraphs[0]
    p_tc_h.text = "PRODUCTION TECHNOLOGY STACK"
    p_tc_h.font.name = "Arial"
    p_tc_h.font.bold = True
    p_tc_h.font.size = Pt(9.5)
    p_tc_h.font.color.rgb = CLR_CYAN_ACCENT
    p_tc_h.space_after = Pt(3)

    t_stack = [
        ("Frontend / UI", "React 18, TypeScript, Vite, MapLibre GL, WebGL, TailwindCSS, PWA"),
        ("Backend / API", "FastAPI, Python 3.10+, Asynchronous DAG Orchestrator, SQLite/WAL"),
        ("Intelligence", "NumPy, SciPy, XGBoost (HSI Matrix), Modified A* Graph Pathfinder"),
        ("IoT / Telecom", "ESP32 LoRa Gateway (868/433 MHz), Multi-Hop Mesh (OLSR/BATMAN)")
    ]
    for cat, spec in t_stack:
        p = tf_tc.add_paragraph()
        p.space_after = Pt(2.5)
        r1 = p.add_run()
        r1.text = f"• {cat}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = CLR_CYAN_ACCENT
        r2 = p.add_run()
        r2.text = spec
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = RGBColor(203, 213, 225)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides[3]
    print("Formatting Slide 4: FEASIBILITY AND VIABILITY...")
    setup_slide_header(slide4, "FEASIBILITY AND VIABILITY: Build, Mitigate & Deploy",
                       [("FEASIBILITY ANALYSIS", CLR_OCEAN_BLUE),
                        ("POTENTIAL RISKS & CHALLENGES", CLR_AMBER),
                        ("STRATEGIES FOR OVERCOMING", CLR_EMERALD),
                        ("DEPLOYMENT & ECONOMICS", CLR_NAVY_TITLE)])

    # Column 1: Implementation Roadmap
    _, tf_rd = add_card(slide4, Inches(0.6), Inches(1.38), Inches(3.0), Inches(5.32),
                        bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_rh = tf_rd.paragraphs[0]
    p_rh.text = "IMPLEMENTATION ROADMAP"
    p_rh.font.name = "Arial"
    p_rh.font.bold = True
    p_rh.font.size = Pt(10)
    p_rh.font.color.rgb = CLR_CYAN_ACCENT
    p_rh.space_after = Pt(6)

    milestones = [
        ("✓ PROTOTYPE", "End-to-End FastAPI & React UI working", CLR_EMERALD),
        ("✓ LIVE DATA", "Open-Meteo & IMD ingestion active", CLR_EMERALD),
        ("✓ RISK ENGINE", "Deterministic safety breaker active", CLR_EMERALD),
        ("✓ ROUTE ENGINE", "A* dynamic pathfinder & HSI matrix", CLR_EMERALD),
        ("◐ HARBOR PILOT", "Kanyakumari & Munambam trial", CLR_AMBER),
        ("◌ SCALE PHASE", "9 Coastal States & ISRO Satcom", RGBColor(148, 163, 184))
    ]
    for title, desc, col in milestones:
        p = tf_rd.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = col
        r2 = p.add_run()
        r2.text = desc
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_WHITE

    # Column 2: Challenge ➔ Mitigation
    _, tf_ch = add_card(slide4, Inches(3.75), Inches(1.38), Inches(5.15), Inches(5.32),
                        bg_color=CLR_CARD_BG_OCEAN, border_color=CLR_NAVY_TITLE, border_width=1.0)
    p_chh = tf_ch.paragraphs[0]
    p_chh.text = "CRITICAL RISKS ➔ PROVEN MITIGATION STRATEGIES"
    p_chh.font.name = "Arial"
    p_chh.font.bold = True
    p_chh.font.size = Pt(10)
    p_chh.font.color.rgb = CLR_NAVY_TITLE
    p_chh.space_after = Pt(6)

    challenges = [
        ("Data Latency / Stale Feeds", "Automated TTL monitoring (<30m) & cached offline fallback"),
        ("API Provider Outages", "Multi-source failover cascade (INCOIS ➔ Open-Meteo ➔ MOSDAC)"),
        ("Source Disagreements", "Provenance-weighted ensemble & explicit uncertainty bounds"),
        ("Offshore Signal Blackout", "Multi-Hop LoRa Fleet Mesh (up to 50 km) + Offline PWA Caching"),
        ("Vessel Hull Diversity", "Parametric Digital Twin model library (Catamaran, FRP, Trawler)"),
        ("Model Hallucination", "100% Deterministic Safety Rules (Zero AI in safety path)")
    ]
    for ct, md in challenges:
        p = tf_ch.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"⚠️ {ct}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = CLR_ROSE
        r2 = p.add_run()
        r2.text = f"  ➔ 🛡️ {md}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_MAIN

    # Column 3 Top: Deployment Architecture
    _, tf_dp = add_card(slide4, Inches(9.05), Inches(1.38), Inches(3.68), Inches(2.60),
                        bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_dph = tf_dp.paragraphs[0]
    p_dph.text = "DEPLOYMENT ARCHITECTURE"
    p_dph.font.name = "Arial"
    p_dph.font.bold = True
    p_dph.font.size = Pt(9.5)
    p_dph.font.color.rgb = CLR_CYAN_ACCENT
    p_dph.space_after = Pt(4)

    dep_steps = [
        ("1. Vessel Edge", "GPS / IMU + ESP32 LoRa Node"),
        ("2. Fleet Relay", "Multi-Hop Boat Mesh (3-8 km hops)"),
        ("3. Coastal Mast", "Lighthouse LoRa Gateway (60m mast)"),
        ("4. Cloud Engine", "FastAPI / SQLite / Orchestrator"),
        ("5. User Tiers", "Fisher PWA | Fleet Web | Coast Guard")
    ]
    for ds, dv in dep_steps:
        p = tf_dp.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"{ds}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = CLR_CYAN_ACCENT
        r2 = p.add_run()
        r2.text = dv
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_WHITE

    # Column 3 Bottom: Deployment Economics & BOM
    _, tf_cs = add_card(slide4, Inches(9.05), Inches(4.10), Inches(3.68), Inches(2.60),
                        bg_color=CLR_CARD_BG_NAVY, border_color=CLR_EMERALD, border_width=1.0)
    p_csh = tf_cs.paragraphs[0]
    p_csh.text = "DEPLOYMENT ECONOMICS & BOM"
    p_csh.font.name = "Arial"
    p_csh.font.bold = True
    p_csh.font.size = Pt(9.5)
    p_csh.font.color.rgb = CLR_EMERALD
    p_csh.space_after = Pt(4)

    cost_items = [
        ("ESP32 MCU + LoRa SX1262", "₹1,800 (~$22)"),
        ("GPS NEO-6M Receiver", "₹600 (~$7)"),
        ("IP67 Enclosure + Solar", "₹1,100 (~$13)"),
        ("Total Hardware Cost / Boat", "≈ ₹3,500 ($42)"),
        ("Data & Server API Cost", "₹0 (Open Feeds / Serverless)")
    ]
    for cn, cv in cost_items:
        p = tf_cs.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"• {cn}: "
        r1.font.name = "Arial"
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = CLR_TEXT_WHITE
        r2 = p.add_run()
        r2.text = cv
        r2.font.name = "Arial"
        r2.font.bold = True
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_EMERALD if "Total" in cn or "₹0" in cv else CLR_CYAN_ACCENT

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides[4]
    print("Formatting Slide 5: IMPACT AND BENEFITS...")
    setup_slide_header(slide5, "IMPACT AND BENEFITS: Empowering India's Blue Economy",
                       [("TARGET AUDIENCE IMPACT", CLR_OCEAN_BLUE),
                        ("SOCIAL & HUMAN BENEFITS", CLR_EMERALD),
                        ("ECONOMIC & OPERATIONAL ROI", CLR_AMBER),
                        ("NATIONAL STRATEGIC SCALE", CLR_NAVY_TITLE)])

    # Left Container: ORCA Impact Ecosystem
    _, tf_ec = add_card(slide5, Inches(0.6), Inches(1.38), Inches(5.5), Inches(4.75),
                        bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_ech = tf_ec.paragraphs[0]
    p_ech.text = "ORCA MULTI-STAKEHOLDER ECOSYSTEM"
    p_ech.font.name = "Arial"
    p_ech.font.bold = True
    p_ech.font.size = Pt(10.5)
    p_ech.font.color.rgb = CLR_CYAN_ACCENT
    p_ech.space_after = Pt(6)

    stakeholders = [
        ("🎣 Artisanal Fishermen (4.5M Lives)", "Zero preventable capsizings, localized vernacular audio safety alerts, precision PFZ fishing hotspots."),
        ("⚓ Harbors & Port Authorities", "Live harbor fleet tracking, automated departure safety clearances, harbor congestion mitigation."),
        ("🚢 Commercial Fishing Fleets", "Optimal fuel-saving route planning, dynamic weather avoidance, higher catch yield ROI."),
        ("🛡️ Coast Guard & NDRF", "Real-time Dark-Fleet anomaly alerts, 1,000-particle Monte Carlo SAR drift tracking."),
        ("🔬 Oceanographic Institutes", "Crowdsourced ground-truth sea state telemetry & high-resolution model validation.")
    ]
    for st, sd in stakeholders:
        p = tf_ec.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"{st}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = CLR_CYAN_ACCENT
        r2 = p.add_run()
        r2.text = f"  • {sd}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = RGBColor(203, 213, 225)

    # Right Container: 4 Pillars of National Impact
    _, tf_pl = add_card(slide5, Inches(6.25), Inches(1.38), Inches(6.48), Inches(4.75),
                        bg_color=CLR_CARD_BG_OCEAN, border_color=CLR_NAVY_TITLE, border_width=1.0)
    p_plh = tf_pl.paragraphs[0]
    p_plh.text = "THE FOUR PILLARS OF NATIONAL IMPACT"
    p_plh.font.name = "Arial"
    p_plh.font.bold = True
    p_plh.font.size = Pt(10.5)
    p_plh.font.color.rgb = CLR_NAVY_TITLE
    p_plh.space_after = Pt(6)

    pillars = [
        ("👥 SOCIAL IMPACT (Saving Lives)", 
         "Eliminates fatal capsizings by calculating real-time wave-vessel physical resonance. Delivers plain-language voice alerts in Tamil, Malayalam, Bengali, Telugu, and Hindi for non-literate crews.",
         CLR_EMERALD),
        ("💰 ECONOMIC IMPACT (Fuel & Yield)", 
         "Reduces diesel expenditure by 20–30% through current-assisted dynamic A* pathfinding. Maximizes high-value pelagic catch yields via multi-species Habitat Suitability Index (HSI) mapping.",
         CLR_AMBER),
        ("🌿 ENVIRONMENTAL IMPACT (Ocean Health)", 
         "Substantially lowers maritime diesel emissions per fishing voyage. Provides automated geofencing to prevent accidental intrusion into Marine Protected Areas (MPAs) and international boundaries.",
         CLR_OCEAN_BLUE),
        ("🇮🇳 STRATEGIC IMPACT (Digital Sovereignty)", 
         "Delivers end-to-end sovereign maritime domain awareness across India's 7,516 km coastline and 2.37M sq km EEZ, creating an integrated national ocean safety infrastructure.",
         CLR_NAVY_TITLE)
    ]
    for p_name, p_body, p_col in pillars:
        p = tf_pl.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = f"{p_name}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = p_col
        r2 = p.add_run()
        r2.text = p_body
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_MAIN

    # Bottom Progression Banner
    _, tf_btm5 = add_card(slide5, Inches(0.6), Inches(6.25), Inches(12.13), Inches(0.45),
                          bg_color=CLR_CARD_BG_DARK, border_color=CLR_CYAN_ACCENT, border_width=1.0)
    p_b5 = tf_btm5.paragraphs[0]
    p_b5.text = "1 ARTISANAL VESSEL   ➔   HARBOR FLEET (500+)   ➔   COASTAL STATE   ➔   PAN-INDIA MARITIME DIGITAL TWIN LAYER"
    p_b5.font.name = "Arial"
    p_b5.font.bold = True
    p_b5.font.size = Pt(9.5)
    p_b5.font.color.rgb = CLR_CYAN_ACCENT
    p_b5.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides[5]
    print("Formatting Slide 6: RESEARCH AND REFERENCES...")
    setup_slide_header(slide6, "RESEARCH  AND REFERENCES: Science to Implementation",
                       [("DETAILS OF RESEARCH WORK", CLR_OCEAN_BLUE),
                        ("IMO STABILITY CRITERIA", CLR_NAVY_TITLE),
                        ("HYDRODYNAMIC WAVE PHYSICS", CLR_NAVY_TITLE),
                        ("VERIFIED ACADEMIC CITATIONS", CLR_EMERALD)])

    # 4 Core Research Cards
    research_cards = [
        ("IMO Code on Intact Stability (2008)",
         "Resolution MSC.267(85) / Small Craft Capsizing Dynamics",
         "Small vessel dynamic instability under steep wave action, beam-sea broaching, and parametric roll resonance.",
         "Deterministic Safety Engine: Computes critical wave threshold (Hcrit ≈ 0.6 · Lhull) and enforces non-bypassable hazard overrides (<10ms).",
         CLR_CYAN_ACCENT),
        ("Marine Hydrodynamics & Encounter Physics",
         "O.M. Faltinsen (1990) & J.N. Newman (1977)",
         "Wave steepness ratio (S = Hs / λ), deep-water dispersion relations (λ = g·Tp² / 2π), and encounter angle frequency shifts.",
         "Vessel Digital Twin: Dynamic calculation of wave steepness & vessel-wave interaction rather than treating sea states as static point forecasts.",
         CLR_EMERALD),
        ("Satellite Oceanography & Pelagic Habitats",
         "INCOIS PFZ Mission & ISRO MOSDAC Telemetry",
         "Aggregation of pelagic marine species along thermal gradients, SST frontal boundaries, and chlorophyll-a density concentration zones.",
         "Multi-Species Habitat Index (HSI): XGBoost model integrated with 2D Sobel spatial filter for high-precision fishery zone recommendations.",
         CLR_AMBER),
        ("Dynamic Constrained Pathfinding & Graph Optimization",
         "Modified A* over Time-Varying Navigational Cost Surfaces",
         "Optimization of multiobjective vessel routes under dynamic weather constraints, geofenced boundaries, and risk cost penalization.",
         "ORCA Navigational Pathfinder: Computes 4-point optimized waypoints balancing weather safety, fuel consumption, and dock proximity.",
         CLR_OCEAN_BLUE)
    ]

    for idx, (title, src, insight, app, bcol) in enumerate(research_cards):
        col = idx % 2
        row = idx // 2
        x = Inches(0.6 + col * 6.18)
        y = Inches(1.38 + row * 2.38)
        
        _, tf_rc = add_card(slide6, x, y, Inches(5.95), Inches(2.28),
                            bg_color=CLR_CARD_BG_DARK, border_color=bcol, border_width=1.0)
        p1 = tf_rc.paragraphs[0]
        p1.text = f"📚 {title}"
        p1.font.name = "Arial"
        p1.font.bold = True
        p1.font.size = Pt(9)
        p1.font.color.rgb = CLR_CYAN_ACCENT
        
        p_src = tf_rc.add_paragraph()
        p_src.text = f"Source: {src}"
        p_src.font.name = "Arial"
        p_src.font.italic = True
        p_src.font.size = Pt(7)
        p_src.font.color.rgb = RGBColor(148, 163, 184)
        p_src.space_after = Pt(2)
        
        p_ins = tf_rc.add_paragraph()
        r1 = p_ins.add_run()
        r1.text = "Key Insight: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = CLR_AMBER
        r2 = p_ins.add_run()
        r2.text = insight
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CLR_TEXT_WHITE
        p_ins.space_after = Pt(2)

        p_app = tf_rc.add_paragraph()
        r3 = p_app.add_run()
        r3.text = "ORCA Implementation: "
        r3.font.name = "Arial"
        r3.font.bold = True
        r3.font.size = Pt(7.5)
        r3.font.color.rgb = CLR_EMERALD
        r4 = p_app.add_run()
        r4.text = app
        r4.font.name = "Arial"
        r4.font.size = Pt(7.5)
        r4.font.color.rgb = RGBColor(203, 213, 225)

    # Bottom Citations Container
    _, tf_rf = add_card(slide6, Inches(0.6), Inches(6.20), Inches(12.13), Inches(0.52),
                        bg_color=CLR_CARD_BG_OCEAN, border_color=CLR_NAVY_TITLE, border_width=1.0)
    p_rf = tf_rf.paragraphs[0]
    p_rf.text = "AUTHENTIC OPERATIONAL & SCIENTIFIC REFERENCES"
    p_rf.font.name = "Arial"
    p_rf.font.bold = True
    p_rf.font.size = Pt(8)
    p_rf.font.color.rgb = CLR_NAVY_TITLE
    
    p_cits = tf_rf.add_paragraph()
    p_cits.text = (
        "[1] IMO Intact Stability Code (2008). Res. MSC.267(85). | "
        "[2] INCOIS Ocean State Forecast & PFZ Protocols. | "
        "[3] ISRO MOSDAC Oceansat-3 & INSAT-3DR Products. | "
        "[4] Copernicus Marine CMEMS Analysis. | "
        "[5] WMO-No. 558 Marine Services."
    )
    p_cits.font.name = "Arial"
    p_cits.font.size = Pt(7)
    p_cits.font.color.rgb = CLR_TEXT_MAIN

    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Blended Presentation generated successfully!")
    print(f"Saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_presentation()
