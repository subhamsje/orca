import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = "/Users/subham/Downloads/test_6slides.pptx"
OUTPUT_PATH = "/Users/subham/Downloads/ORCA_4.0_SIH2026_Official_10_Slide_Deck.pptx"
LOGO_PATH = "/Users/subham/code/orca/sih_logo_extracted.png"

# ==============================================================================
# DESIGN SYSTEM PALETTE: PALANTIR x MARITIME COMMAND x MISSION CONTROL
# ==============================================================================
C_WHITE        = RGBColor(255, 255, 255)
C_SLATE_50     = RGBColor(248, 250, 252)   # #f8fafc - Base slide canvas
C_SLATE_100    = RGBColor(241, 245, 249)   # #f1f5f9 - Card fill
C_SLATE_200    = RGBColor(226, 232, 240)   # #e2e8f0 - Crisp border
C_SLATE_700    = RGBColor(51, 65, 85)      # #334155 - Body typography
C_SLATE_900    = RGBColor(15, 23, 42)      # #0f172a - Dark bold text
C_MUTED        = RGBColor(100, 116, 139)   # #64748b - Subtitle / labels

C_NAVY_HERO    = RGBColor(10, 37, 64)      # #0a2540 - Deep maritime hero navy
C_NAVY_CARD    = RGBColor(15, 34, 64)      # #0f2240 - Navy inner cards
C_BLUE_PRIMARY = RGBColor(2, 132, 199)     # #0284c7 - Ocean blue
C_BLUE_LIGHT   = RGBColor(240, 249, 255)   # #f0f9ff - Ice blue fill
C_BLUE_BORDER  = RGBColor(186, 230, 253)   # #bae6fd - Ice blue border
C_CYAN_ACCENT  = RGBColor(6, 182, 212)     # #06b6d4 - Vibrant cyan

C_EMERALD      = RGBColor(16, 185, 129)    # #10b981 - Green Safe
C_EMERALD_BG   = RGBColor(236, 253, 245)   # #ecfdf5
C_EMERALD_BORDER = RGBColor(167, 243, 208)

C_AMBER        = RGBColor(217, 119, 6)     # #d97706 - Warning Amber
C_AMBER_BG     = RGBColor(254, 243, 199)   # #fef3c7
C_AMBER_BORDER = RGBColor(253, 230, 138)

C_ROSE         = RGBColor(225, 29, 72)     # #e11d48 - Danger Rose
C_ROSE_BG      = RGBColor(255, 241, 242)   # #fff1f2

def set_shape_flat(shape, fill_color, border_color=None, border_width=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

def create_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_flat(shape, bg_color, border_color, border_width)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return shape, tf

def apply_official_slide_frame(slide, slide_num, total_slides, section_badge, title_text, subtitle_text=""):
    # Clear any default placeholders from the layout
    for s in list(slide.shapes):
        if s.name.startswith("Title") or s.name.startswith("TextBox") or s.name.startswith("Oval") or s.name.startswith("Rectangle") or s.name.startswith("Picture") or s.name.startswith("Footer") or s.name.startswith("Slide Number"):
            # We will build clean official template chrome
            pass

    # Top Team Oval Badge (Left: 0.36, top: 0.18)
    oval_shape, oval_tf = create_card(slide, Inches(0.36), Inches(0.18), Inches(1.35), Inches(0.72),
                                      bg_color=C_NAVY_HERO, border_color=None, shape_type=MSO_SHAPE.OVAL)
    p_ov = oval_tf.paragraphs[0]
    p_ov.text = "ORCA 4.0"
    p_ov.font.name = "Arial"
    p_ov.font.size = Pt(11)
    p_ov.font.bold = True
    p_ov.font.color.rgb = C_WHITE
    p_ov.alignment = PP_ALIGN.CENTER

    # Official SIH Top Right Logo Picture
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(10.70), Inches(0.02), Inches(2.46), Inches(1.10))

    # Center Title & Subtitle (Left: 1.85, width: 8.75)
    t_box = slide.shapes.add_textbox(Inches(1.85), Inches(0.10), Inches(8.75), Inches(0.85))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
    
    p_title = tf_t.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(17)
    p_title.font.bold = True
    p_title.font.color.rgb = C_NAVY_HERO

    if subtitle_text:
        p_sub = tf_t.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(9)
        p_sub.font.color.rgb = C_MUTED

    # Top Section Tag & Flow Indicator Bar (top: 0.96, height: 0.28)
    ribbon_w = Inches(12.13)
    _, tf_rb = create_card(slide, Inches(0.6), Inches(0.96), ribbon_w, Inches(0.28),
                           bg_color=C_SLATE_100, border_color=C_SLATE_200, border_width=1.0)
    p_rb = tf_rb.paragraphs[0]
    
    r_tag = p_rb.add_run()
    r_tag.text = f"{slide_num:02d}/{total_slides:02d}  ·  {section_badge.upper()}    |    "
    r_tag.font.name = "Arial"
    r_tag.font.size = Pt(8)
    r_tag.font.bold = True
    r_tag.font.color.rgb = C_BLUE_PRIMARY

    chain_steps = ["OCEAN", "WEATHER", "VESSEL", "ROUTE", "ORCA", "RISK", "DECISION"]
    active_idx = min(slide_num - 1, len(chain_steps) - 1)
    
    r_fl = p_rb.add_run()
    r_fl.text = "FLOW: "
    r_fl.font.name = "Arial"
    r_fl.font.size = Pt(7.5)
    r_fl.font.bold = True
    r_fl.font.color.rgb = C_MUTED

    for c_i, c_name in enumerate(chain_steps):
        r_step = p_rb.add_run()
        r_step.text = c_name
        r_step.font.name = "Arial"
        r_step.font.size = Pt(7.5)
        r_step.font.bold = (c_i == active_idx)
        r_step.font.color.rgb = C_BLUE_PRIMARY if (c_i == active_idx) else C_MUTED
        if c_i < len(chain_steps) - 1:
            r_arr = p_rb.add_run()
            r_arr.text = " ➔ "
            r_arr.font.name = "Arial"
            r_arr.font.size = Pt(7)
            r_arr.font.color.rgb = C_MUTED

    # Official Bottom Footer Bar (top: 6.95, height: 0.40)
    btm_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), Inches(13.333), Inches(0.55))
    set_shape_flat(btm_bar, C_WHITE, C_SLATE_200, 0.5)

    # Footer Text
    ft_box = slide.shapes.add_textbox(Inches(4.5), Inches(7.02), Inches(4.33), Inches(0.35))
    tf_ft = ft_box.text_frame
    tf_ft.margin_left = tf_ft.margin_right = tf_ft.margin_top = tf_ft.margin_bottom = 0
    p_ft = tf_ft.paragraphs[0]
    p_ft.text = "@SIH Idea submission- Template  ·  ORCA 4.0"
    p_ft.font.name = "Arial"
    p_ft.font.size = Pt(8.5)
    p_ft.font.color.rgb = C_MUTED
    p_ft.alignment = PP_ALIGN.CENTER

    # Slide Number
    num_box = slide.shapes.add_textbox(Inches(11.0), Inches(7.02), Inches(1.7), Inches(0.35))
    tf_num = num_box.text_frame
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    p_num = tf_num.paragraphs[0]
    p_num.text = f"{slide_num}"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(9)
    p_num.font.bold = True
    p_num.font.color.rgb = C_NAVY_HERO
    p_num.alignment = PP_ALIGN.RIGHT

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    print("Building Official 10-Slide Template Presentation for ORCA 4.0...")

    # =========================================================================
    # SLIDE 1: TITLE PAGE (WHO?)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    print("Formatting Slide 1: TITLE PAGE...")

    # Slide 1 Header Bar
    top_bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.15))
    set_shape_flat(top_bar1, C_WHITE, C_SLATE_200, 0.5)

    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(10.70), Inches(0.02), Inches(2.46), Inches(1.10))

    t_box1 = slide1.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(9.8), Inches(0.85))
    tf_t1 = t_box1.text_frame
    tf_t1.margin_left = tf_t1.margin_right = tf_t1.margin_top = tf_t1.margin_bottom = 0
    p_sih_title = tf_t1.paragraphs[0]
    p_sih_title.text = "SMART INDIA HACKATHON 2026"
    p_sih_title.font.name = "Arial"
    p_sih_title.font.size = Pt(20)
    p_sih_title.font.bold = True
    p_sih_title.font.color.rgb = C_NAVY_HERO
    
    p_sih_sub = tf_t1.add_paragraph()
    p_sih_sub.text = "OFFICIAL IDEA PRESENTATION  ·  NATIONAL FINALS"
    p_sih_sub.font.name = "Arial"
    p_sih_sub.font.size = Pt(10)
    p_sih_sub.font.bold = True
    p_sih_sub.font.color.rgb = C_BLUE_PRIMARY

    # Left: Official Form Details Card (width: 5.8)
    _, tf_sih_form = create_card(slide1, Inches(0.6), Inches(1.30), Inches(5.8), Inches(5.50),
                                 bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    
    p_fh = tf_sih_form.paragraphs[0]
    p_fh.text = "PROPOSAL METADATA & TEAM REGISTRATION"
    p_fh.font.name = "Arial"
    p_fh.font.bold = True
    p_fh.font.size = Pt(10.5)
    p_fh.font.color.rgb = C_NAVY_HERO
    p_fh.space_after = Pt(8)

    sih_fields = [
        ("Problem Statement ID", "SIH26176"),
        ("Problem Statement Title", "AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform for Artisanal Fishermen & Coastal Fleets"),
        ("Sponsoring Ministry", "ISRO / Department of Space & INCOIS"),
        ("Theme", "Space Technology / Blue Economy / Disaster Management"),
        ("PS Category", "Software (with Edge Hardware compatibility)"),
        ("Team ID / Registered Name", "ORCA 4.0")
    ]
    for label, val in sih_fields:
        p_f = tf_sih_form.add_paragraph()
        p_f.space_after = Pt(5)
        r_lbl = p_f.add_run()
        r_lbl.text = f"{label}:\n"
        r_lbl.font.name = "Arial"
        r_lbl.font.bold = True
        r_lbl.font.size = Pt(8.5)
        r_lbl.font.color.rgb = C_NAVY_HERO
        r_val = p_f.add_run()
        r_val.text = val
        r_val.font.name = "Arial"
        r_val.font.size = Pt(9.5)
        r_val.font.color.rgb = C_SLATE_700

    # Right: Dark Navy ORCA 4.0 Showcase Panel (width: 6.1)
    _, tf_hero = create_card(slide1, Inches(6.6), Inches(1.30), Inches(6.13), Inches(5.50),
                             bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.5)
    
    p_h_badge = tf_hero.paragraphs[0]
    p_h_badge.text = "MARITIME DECISION OPERATING SYSTEM"
    p_h_badge.font.name = "Arial"
    p_h_badge.font.size = Pt(9.5)
    p_h_badge.font.bold = True
    p_h_badge.font.color.rgb = C_CYAN_ACCENT
    p_h_badge.space_after = Pt(4)

    p_h_title = tf_hero.add_paragraph()
    p_h_title.text = "ORCA 4.0"
    p_h_title.font.name = "Arial"
    p_h_title.font.size = Pt(36)
    p_h_title.font.bold = True
    p_h_title.font.color.rgb = C_WHITE
    p_h_title.space_after = Pt(2)

    p_h_tag = tf_hero.add_paragraph()
    p_h_tag.text = "“FROM OCEAN DATA TO SAFER DECISIONS”"
    p_h_tag.font.name = "Arial"
    p_h_tag.font.size = Pt(13)
    p_h_tag.font.bold = True
    p_h_tag.font.color.rgb = RGBColor(56, 189, 248)
    p_h_tag.space_after = Pt(6)

    p_h_desc = tf_hero.add_paragraph()
    p_h_desc.text = "Transforming fragmented satellite, wave & atmospheric telemetry into deterministic, vessel-aware voyage safety decisions."
    p_h_desc.font.name = "Arial"
    p_h_desc.font.size = Pt(9.5)
    p_h_desc.font.color.rgb = RGBColor(203, 213, 225)
    p_h_desc.space_after = Pt(10)

    # Hero Decision Chain Ribbon
    p_dc_lbl = tf_hero.add_paragraph()
    p_dc_lbl.text = "CORE DECISION CHAIN:"
    p_dc_lbl.font.name = "Arial"
    p_dc_lbl.font.size = Pt(8.5)
    p_dc_lbl.font.bold = True
    p_dc_lbl.font.color.rgb = C_EMERALD
    
    p_dc_flow = tf_hero.add_paragraph()
    p_dc_flow.text = "OCEAN  ➔  WEATHER  ➔  VESSEL  ➔  ROUTE  ➔  ORCA  ➔  RISK  ➔  DECISION"
    p_dc_flow.font.name = "Arial"
    p_dc_flow.font.size = Pt(8)
    p_dc_flow.font.bold = True
    p_dc_flow.font.color.rgb = C_WHITE
    p_dc_flow.space_after = Pt(10)

    s1_pillars = [
        ("🛡️ Deterministic Safety Shield", "Hard physical capsizing rules (Hcrit = 0.6·L) with zero AI hallucination risk", C_EMERALD),
        ("🚢 Parametric Vessel Twin", "Dynamic wave-vessel encounter physics tailored to hull length, draft & beam", RGBColor(56, 189, 248)),
        ("📻 Multi-Hop LoRa Mesh", "Overcomes 12 NM 4G blackout for zero-cost offshore telemetry up to 50 km", C_AMBER)
    ]
    for p_t, p_d, p_c in s1_pillars:
        p_pil = tf_hero.add_paragraph()
        p_pil.space_after = Pt(4)
        r_pt = p_pil.add_run()
        r_pt.text = f"{p_t}: "
        r_pt.font.name = "Arial"
        r_pt.font.bold = True
        r_pt.font.size = Pt(9)
        r_pt.font.color.rgb = p_c
        r_pd = p_pil.add_run()
        r_pd.text = p_d
        r_pd.font.name = "Arial"
        r_pd.font.size = Pt(8)
        r_pd.font.color.rgb = RGBColor(203, 213, 225)

    # Footer on Slide 1
    btm_bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), Inches(13.333), Inches(0.55))
    set_shape_flat(btm_bar1, C_WHITE, C_SLATE_200, 0.5)
    ft_box1 = slide1.shapes.add_textbox(Inches(4.5), Inches(7.02), Inches(4.33), Inches(0.35))
    tf_ft1 = ft_box1.text_frame
    p_ft1 = tf_ft1.paragraphs[0]
    p_ft1.text = "@SIH Idea submission- Template  ·  ORCA 4.0"
    p_ft1.font.name = "Arial"
    p_ft1.font.size = Pt(8.5)
    p_ft1.font.color.rgb = C_MUTED
    p_ft1.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: THE PROBLEM (WHY?)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide2, 2, 10, "Problem Discovery", 
                               "“THE OCEAN IS FORECAST. THE VOYAGE IS NOT.”", 
                               "Marine data is heavily fragmented across agencies. Regional weather forecasts fail to answer: What should THIS vessel do right now?")
    print("Formatting Slide 2: The Problem...")

    # Hero Thesis Callout Box
    _, tf_s2_th = create_card(slide2, Inches(0.6), Inches(1.30), Inches(12.13), Inches(0.52),
                              bg_color=C_ROSE_BG, border_color=C_ROSE, border_width=1.0)
    p_s2_th = tf_s2_th.paragraphs[0]
    r_th1 = p_s2_th.add_run()
    r_th1.text = "THE FATAL DECISION GAP:  FORECAST ≠ OPERATIONAL RISK.  "
    r_th1.font.name = "Arial"
    r_th1.font.bold = True
    r_th1.font.size = Pt(9.5)
    r_th1.font.color.rgb = C_ROSE
    r_th2 = p_s2_th.add_run()
    r_th2.text = "General weather reports broadcast broad regional numbers (e.g. 'Waves 2.5m in South Arabian Sea'). They cannot calculate whether a 9m artisanal craft will capsize under beam-sea resonance on its specific waypoint track."
    r_th2.font.name = "Arial"
    r_th2.font.size = Pt(8.5)
    r_th2.font.color.rgb = C_SLATE_900

    # Left Container: 8 Fragmented Data Silos (width: 4.8)
    _, tf_s2_l = create_card(slide2, Inches(0.6), Inches(1.88), Inches(4.8), Inches(4.95),
                             bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_l_h = tf_s2_l.paragraphs[0]
    p_l_h.text = "FRAGMENTED RAW DATA SILOS"
    p_l_h.font.name = "Arial"
    p_l_h.font.bold = True
    p_l_h.font.size = Pt(10)
    p_l_h.font.color.rgb = C_NAVY_HERO
    p_l_h.space_after = Pt(4)

    s2_silos = [
        ("🌊 INCOIS Wave Spectrum", "Significant wave height (Hs), swell period (Tp), swell direction"),
        ("💨 IMD Weather & Gusts", "Surface wind velocity, storm track vectors, cyclone alerts"),
        ("🛰️ ISRO / Copernicus SST", "Sea surface temperature gradients & ocean color rasters"),
        ("🌊 Ocean Surface Currents", "Dynamic drift velocity (u, v) from oceanographic models"),
        ("📡 AIS Fleet Telemetry", "Real-time vessel positions, speeds & collision vectors"),
        ("🚢 Vessel Physical Geometry", "Hull length (Lwl), draft, beam, engine power, freeboard"),
        ("🗺️ Planned Voyage Route", "Departure harbor, targeted fishing grounds & arrival dock"),
        ("📻 Offshore Telemetry", "0% mobile connectivity beyond 12 NM territorial waters")
    ]
    for title, desc in s2_silos:
        p_s = tf_s2_l.add_paragraph()
        p_s.space_after = Pt(2.5)
        r_st = p_s.add_run()
        r_st.text = f"• {title}: "
        r_st.font.name = "Arial"
        r_st.font.bold = True
        r_st.font.size = Pt(8)
        r_st.font.color.rgb = C_BLUE_PRIMARY
        r_sd = p_s.add_run()
        r_sd.text = desc
        r_sd.font.name = "Arial"
        r_sd.font.size = Pt(7.5)
        r_sd.font.color.rgb = C_SLATE_700

    # Center: The Existing Broken Decision Chain (width: 3.2)
    _, tf_s2_c = create_card(slide2, Inches(5.6), Inches(1.88), Inches(3.2), Inches(4.95),
                             bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_c_h = tf_s2_c.paragraphs[0]
    p_c_h.text = "TODAY: MANUAL GUESSWORK"
    p_c_h.font.name = "Arial"
    p_c_h.font.bold = True
    p_c_h.font.size = Pt(9.5)
    p_c_h.font.color.rgb = C_ROSE
    p_c_h.space_after = Pt(6)

    broken_steps = [
        ("1. Raw Weather Broadcast", "SMS text or radio broadcast with generic regional numbers"),
        ("2. Zero Vessel Context", "Identical warning sent to a 40m trawler and a 7m canoe"),
        ("3. Mental Interpretation", "Fisherman forced to guess if boat can withstand sea state"),
        ("4. Offshore Blackout", "Mobile data dies at 12 NM; crew blinded in open sea"),
        ("5. Preventable Accidents", "Capsizings, blind fuel burn (₹45k Cr lost), search delays")
    ]
    for st, sd in broken_steps:
        p_bs = tf_s2_c.add_paragraph()
        p_bs.space_after = Pt(5)
        r_bt = p_bs.add_run()
        r_bt.text = f"{st}\n"
        r_bt.font.name = "Arial"
        r_bt.font.bold = True
        r_bt.font.size = Pt(8)
        r_bt.font.color.rgb = C_SLATE_900
        r_bd = p_bs.add_run()
        r_bd.text = sd
        r_bd.font.name = "Arial"
        r_bd.font.size = Pt(7.5)
        r_bd.font.color.rgb = C_MUTED

    # Right: The ORCA Paradigm (width: 3.73)
    _, tf_s2_r = create_card(slide2, Inches(9.0), Inches(1.88), Inches(3.73), Inches(4.95),
                             bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.5)
    p_r_h = tf_s2_r.paragraphs[0]
    p_r_h.text = "ORCA DECISION INTELLIGENCE"
    p_r_h.font.name = "Arial"
    p_r_h.font.bold = True
    p_r_h.font.size = Pt(10)
    p_r_h.font.color.rgb = C_EMERALD
    p_r_h.space_after = Pt(6)

    orca_paradigm = [
        ("DATA CONTEXTUALIZATION", "Fuses 8 disparate feeds into a unified 0.083° H3 spatial world model."),
        ("VESSEL-AWARE PHYSICS", "Evaluates real-time wave encounter angle, period resonance & critical capsizing height (Hcrit)."),
        ("DYNAMIC ROUTE RISK", "Computes A* safe navigational path avoiding localized danger pockets."),
        ("ACTIONABLE VERDICT", "Plain-language vernacular voice advisory: GO, CAUTION, REROUTE or RETURN.")
    ]
    for pt, pd in orca_paradigm:
        p_op = tf_s2_r.add_paragraph()
        p_op.space_after = Pt(5)
        r_opt = p_op.add_run()
        r_opt.text = f"✓ {pt}\n"
        r_opt.font.name = "Arial"
        r_opt.font.bold = True
        r_opt.font.size = Pt(8)
        r_opt.font.color.rgb = C_NAVY_HERO
        r_opd = p_op.add_run()
        r_opd.text = pd
        r_opd.font.name = "Arial"
        r_opd.font.size = Pt(7.5)
        r_opd.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 3: THE ORCA SOLUTION (WHAT?)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide3, 3, 10, "Solution Overview", 
                               "“ORCA TURNS SEA CONDITIONS INTO A VOYAGE DECISION”", 
                               "A unified decision layer that harmonizes raw environmental state, vessel digital twins, and dynamic risk routing.")
    print("Formatting Slide 3: The ORCA Solution...")

    # Left: Inputs Funnel (width: 3.2)
    _, tf_s3_in = create_card(slide3, Inches(0.6), Inches(1.30), Inches(3.2), Inches(5.55),
                              bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_s3_ih = tf_s3_in.paragraphs[0]
    p_s3_ih.text = "1. INCOMING DATA STREAMS"
    p_s3_ih.font.name = "Arial"
    p_s3_ih.font.bold = True
    p_s3_ih.font.size = Pt(9.5)
    p_s3_ih.font.color.rgb = C_BLUE_PRIMARY
    p_s3_ih.space_after = Pt(4)

    in_chips = [
        ("INCOIS ERDDAP", "Wave spectrum & ocean currents"),
        ("IMD Cyclone Feeds", "Storm alerts & gust velocities"),
        ("ISRO / MOSDAC", "Oceansat-3 & INSAT-3DR SST"),
        ("Copernicus CMEMS", "0.083° Physical Marine Models"),
        ("AIS Transponders", "Fleet telemetry & coordinates"),
        ("Onboard NMEA", "Vessel GPS / IMU pitch-roll"),
        ("Vessel Registry", "Hull length, beam, draft profile")
    ]
    for c_name, c_sub in in_chips:
        p_c = tf_s3_in.add_paragraph()
        p_c.space_after = Pt(3.5)
        r1 = p_c.add_run()
        r1.text = f"• {c_name}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p_c.add_run()
        r2.text = f"  {c_sub}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_MUTED

    # Center: Central ORCA Engine (width: 5.2)
    _, tf_s3_eng = create_card(slide3, Inches(4.0), Inches(1.30), Inches(5.2), Inches(5.55),
                               bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.5)
    p_s3_eh = tf_s3_eng.paragraphs[0]
    p_s3_eh.text = "2. ORCA 4.0 DECISION ENGINE"
    p_s3_eh.font.name = "Arial"
    p_s3_eh.font.bold = True
    p_s3_eh.font.size = Pt(11)
    p_s3_eh.font.color.rgb = C_CYAN_ACCENT
    p_s3_eh.space_after = Pt(6)

    eng_modules = [
        ("🌐 Environmental Digital State", "Continuous 0.083° spatial grid fusing wave height, period, wind & thermal front gradients."),
        ("🚢 Parametric Vessel Twin", "Dynamic hydrodynamics model computing critical capsizing wave height (Hcrit = 0.6 · Lhull)."),
        ("🛡️ Safety Circuit Breaker", "Non-bypassable deterministic rule gate (<10ms) providing hard override on extreme sea states."),
        ("🗺️ A* Dynamic Route Risk", "Navigational cost optimizer that routes around localized high-risk sea state pockets."),
        ("🐟 Multi-Species HSI Matrix", "XGBoost machine learning predicting pelagic habitat suitability for maximum fishing yield.")
    ]
    for m_title, m_desc in eng_modules:
        p_m = tf_s3_eng.add_paragraph()
        p_m.space_after = Pt(5)
        r1 = p_m.add_run()
        r1.text = f"{m_title}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = C_WHITE
        r2 = p_m.add_run()
        r2.text = m_desc
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = RGBColor(203, 213, 225)

    # Right: 3 Actionable Outputs (width: 3.33)
    _, tf_s3_out = create_card(slide3, Inches(9.4), Inches(1.30), Inches(3.33), Inches(5.55),
                               bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_s3_oh = tf_s3_out.paragraphs[0]
    p_s3_oh.text = "3. OPERATIONAL OUTPUTS"
    p_s3_oh.font.name = "Arial"
    p_s3_oh.font.bold = True
    p_s3_oh.font.size = Pt(9.5)
    p_s3_oh.font.color.rgb = C_EMERALD
    p_s3_oh.space_after = Pt(6)

    out_verdicts = [
        ("🟢 GO / SAFE DEPARTURE", "Sea state within vessel stability limits. Direct optimal route clear.", C_EMERALD_BG, C_EMERALD),
        ("🟡 CAUTION / REROUTE", "Localized wave danger detected. A* dynamic diversion generated.", C_AMBER_BG, C_AMBER),
        ("🔴 DANGER / STAY ASHORE", "Wave height exceeds Hcrit or active cyclone alert. Deterministic hard stop.", C_ROSE_BG, C_ROSE)
    ]
    for v_t, v_d, v_bg, v_c in out_verdicts:
        p_v = tf_s3_out.add_paragraph()
        p_v.space_after = Pt(5)
        r1 = p_v.add_run()
        r1.text = f"{v_t}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = v_c
        r2 = p_v.add_run()
        r2.text = v_d
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    p_voice = tf_s3_out.add_paragraph()
    p_voice.space_after = Pt(3)
    r_vh = p_voice.add_run()
    r_vh.text = "📢 Vernacular Audio Advisory\n"
    r_vh.font.name = "Arial"
    r_vh.font.bold = True
    r_vh.font.size = Pt(8)
    r_vh.font.color.rgb = C_BLUE_PRIMARY
    r_vb = p_voice.add_run()
    r_vb.text = "Structured JSON converted to natural voice in Tamil, Malayalam, Bengali, Telugu & Hindi for non-literate crews."
    r_vb.font.name = "Arial"
    r_vb.font.size = Pt(7.5)
    r_vb.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 4: WHY ORCA IS DIFFERENT (DIFFERENT FROM WHAT?)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide4, 4, 10, "Competitive Landscape", 
                               "“OTHERS SHOW THE SEA. ORCA UNDERSTANDS THE VOYAGE.”", 
                               "The fundamental difference is not collecting more weather data — it is injecting vessel context and deterministic physics.")
    print("Formatting Slide 4: Competitive Differentiation...")

    # Left: Comparison Matrix (width: 7.6)
    _, tf_s4_m = create_card(slide4, Inches(0.6), Inches(1.30), Inches(7.6), Inches(5.55),
                             bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_s4_mh = tf_s4_m.paragraphs[0]
    p_s4_mh.text = "CAPABILITY EVALUATION MATRIX"
    p_s4_mh.font.name = "Arial"
    p_s4_mh.font.bold = True
    p_s4_mh.font.size = Pt(10)
    p_s4_mh.font.color.rgb = C_NAVY_HERO
    p_s4_mh.space_after = Pt(4)

    comp_rows = [
        ("System Capability", "Standard Weather Apps", "Marine Port Portals", "ORCA 4.0 Platform"),
        ("Weather & Wave Forecast", "● Broad Regional", "● Coarse Grid", "✓ High-Res (0.083° H3)"),
        ("Cyclone & Hazard Alerts", "● Text / SMS Alert", "● Static Bulletin", "✓ Realtime Geo-Fenced"),
        ("Vessel Hull & Draft Context", "— None (Ignored)", "— None (Ignored)", "✓ Parametric Vessel Twin"),
        ("Wave-Vessel Dynamic Physics", "— None", "— None", "✓ Hcrit & Resonance Math"),
        ("Dynamic Route Pathfinder", "— None (Straight Line)", "— Fixed Port Lanes", "✓ A* Risk Cost Avoidance"),
        ("Deterministic Safety Guard", "— None (Unverified)", "— Manual Port Call", "✓ Hard Non-Bypassable Rule"),
        ("Multi-Species Habitat (PFZ)", "— None", "● Coarse PFZ Bulletin", "✓ XGBoost Habitat Matrix"),
        ("Offshore Reach Beyond 12 NM", "— 0% (GSM Blind)", "● Expensive Satellite", "✓ Multi-Hop LoRa Mesh (50 km)")
    ]

    for idx, (c1, c2, c3, c4) in enumerate(comp_rows):
        p = tf_s4_m.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"{c1.ljust(26)} "
        r1.font.name = "Arial"
        r1.font.bold = True if idx == 0 else False
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_NAVY_HERO if idx == 0 else C_SLATE_900

        r2 = p.add_run()
        r2.text = f"{c2.ljust(18)} "
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_ROSE if "—" in c2 else C_MUTED

        r3 = p.add_run()
        r3.text = f"{c3.ljust(18)} "
        r3.font.name = "Arial"
        r3.font.size = Pt(7)
        r3.font.color.rgb = C_ROSE if "—" in c3 else C_MUTED

        r4 = p.add_run()
        r4.text = c4
        r4.font.name = "Arial"
        r4.font.bold = True
        r4.font.size = Pt(7.5)
        r4.font.color.rgb = C_EMERALD if "✓" in c4 else C_NAVY_HERO

    # Right: Context Difference Box (width: 4.33)
    _, tf_s4_r = create_card(slide4, Inches(8.4), Inches(1.30), Inches(4.33), Inches(5.55),
                             bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER, border_width=1.0)
    p_s4_rh = tf_s4_r.paragraphs[0]
    p_s4_rh.text = "THE ARCHITECTURAL SHIFT"
    p_s4_rh.font.name = "Arial"
    p_s4_rh.font.bold = True
    p_s4_rh.font.size = Pt(10)
    p_s4_rh.font.color.rgb = C_BLUE_PRIMARY
    p_s4_rh.space_after = Pt(6)

    p_quote = tf_s4_r.add_paragraph()
    p_quote.text = "“THE DIFFERENCE IS NOT MORE DATA. THE DIFFERENCE IS CONTEXT.”"
    p_quote.font.name = "Arial"
    p_quote.font.size = Pt(12)
    p_quote.font.bold = True
    p_quote.font.color.rgb = C_NAVY_HERO
    p_quote.space_after = Pt(8)

    p_exp = tf_s4_r.add_paragraph()
    p_exp.text = "Existing platforms treat all mariners as identical passive readers of raw weather maps. ORCA treats every fishing voyage as an active physical system with specific hydrodynamic limits, route deadlines, and economic constraints."
    p_exp.font.name = "Arial"
    p_exp.font.size = Pt(8)
    p_exp.font.color.rgb = C_SLATE_700
    p_exp.space_after = Pt(10)

    p_flow_h = tf_s4_r.add_paragraph()
    p_flow_h.text = "ORCA PROGRESSION:"
    p_flow_h.font.name = "Arial"
    p_flow_h.font.bold = True
    p_flow_h.font.size = Pt(8)
    p_flow_h.font.color.rgb = C_EMERALD

    p_flow = tf_s4_r.add_paragraph()
    p_flow.text = "DATA ➔ CONTEXT ➔ PHYSICS ➔ RISK ➔ ROUTE ➔ DECISION"
    p_flow.font.name = "Arial"
    p_flow.font.bold = True
    p_flow.font.size = Pt(7.5)
    p_flow.font.color.rgb = C_NAVY_HERO

    # =========================================================================
    # SLIDE 5: END-TO-END TECHNICAL ARCHITECTURE (HOW?)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide5, 5, 10, "System Architecture", 
                               "“ONE PIPELINE. ONE WORLD MODEL.”", 
                               "End-to-end layered architecture: Ingestion ➔ Normalization ➔ Dual Digital Twin ➔ Computing Cores ➔ Safety Circuit Breaker.")
    print("Formatting Slide 5: Technical Architecture...")

    # Layer 1: Ingestion Sources
    _, tf_l1 = create_card(slide5, Inches(0.6), Inches(1.30), Inches(12.13), Inches(0.85),
                           bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_l1_h = tf_l1.paragraphs[0]
    p_l1_h.text = "TIER 1: MULTI-SOURCE TELEMETRY INGESTION LAYER"
    p_l1_h.font.name = "Arial"
    p_l1_h.font.bold = True
    p_l1_h.font.size = Pt(8.5)
    p_l1_h.font.color.rgb = C_BLUE_PRIMARY
    
    p_l1_b = tf_l1.add_paragraph()
    p_l1_b.text = "📡 INCOIS ERDDAP (Wave Spectra)  |  📡 IMD Alerts (Cyclones)  |  📡 Copernicus CMEMS (0.083° Physical Grids)  |  📡 ISRO MOSDAC (SST)  |  📡 AIS Feeds  |  📡 NMEA Sensors"
    p_l1_b.font.name = "Arial"
    p_l1_b.font.size = Pt(8)
    p_l1_b.font.color.rgb = C_SLATE_900

    # Layer 2: Normalization & Dual Digital State
    _, tf_l2 = create_card(slide5, Inches(0.6), Inches(2.22), Inches(12.13), Inches(0.95),
                           bg_color=C_SLATE_50, border_color=C_BLUE_BORDER, border_width=1.0)
    p_l2_h = tf_l2.paragraphs[0]
    p_l2_h.text = "TIER 2: DATA NORMALIZATION & DUAL DIGITAL STATE"
    p_l2_h.font.name = "Arial"
    p_l2_h.font.bold = True
    p_l2_h.font.size = Pt(8.5)
    p_l2_h.font.color.rgb = C_NAVY_HERO
    
    p_l2_b = tf_l2.add_paragraph()
    p_l2_b.text = "• Environmental Digital State: Spatial H3 Indexing (0.083° / ~9km) + Temporal Alignment + Sensor Calibration (<50ms)\n• Parametric Vessel Digital Twin: Dynamic hydrodynamics model incorporating Hull Length (Lwl), Draft (Td), Beam (B), Freeboard & Engine Power"
    p_l2_b.font.name = "Arial"
    p_l2_b.font.size = Pt(7.5)
    p_l2_b.font.color.rgb = C_SLATE_700

    # Layer 3: Tripartite Computing Engines
    _, tf_l3 = create_card(slide5, Inches(0.6), Inches(3.24), Inches(12.13), Inches(1.45),
                           bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.0)
    p_l3_h = tf_l3.paragraphs[0]
    p_l3_h.text = "TIER 3: TRIPARTITE COMPUTING CORE (STRICT ARCHITECTURAL SEPARATION)"
    p_l3_h.font.name = "Arial"
    p_l3_h.font.bold = True
    p_l3_h.font.size = Pt(9)
    p_l3_h.font.color.rgb = C_CYAN_ACCENT
    p_l3_h.space_after = Pt(2)

    p_l3_b1 = tf_l3.add_paragraph()
    p_l3_b1.text = "⚙️ DETERMINISTIC PHYSICS: Wave encounter angle math, Critical wave height (Hcrit = 0.6 · L), wave steepness ratio & roll resonance."
    p_l3_b1.font.name = "Arial"
    p_l3_b1.font.bold = True
    p_l3_b1.font.size = Pt(8)
    p_l3_b1.font.color.rgb = C_EMERALD

    p_l3_b2 = tf_l3.add_paragraph()
    p_l3_b2.text = "🧠 PREDICTIVE ML LAYER: Multi-Species Habitat Suitability Index (HSI XGBoost), 2D Sobel thermal fronts & 1,000-particle SAR drift."
    p_l3_b2.font.name = "Arial"
    p_l3_b2.font.bold = True
    p_l3_b2.font.size = Pt(8)
    p_l3_b2.font.color.rgb = RGBColor(56, 189, 248)

    p_l3_b3 = tf_l3.add_paragraph()
    p_l3_b3.text = "📍 DYNAMIC OPTIMIZATION: Modified A* navigational pathfinder navigating through time-varying 2D risk-cost surface meshes."
    p_l3_b3.font.name = "Arial"
    p_l3_b3.font.bold = True
    p_l3_b3.font.size = Pt(8)
    p_l3_b3.font.color.rgb = C_AMBER

    # Layer 4: Hard Safety Circuit Breaker Shield
    _, tf_l4 = create_card(slide5, Inches(0.6), Inches(4.76), Inches(12.13), Inches(0.85),
                           bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.5)
    p_l4_h = tf_l4.paragraphs[0]
    p_l4_h.text = "TIER 4: NON-BYPASSABLE SAFETY CIRCUIT BREAKER (<10ms Pure Python, Zero Hallucination)"
    p_l4_h.font.name = "Arial"
    p_l4_h.font.bold = True
    p_l4_h.font.size = Pt(8.5)
    p_l4_h.font.color.rgb = C_EMERALD
    
    p_l4_b = tf_l4.add_paragraph()
    p_l4_b.text = "All physical safety checks and active cyclone overrides execute as deterministic compiled code. LLM models sit strictly downstream for plain-language vernacular voice translation only."
    p_l4_b.font.name = "Arial"
    p_l4_b.font.size = Pt(7.5)
    p_l4_b.font.color.rgb = C_SLATE_900

    # Layer 5: Operational Deliverables
    _, tf_l5 = create_card(slide5, Inches(0.6), Inches(5.68), Inches(12.13), Inches(1.10),
                           bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_l5_h = tf_l5.paragraphs[0]
    p_l5_h.text = "TIER 5: MULTI-CHANNEL OPERATIONAL DISPATCH"
    p_l5_h.font.name = "Arial"
    p_l5_h.font.bold = True
    p_l5_h.font.size = Pt(8.5)
    p_l5_h.font.color.rgb = C_NAVY_HERO
    
    p_l5_b = tf_l5.add_paragraph()
    p_l5_b.text = "• Fisherman PWA & LoRa Voice Device: Actionable safety verdict dial, optimized 4-point waypoints, offline voice guidance.\n• Fleet Web & Coast Guard Command: Real-time dark-fleet AIS anomaly detection, harbor fleet congestion & SAR drift tracking."
    p_l5_b.font.name = "Arial"
    p_l5_b.font.size = Pt(7.5)
    p_l5_b.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 6: INTELLIGENCE ENGINE (HOW INTELLIGENT?)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide6, 6, 10, "Intelligence Separation", 
                               "“PHYSICS + AI + OPTIMIZATION. EACH DOES WHAT IT IS GOOD AT.”", 
                               "Strict architectural boundary: Deterministic physics guarantees safety; ML discovers patterns; LLM translates language.")
    print("Formatting Slide 6: Intelligence Engine...")

    tri_data = [
        ("⚙️ DETERMINISTIC PHYSICS", 
         "Non-Bypassable Safety Law", 
         [("Critical Wave Limit", "Hcrit ≈ 0.6 · Lhull threshold prevents capsizing"),
          ("Wave Steepness Ratio", "S = Hs / λ calculated dynamically from wave spectra"),
          ("Beam-Sea Encounter", "Identifies roll resonance & parametric broaching risk"),
          ("Geofenced Boundaries", "Hard stop on Marine Protected Areas & international borders")],
         C_EMERALD_BG, C_EMERALD),
        ("🧠 PREDICTIVE MACHINE LEARNING", 
         "Pattern Recognition & Analytics", 
         [("Multi-Species HSI", "XGBoost model predicts Bangda, Surmai & Pomfret zones"),
          ("Sobel Thermal Fronts", "2D spatial filtering identifies productive SST boundaries"),
          ("SAR Anomaly Matching", "Flags dark vessels turning off AIS transponders"),
          ("Monte Carlo SAR Drift", "1,000-particle stochastic simulation for rescue targets")],
         C_BLUE_LIGHT, C_BLUE_PRIMARY),
        ("📍 GRAPH OPTIMIZATION (A*)", 
         "Navigational Pathfinding", 
         [("2D Risk Cost Surface", "Dynamic mesh where hazardous sea cells have infinite cost"),
          ("4-Point Waypoints", "Generates fuel-optimal path circumnavigating storm pockets"),
          ("Current Assist Optimization", "Leverages ocean surface drift velocity to save diesel"),
          ("Harbor Dock Allocation", "Optimizes arrival harbor based on catch ROI & congestion")],
         C_AMBER_BG, C_AMBER)
    ]

    for idx, (t_title, t_sub, t_items, bg_c, b_col) in enumerate(tri_data):
        _, tf_tri = create_card(slide6, Inches(0.6 + idx * 4.14), Inches(1.30), Inches(3.85), Inches(4.35),
                                bg_color=bg_c, border_color=b_col, border_width=1.0)
        p_th = tf_tri.paragraphs[0]
        p_th.text = t_title
        p_th.font.name = "Arial"
        p_th.font.bold = True
        p_th.font.size = Pt(9.5)
        p_th.font.color.rgb = b_col
        p_th.space_after = Pt(2)

        p_ts = tf_tri.add_paragraph()
        p_ts.text = t_sub
        p_ts.font.name = "Arial"
        p_ts.font.italic = True
        p_ts.font.size = Pt(7.5)
        p_ts.font.color.rgb = C_MUTED
        p_ts.space_after = Pt(6)

        for iname, idesc in t_items:
            p_i = tf_tri.add_paragraph()
            p_i.space_after = Pt(3.5)
            r1 = p_i.add_run()
            r1.text = f"• {iname}: "
            r1.font.name = "Arial"
            r1.font.bold = True
            r1.font.size = Pt(8)
            r1.font.color.rgb = C_SLATE_900
            r2 = p_i.add_run()
            r2.text = idesc
            r2.font.name = "Arial"
            r2.font.size = Pt(7.5)
            r2.font.color.rgb = C_SLATE_700

    # Bottom Banner: The Constrained LLM Boundary
    _, tf_llm = create_card(slide6, Inches(0.6), Inches(5.72), Inches(12.13), Inches(1.10),
                            bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.0)
    p_llm_h = tf_llm.paragraphs[0]
    p_llm_h.text = "THE LLM SAFETY GOVERNANCE DIRECTIVE"
    p_llm_h.font.name = "Arial"
    p_llm_h.font.bold = True
    p_llm_h.font.size = Pt(9)
    p_llm_h.font.color.rgb = C_CYAN_ACCENT
    p_llm_h.space_after = Pt(2)

    p_llm_b = tf_llm.add_paragraph()
    r_l1 = p_llm_b.add_run()
    r_l1.text = "LLM PERMITTED ROLES: "
    r_l1.font.name = "Arial"
    r_l1.font.bold = True
    r_l1.font.size = Pt(7.5)
    r_l1.font.color.rgb = C_EMERALD
    r_l2 = p_llm_b.add_run()
    r_l2.text = "Voice intent parsing (Speech ➔ JSON) and multilingual explanation synthesis (JSON ➔ Vernacular Audio).\n"
    r_l2.font.name = "Arial"
    r_l2.font.size = Pt(7.5)
    r_l2.font.color.rgb = C_WHITE

    r_l3 = p_llm_b.add_run()
    r_l3.text = "LLM STRICT PROHIBITIONS: "
    r_l3.font.name = "Arial"
    r_l3.font.bold = True
    r_l3.font.size = Pt(7.5)
    r_l3.font.color.rgb = C_ROSE
    r_l4 = p_llm_b.add_run()
    r_l4.text = "LLM CANNOT calculate safety, override physics circuit breakers, invent maritime data, or modify A* waypoints."
    r_l4.font.name = "Arial"
    r_l4.font.size = Pt(7.5)
    r_l4.font.color.rgb = C_WHITE

    # =========================================================================
    # SLIDE 7: LIVE WORKING PROTOTYPE (DOES IT WORK?)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide7, 7, 10, "Prototype Validation", 
                               "“FROM LOCATION TO DECISION: LIVE SIMULATION PIPELINE”", 
                               "End-to-end execution simulation demonstrating how ORCA processes a 9m FRP vessel navigating out of Kanyakumari harbor.")
    print("Formatting Slide 7: Live Working Prototype...")

    # Left: 6-Stage Visual Execution Pipeline (width: 7.2)
    _, tf_s7_p = create_card(slide7, Inches(0.6), Inches(1.30), Inches(7.2), Inches(5.55),
                             bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_s7_ph = tf_s7_p.paragraphs[0]
    p_s7_ph.text = "6-STAGE PROTOTYPE EXECUTION FLOW"
    p_s7_ph.font.name = "Arial"
    p_s7_ph.font.bold = True
    p_s7_ph.font.size = Pt(9.5)
    p_s7_ph.font.color.rgb = C_NAVY_HERO
    p_s7_ph.space_after = Pt(4)

    pipe_stages = [
        ("STAGE 01: GPS VOYAGE QUERY", "POST /api/v1/assess-trip (Lat: 8.08°N, Lon: 77.55°E, Kanyakumari Port)", C_BLUE_PRIMARY),
        ("STAGE 02: LIVE OCEAN INGESTION", "Open-Meteo & IMD Feeds: Wave Hs = 2.8m, Swell Tp = 8.2s, Wind Gusts = 24 kts", C_NAVY_HERO),
        ("STAGE 03: VESSEL TWIN INSTANTIATION", "Parametric Profile: 9m FRP artisanal gillnetter, 1.1m draft, 2.4m beam", C_NAVY_HERO),
        ("STAGE 04: DETERMINISTIC PHYSICS CHECK", "Hcrit = 0.6 · 9m = 5.4m > 2.8m (Safe from immediate open-water capsize)", C_EMERALD),
        ("STAGE 05: A* ROUTE RISK DIVERSION", "Detects shallow shoal wave breaker on direct path ➔ Diverts route 3.2 NM East", C_AMBER),
        ("STAGE 06: MULTILINGUAL VERDICT & AUDIO", "Verdict: PROCEED WITH CAUTION  ·  Synthesizes Tamil Voice Advisory & Map", C_EMERALD)
    ]
    for st, sd, sc in pipe_stages:
        p_stg = tf_s7_p.add_paragraph()
        p_stg.space_after = Pt(3.5)
        r1 = p_stg.add_run()
        r1.text = f"{st}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = sc
        r2 = p_stg.add_run()
        r2.text = f"  {sd}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Right Top: Map Navigation Simulation Diagram
    _, tf_s7_map = create_card(slide7, Inches(8.0), Inches(1.30), Inches(4.73), Inches(3.20),
                               bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER, border_width=1.0)
    p_s7_mh = tf_s7_map.paragraphs[0]
    p_s7_mh.text = "MARITIME COMMAND MAP: DYNAMIC ROUTE"
    p_s7_mh.font.name = "Arial"
    p_s7_mh.font.bold = True
    p_s7_mh.font.size = Pt(9)
    p_s7_mh.font.color.rgb = C_BLUE_PRIMARY
    p_s7_mh.space_after = Pt(3)

    p_map_v = tf_s7_map.add_paragraph()
    p_map_v.text = "📍 KANYAKUMARI HARBOR (Start Point)\n"
    p_map_v.font.name = "Arial"
    p_map_v.font.bold = True
    p_map_v.font.size = Pt(8)
    p_map_v.font.color.rgb = C_NAVY_HERO

    p_map_d = tf_s7_map.add_paragraph()
    p_map_d.text = "   |─── ❌ Direct Path: Crosses High-Risk Shoal Breaker Zone (Risk Score: 0.88)\n   └─── ✅ ORCA A* Diversion: +3.2 NM East (Risk Score: 0.24, Fuel Est: 14.2 L)\n"
    p_map_d.font.name = "Arial"
    p_map_d.font.size = Pt(7.5)
    p_map_d.font.color.rgb = C_SLATE_700

    p_map_e = tf_s7_map.add_paragraph()
    p_map_e.text = "🏁 PELAGIC FISHING GROUND (Target Zone: HSI Score 0.82)"
    p_map_e.font.name = "Arial"
    p_map_e.font.bold = True
    p_map_e.font.size = Pt(8)
    p_map_e.font.color.rgb = C_EMERALD

    # Right Bottom: Live Telemetry Metric Badges
    _, tf_s7_met = create_card(slide7, Inches(8.0), Inches(4.65), Inches(4.73), Inches(2.20),
                               bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_s7_th = tf_s7_met.paragraphs[0]
    p_s7_th.text = "BENCHMARK PERFORMANCE (SIMULATION)"
    p_s7_th.font.name = "Arial"
    p_s7_th.font.bold = True
    p_s7_th.font.size = Pt(9)
    p_s7_th.font.color.rgb = C_NAVY_HERO
    p_s7_th.space_after = Pt(2)

    metrics_s7 = [
        ("DAG Pipeline Execution Latency", "< 85ms (FastAPI + SQLite WAL)"),
        ("Safety Circuit Breaker Eval", "< 8ms (Compiled Deterministic Pure Python)"),
        ("A* Pathfinder Grid Compute", "< 35ms (100x100 Spatial Cost Surface)"),
        ("Offline Voice Audio Generation", "Pre-cached vernacular phoneme synthesis")
    ]
    for mn, mv in metrics_s7:
        p_m = tf_s7_met.add_paragraph()
        p_m.space_after = Pt(1.5)
        r1 = p_m.add_run()
        r1.text = f"• {mn}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p_m.add_run()
        r2.text = mv
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 8: SAFETY + FAILURE RESILIENCE (WHAT IF THINGS FAIL?)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide8, 8, 10, "Resilience & Failsafes", 
                               "“DESIGNED FOR WHEN THE SYSTEM IS WRONG”", 
                               "Maritime environments are hostile and unpredictable. ORCA assumes data will be stale, APIs will fail, and signals will drop.")
    print("Formatting Slide 8: Safety & Failure Resilience...")

    # Left: 6 Real-World Failure Modes (width: 5.6)
    _, tf_s8_l = create_card(slide8, Inches(0.6), Inches(1.30), Inches(5.6), Inches(5.55),
                             bg_color=C_ROSE_BG, border_color=C_ROSE, border_width=1.0)
    p_s8_lh = tf_s8_l.paragraphs[0]
    p_s8_lh.text = "HOSTILE SEA FAILURE CONDITIONS"
    p_s8_lh.font.name = "Arial"
    p_s8_lh.font.bold = True
    p_s8_lh.font.size = Pt(10)
    p_s8_lh.font.color.rgb = C_ROSE
    p_s8_lh.space_after = Pt(4)

    fail_modes = [
        ("1. Data Stale / Cache Expired", "Satellite pass delayed or forecast timestamp older than 3 hours."),
        ("2. API Provider Downtime", "INCOIS or Open-Meteo external REST gateway timeout or HTTP 500."),
        ("3. Source Disagreement", "IMD predicts 18 kt wind while ECMWF model indicates 32 kt gusts."),
        ("4. Offshore Signal Blackout", "Vessel sails 15 NM offshore, 4G/5G mobile signal drops to 0%."),
        ("5. Onboard Sensor Failure", "GPS or pitch-roll IMU disconnects or outputs erratic noisy data."),
        ("6. ML Uncertainty Spikes", "Habitat model or drift simulator encounters out-of-distribution state.")
    ]
    for fn, fd in fail_modes:
        p_f = tf_s8_l.add_paragraph()
        p_f.space_after = Pt(3.5)
        r1 = p_f.add_run()
        r1.text = f"⚠️ {fn}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_SLATE_900
        r2 = p_f.add_run()
        r2.text = f"  {fd}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Right: ORCA Engineering Mitigations (width: 6.33)
    _, tf_s8_r = create_card(slide8, Inches(6.4), Inches(1.30), Inches(6.33), Inches(5.55),
                             bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.5)
    p_s8_rh = tf_s8_r.paragraphs[0]
    p_s8_rh.text = "ORCA RESILIENCE & SAFETY ARCHITECTURE"
    p_s8_rh.font.name = "Arial"
    p_s8_rh.font.bold = True
    p_s8_rh.font.size = Pt(10)
    p_s8_rh.font.color.rgb = C_EMERALD
    p_s8_rh.space_after = Pt(4)

    mitigations = [
        ("🛡️ Automated TTL Freshness Monitoring", "Flags data older than 30 min; falls back to conservative historical wave bounds."),
        ("🛡️ Multi-Source Cascade Circuit Breaker", "Automated failover: INCOIS ERDDAP ➔ Open-Meteo Marine ➔ MOSDAC Backup."),
        ("🛡️ Provenance & Confidence Weighting", "Calculates ensemble disagreement bounds; defaults to highest-risk scenario."),
        ("🛡️ Multi-Hop LoRa Fleet Mesh (OLSR)", "Vessels relay 16-byte emergency packets across 3-8 km hops to 60m lighthouse masts."),
        ("🛡️ Kalman Filter Sensor Sanity Gate", "Rejects IMU spikes > 3 sigma; defaults to vessel class static stability envelope."),
        ("🛡️ 100% Deterministic Safety Override", "Zero AI in safety gate; pure compiled physical logic enforces non-negotiable stop.")
    ]
    for mn, md in mitigations:
        p_m = tf_s8_r.add_paragraph()
        p_m.space_after = Pt(3.5)
        r1 = p_m.add_run()
        r1.text = f"{mn}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p_m.add_run()
        r2.text = f"  {md}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 9: FEASIBILITY + DEPLOYMENT + ECONOMICS (CAN IT SCALE?)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide9, 9, 10, "Execution & Scale", 
                               "“FROM PROTOTYPE TO COASTAL INFRASTRUCTURE”", 
                               "Verified engineering status, multi-hop LoRa hardware mesh, and low-cost deployment economics.")
    print("Formatting Slide 9: Feasibility, Deployment & Economics...")

    # Zone 1: Verified Roadmap
    _, tf_s9_z1 = create_card(slide9, Inches(0.6), Inches(1.30), Inches(3.6), Inches(5.55),
                              bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_z1_h = tf_s9_z1.paragraphs[0]
    p_z1_h.text = "ZONE 1: VERIFIED ROADMAP"
    p_z1_h.font.name = "Arial"
    p_z1_h.font.bold = True
    p_z1_h.font.size = Pt(9.5)
    p_z1_h.font.color.rgb = C_NAVY_HERO
    p_z1_h.space_after = Pt(4)

    milestones_s9 = [
        ("✓ END-TO-END PROTOTYPE", "FastAPI + React MapLibre working architecture [BUILT]"),
        ("✓ MULTI-SOURCE INGESTION", "Open-Meteo, IMD, MOSDAC pipelines active [BUILT]"),
        ("✓ DETERMINISTIC RISK ENGINE", "Physical capsizing & Hcrit circuit breaker [BUILT]"),
        ("✓ DYNAMIC A* PATHFINDER", "Time-varying 2D risk surface routing [BUILT]"),
        ("◐ HARBOR PILOT TRIAL", "Kanyakumari & Munambam field trials [IN PROGRESS]"),
        ("◌ PAN-INDIA EXPANSION", "Integration across 9 coastal states & Satcom [PLANNED]")
    ]
    for mt, md in milestones_s9:
        p_m = tf_s9_z1.add_paragraph()
        p_m.space_after = Pt(3.5)
        r1 = p_m.add_run()
        r1.text = f"{mt}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_EMERALD if "✓" in mt else (C_AMBER if "◐" in mt else C_MUTED)
        r2 = p_m.add_run()
        r2.text = md
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Zone 2: Mesh Deployment Architecture
    _, tf_s9_z2 = create_card(slide9, Inches(4.4), Inches(1.30), Inches(4.2), Inches(5.55),
                              bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_z2_h = tf_s9_z2.paragraphs[0]
    p_z2_h.text = "ZONE 2: DEPLOYMENT ARCHITECTURE"
    p_z2_h.font.name = "Arial"
    p_z2_h.font.bold = True
    p_z2_h.font.size = Pt(9.5)
    p_z2_h.font.color.rgb = C_BLUE_PRIMARY
    p_z2_h.space_after = Pt(4)

    dep_layers = [
        ("1. Vessel Edge Node", "Low-power ESP32 MCU + LoRa SX1262 transceiver + GPS NEO-6M."),
        ("2. Fleet Ad-Hoc Relay", "Boats spaced 3–8 km apart form a dynamic peer-to-peer radio mesh."),
        ("3. Coastal Lighthouse Mast", "60m lighthouse receiver with line-of-sight reach of 33.8 km to boats."),
        ("4. ORCA Fast Cloud Engine", "Ingests 16-byte binary packets; executes A* & HSI in <85ms."),
        ("5. Dual User Interface", "Offline PWA for fishers; real-time dashboard for Coast Guard.")
    ]
    for dt, dd in dep_layers:
        p_d = tf_s9_z2.add_paragraph()
        p_d.space_after = Pt(4)
        r1 = p_d.add_run()
        r1.text = f"{dt}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p_d.add_run()
        r2.text = dd
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Zone 3: BOM Economics & Scaling
    _, tf_s9_z3 = create_card(slide9, Inches(8.8), Inches(1.30), Inches(3.93), Inches(5.55),
                              bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.0)
    p_z3_h = tf_s9_z3.paragraphs[0]
    p_z3_h.text = "ZONE 3: ECONOMICS & BOM"
    p_z3_h.font.name = "Arial"
    p_z3_h.font.bold = True
    p_z3_h.font.size = Pt(9.5)
    p_z3_h.font.color.rgb = C_EMERALD
    p_z3_h.space_after = Pt(4)

    bom_items = [
        ("ESP32 MCU + LoRa SX1262", "₹1,800 (~$22)"),
        ("GPS NEO-6M Receiver", "₹600 (~$7)"),
        ("IP67 Enclosure + Solar Panel", "₹1,100 (~$13)"),
        ("Total Hardware Cost / Boat", "≈ ₹3,500 ($42) [INDICATIVE]"),
        ("Public Data Ingestion Cost", "₹0 (Open INCOIS & MOSDAC)"),
        ("Serverless Compute / Trip", "₹0.04 per assessment [ESTIMATE]")
    ]
    for bn, bv in bom_items:
        p_b = tf_s9_z3.add_paragraph()
        p_b.space_after = Pt(2.5)
        r1 = p_b.add_run()
        r1.text = f"• {bn}: "
        r1.font.name = "Arial"
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_SLATE_900
        r2 = p_b.add_run()
        r2.text = bv
        r2.font.name = "Arial"
        r2.font.bold = True
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_EMERALD if "Total" in bn or "₹0" in bv else C_BLUE_PRIMARY

    p_scl = tf_s9_z3.add_paragraph()
    p_scl.space_after = Pt(2)
    r_scl_h = p_scl.add_run()
    r_scl_h.text = "\nSCALE PROGRESSION:\n"
    r_scl_h.font.name = "Arial"
    r_scl_h.font.bold = True
    r_scl_h.font.size = Pt(8)
    r_scl_h.font.color.rgb = C_NAVY_HERO
    r_scl_b = p_scl.add_run()
    r_scl_b.text = "1 Vessel ➔ Harbor Fleet (500+) ➔ Coastal State ➔ Pan-India Coast (7,516 km)"
    r_scl_b.font.name = "Arial"
    r_scl_b.font.size = Pt(7.5)
    r_scl_b.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 10: IMPACT + RESEARCH + FINAL VISION (WHY DOES IT MATTER?)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_official_slide_frame(slide10, 10, 10, "Impact & Vision", 
                               "“ONE VESSEL. ONE HARBOR. ONE NATIONAL MARITIME INTELLIGENCE LAYER.”", 
                               "Transforming academic oceanographic research and space telemetry into sovereign digital infrastructure for India's Blue Economy.")
    print("Formatting Slide 10: Impact, Research & Final Vision...")

    # Top: 4 Grounded Research Blueprint Cards (width: 2.9 each)
    research_s10 = [
        ("IMO Code on Intact Stability (2008)", "Res. MSC.267(85)", "Hcrit ≈ 0.6·Lhull capsizing threshold", C_BLUE_LIGHT, C_BLUE_PRIMARY),
        ("Wave Hydrodynamics (Faltinsen)", "Encounter Physics", "Wave steepness (S = Hs/λ) & roll resonance", C_EMERALD_BG, C_EMERALD),
        ("Satellite Oceanography (INCOIS)", "MOSDAC Telemetry", "Sobel thermal fronts & Multi-species HSI", C_AMBER_BG, C_AMBER),
        ("Dynamic Pathfinding (Modified A*)", "Risk Cost Surfaces", "4-Point waypoint storm avoidance routing", C_BLUE_LIGHT, C_BLUE_PRIMARY)
    ]
    for idx, (rt, rsrc, rapp, bg_c, b_col) in enumerate(research_s10):
        _, tf_r = create_card(slide10, Inches(0.6 + idx * 3.08), Inches(1.30), Inches(2.9), Inches(1.85),
                              bg_color=bg_c, border_color=b_col, border_width=1.0)
        p1 = tf_r.paragraphs[0]
        p1.text = f"📚 {rt}"
        p1.font.name = "Arial"
        p1.font.bold = True
        p1.font.size = Pt(8)
        p1.font.color.rgb = b_col
        
        p2 = tf_r.add_paragraph()
        p2.text = f"Source: {rsrc}"
        p2.font.name = "Arial"
        p2.font.italic = True
        p2.font.size = Pt(7)
        p2.font.color.rgb = C_MUTED
        p2.space_after = Pt(2)

        p3 = tf_r.add_paragraph()
        r1 = p3.add_run()
        r1.text = "ORCA Implementation: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p3.add_run()
        r2.text = rapp
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Center: 5-Stakeholder Multiplier Ring
    _, tf_stk = create_card(slide10, Inches(0.6), Inches(3.24), Inches(12.13), Inches(1.65),
                            bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_stk_h = tf_stk.paragraphs[0]
    p_stk_h.text = "NATIONAL BLUE ECONOMY STAKEHOLDER MULTIPLIER"
    p_stk_h.font.name = "Arial"
    p_stk_h.font.bold = True
    p_stk_h.font.size = Pt(9.5)
    p_stk_h.font.color.rgb = C_NAVY_HERO
    p_stk_h.space_after = Pt(3)

    stk_items = [
        ("🎣 4.5M Artisanal Fishers", "Zero preventable drownings, timely voice warnings in native dialect, precise fishing spots."),
        ("⚓ Harbors & Port Authorities", "Automated safety clearance, harbor fleet congestion management & departure scheduling."),
        ("🚢 Commercial Trawlers", "20–30% fuel savings via current-assisted A* routing; higher pelagic catch ROI."),
        ("🛡️ Coast Guard & NDRF", "Instant Dark-Fleet anomaly matching & 1,000-particle Monte Carlo SAR drift tracking."),
        ("🔬 Oceanographic Scientists", "Crowdsourced ground-truth sea state telemetry across India's 2.37M sq km EEZ.")
    ]
    for s_title, s_desc in stk_items:
        p_s = tf_stk.add_paragraph()
        p_s.space_after = Pt(2)
        r1 = p_s.add_run()
        r1.text = f"{s_title}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p_s.add_run()
        r2.text = s_desc
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Bottom: Final Heroic Vision Statement
    _, tf_vis = create_card(slide10, Inches(0.6), Inches(4.98), Inches(12.13), Inches(1.85),
                            bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.5)
    p_vis_h = tf_vis.paragraphs[0]
    p_vis_h.text = "“ORCA DOESN’T JUST PREDICT THE OCEAN. IT HELPS PEOPLE DECIDE WHAT TO DO NEXT.”"
    p_vis_h.font.name = "Arial"
    p_vis_h.font.size = Pt(12)
    p_vis_h.font.bold = True
    p_vis_h.font.color.rgb = C_WHITE
    p_vis_h.alignment = PP_ALIGN.CENTER
    p_vis_h.space_after = Pt(3)

    p_vis_sub = tf_vis.add_paragraph()
    p_vis_sub.text = "From a single artisanal fishing boat in Kanyakumari to a sovereign digital intelligence layer for India’s entire 7,516 km coastline."
    p_vis_sub.font.name = "Arial"
    p_vis_sub.font.size = Pt(9)
    p_vis_sub.font.color.rgb = RGBColor(203, 213, 225)
    p_vis_sub.alignment = PP_ALIGN.CENTER
    p_vis_sub.space_after = Pt(5)

    p_vis_flow = tf_vis.add_paragraph()
    p_vis_flow.text = "OCEAN  ➔  WEATHER  ➔  VESSEL  ➔  ROUTE  ➔  ORCA  ➔  RISK  ➔  DECISION"
    p_vis_flow.font.name = "Arial"
    p_vis_flow.font.size = Pt(8.5)
    p_vis_flow.font.bold = True
    p_vis_flow.font.color.rgb = C_CYAN_ACCENT
    p_vis_flow.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Official 10-Slide Template Presentation generated successfully!")
    print(f"Saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_presentation()
