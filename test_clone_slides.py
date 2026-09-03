import pptx
from pptx import Presentation

def duplicate_slide(prs, index):
    # Clones slide at index
    template_slide = prs.slides[index]
    slide_layout = prs.slide_layouts[1]
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy shapes from template_slide
    for shp in template_slide.shapes:
        if shp.name.startswith("Title") or shp.name.startswith("TextBox"):
            continue
        # We can replicate other shapes like pictures, ovals, etc.
    return new_slide

prs = Presentation('/Users/subham/Downloads/test_6slides.pptx')
print("Initial slides:", len(prs.slides))
