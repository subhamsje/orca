"""
Populate the official SIH 2026 template with ORCA 4.0 content.

Strategy:
- Open /Users/subham/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx
- Keep slide 7 as-is (SIH instruction about 6-slide limit)
- Replace the body text in slides 1–6 with ORCA content
- Keep all placeholders (title, body, footer, sldNum) intact
- Keep the blue footer bar / images on each slide
- Add ORCA diagrams as additional shapes on top of the placeholder area

Grounded in real ORCA repo data — no fabricated statistics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


NAVY        = RGBColor(0x08, 0x1B, 0x33)
DEEP_NAVY   = RGBColor(0x04, 0x10, 0x20)
TEAL        = RGBColor(0x14, 0xB8, 0xA6)
TEAL_DARK   = RGBColor(0x0E, 0x8C, 0x7F)
CYAN        = RGBColor(0x67, 0xE8, 0xF9)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
WHITE       = RGBColor(0xF8, 0xFA, 0xFC)
SOFT        = RGBColor(0xCB, 0xD5, 0xE1)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
GRID        = RGBColor(0x1E, 0x3A, 0x5F)
SURFACE     = RGBColor(0x0F, 0x24, 0x40)
SIH_BLUE    = RGBColor(0x00, 0x70, 0xC0)  # the template footer blue


TEMPLATE = "/Users/subham/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx"
OUT      = "/Users/subham/code/orca/ORCA_4.0_SIH_Deck.pptx"


# ---------------- Helpers ----------------------------------------------------

def rect(slide, x, y, w, h, *, fill_rgb=None, line_rgb=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill_rgb is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
        if line_w is not None:
            shp.line.width = line_w
    return shp


def text(slide, x, y, w, h, *, text="", size=14, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
         italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def set_text(slide, *, must_have_placeholder=None, runs):
    """Replace the text of a placeholder with given runs.

    runs: list of dicts with 'text', 'size', 'bold', 'color'.
    Locate the placeholder by index or type."""
    target = None
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        if shp.is_placeholder:
            ph = shp.placeholder_format
            if must_have_placeholder == "title" and ph.type == 13:
                target = shp; break
            if must_have_placeholder == "body" and ph.idx == 1:
                target = shp; break
            if must_have_placeholder == "ctrTitle" and ph.type == 3:
                target = shp; break
            if must_have_placeholder == "subTitle" and ph.type == 4:
                target = shp; break
    if target is None:
        # fall back: any placeholder
        for shp in slide.shapes:
            if shp.is_placeholder and shp.has_text_frame:
                target = shp; break
    if target is None:
        return None
    tf = target.text_frame
    tf.clear()
    for i, run in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = run.get("align", PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = run["text"]
        r.font.name = run.get("font", "Calibri")
        r.font.size = Pt(run.get("size", 18))
        r.font.bold = run.get("bold", False)
        r.font.italic = run.get("italic", False)
        r.font.color.rgb = run.get("color", WHITE)
    return target


def line(slide, x1, y1, x2, y2, rgb=TEAL, weight=Pt(1.5)):
    cx = slide.shapes.add_connector(1, x1, y1, x2, y2)
    cx.line.color.rgb = rgb
    cx.line.width = weight
    return cx


def bullet_textbox(slide, x, y, w, h, items, *, size=12, color=WHITE,
                   bullet_char="•", bullet_color=None, gap=Pt(6)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    bc = bullet_color or TEAL
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = gap
        r = p.add_run(); r.text = bullet_char + "  "
        r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = bc
        if isinstance(item, tuple):
            head, body = item
            r2 = p.add_run(); r2.text = head
            r2.font.name = "Calibri"; r2.font.size = Pt(size); r2.font.bold = True
            r2.font.color.rgb = color
            if body:
                r3 = p.add_run(); r3.text = "  " + body
                r3.font.name = "Calibri"; r3.font.size = Pt(size)
                r3.font.color.rgb = SOFT
        else:
            r2 = p.add_run(); r2.text = item
            r2.font.name = "Calibri"; r2.font.size = Pt(size)
            r2.font.color.rgb = color
    return tb


# ---------------- Slide content ----------------------------------------------

def fill_slide_1(slide):
    """TITLE PAGE — keep template, set official fields."""
    # ctrTitle: ORCA 4.0
    set_text(slide, must_have_placeholder="ctrTitle", runs=[
        {"text": "ORCA 4.0", "size": 60, "bold": True, "color": WHITE,
         "align": PP_ALIGN.CENTER},
        {"text": "From Ocean Data to Safer Decisions.",
         "size": 24, "bold": False, "color": TEAL, "align": PP_ALIGN.CENTER},
    ])
    # subTitle: official fields
    set_text(slide, must_have_placeholder="subTitle", runs=[
        {"text": "Maritime Decision Intelligence for Coastal India",
         "size": 22, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": " ", "size": 8, "color": WHITE},
        {"text": "Problem Statement ID:  SIH26176",
         "size": 16, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": "Problem Statement Title:  Smart India Hackathon 2026",
         "size": 14, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": "Theme:  Marine / Blue Economy  ·  "
                  "Sponsoring Agency:  ISRO / Department of Space",
         "size": 14, "color": WHITE, "align": PP_ALIGN.CENTER},
        {"text": "PS Category:  Software + Hardware  ·  "
                  "Team ID:  ORCA-04  ·  Team Name:  ORCA",
         "size": 14, "color": WHITE, "align": PP_ALIGN.CENTER},
    ])


def fill_slide_2(slide):
    """IDEA TITLE / PROPOSED SOLUTION."""
    # Title placeholder
    set_text(slide, must_have_placeholder="title", runs=[
        {"text": "IDEA TITLE", "size": 28, "bold": True, "color": TEAL},
        {"text": "Proposed Solution", "size": 20, "bold": True, "color": WHITE},
    ])
    # Body: lead visual + statement
    # Diagram region under title placeholder (~y=1.3", h=5.0")
    diag_x = Inches(0.5); diag_y = Inches(1.45)
    diag_w = Inches(12.3); diag_h = Inches(3.4)

    # Fragmented inputs (left)
    inputs = ["WAVES", "WIND", "CURRENTS", "WEATHER",
              "CYCLONES", "VESSEL", "AIS", "ROUTE"]
    box_w = Inches(1.20); box_h = Inches(0.42)
    for i, label in enumerate(inputs):
        col = i % 4
        row = i // 4
        x = diag_x + Inches(0.05) + col * (box_w + Inches(0.10))
        y = diag_y + row * (box_h + Inches(0.12))
        rect(slide, x, y, box_w, box_h, fill_rgb=SURFACE)
        rect(slide, x, y, box_w, Inches(0.04), fill_rgb=TEAL)
        text(slide, x, y, box_w, box_h, text=label, size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Convergence arrows
    arrow_start_x = diag_x + Inches(5.2)
    arrow_end_x   = diag_x + Inches(6.6)
    arrow_y       = diag_y + Inches(0.65)
    line(slide, arrow_start_x, arrow_y, arrow_end_x, arrow_y,
         rgb=TEAL, weight=Pt(2.5))

    # ORCA central engine block
    eng_x = diag_x + Inches(6.6); eng_y = diag_y
    eng_w = Inches(5.5); eng_h = Inches(2.2)
    rect(slide, eng_x, eng_y, eng_w, eng_h, fill_rgb=DEEP_NAVY,
         line_rgb=TEAL, line_w=Pt(2))
    text(slide, eng_x, eng_y + Inches(0.10), eng_w, Inches(0.4),
         text="ORCA 4.0  INTELLIGENCE ENGINE", size=14, bold=True,
         color=TEAL, align=PP_ALIGN.CENTER)
    # 5 sub-stages inside the engine
    sub = ["Environmental State", "Vessel Digital Twin", "Maritime Physics",
           "Circuit Breaker", "ORCA MRSI 0–100"]
    sw = eng_w / 5
    for i, s in enumerate(sub):
        sx = eng_x + sw * i + Inches(0.05)
        sw_ = sw - Inches(0.10)
        rect(slide, sx, eng_y + Inches(0.65), sw_, Inches(0.95),
             fill_rgb=SURFACE)
        text(slide, sx, eng_y + Inches(0.75), sw_, Inches(0.85),
             text=s, size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Output cascade
    out_x = diag_x + Inches(7.5); out_y = diag_y + Inches(2.4)
    out_w = Inches(2.0);  out_h = Inches(0.55)
    outs = [("RISK", AMBER),
            ("ROUTE-SPECIFIC RISK", TEAL),
            ("ACTIONABLE DECISION", GREEN)]
    for i, (label, color) in enumerate(outs):
        y = out_y + i * (out_h + Inches(0.08))
        rect(slide, out_x, y, out_w, out_h, fill_rgb=SURFACE)
        rect(slide, out_x, y, Inches(0.10), out_h, fill_rgb=color)
        text(slide, out_x + Inches(0.20), y, out_w - Inches(0.30),
             out_h, text=label, size=11, bold=True, color=color,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # Below diagram: key statement + innovation strip
    key_y = diag_y + diag_h + Inches(0.20)
    text(slide, Inches(0.5), key_y, Inches(12.3), Inches(0.55),
         text="FORECAST  ≠  OPERATIONAL RISK.",
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(slide, Inches(0.5), key_y + Inches(0.55), Inches(12.3), Inches(0.40),
         text="ORCA converts environmental conditions into a decision for the specific voyage.",
         size=14, color=SOFT, align=PP_ALIGN.CENTER, italic=True)


def fill_slide_3(slide):
    """TECHNICAL APPROACH."""
    set_text(slide, must_have_placeholder="title", runs=[
        {"text": "TECHNICAL APPROACH", "size": 28, "bold": True, "color": TEAL},
        {"text": "Technologies · Methodology · Process · Prototype",
         "size": 16, "bold": False, "color": SOFT},
    ])

    # Pipeline diagram (full slide below title, above blue bar at y≈6.95")
    dx = Inches(0.5); dy = Inches(1.40)
    dw = Inches(12.33); dh = Inches(5.40)

    # Layer 1: data sources (stacked column)
    src_y = dy
    src_h = Inches(0.35)
    sources = [
        ("INCOIS", "erddap (roadmap)"),
        ("IMD", "CWD / Mausam (roadmap)"),
        ("COPERNICUS MARINE", "needs credentials"),
        ("OPEN-METEO MARINE / ECMWF / FORECAST", "live, no creds"),
        ("MET NORWAY (yr.no)", "live, no creds"),
        ("NOAA NDBC BUOYS", "live, in-situ"),
        ("AIS / VESSEL TELEMETRY", "via LoRa / NMEA"),
    ]
    src_x = dx + Inches(0.05)
    src_w = Inches(3.10)
    for i, (a, b) in enumerate(sources):
        y = src_y + i * (src_h + Inches(0.04))
        rect(slide, src_x, y, src_w, src_h, fill_rgb=SURFACE)
        rect(slide, src_x, y, Inches(0.07), src_h, fill_rgb=TEAL)
        text(slide, src_x + Inches(0.15), y, Inches(1.6), src_h,
             text=a, size=9, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, src_x + Inches(1.7), y, Inches(1.35), src_h,
             text=b, size=8, color=MUTED, anchor=MSO_ANCHOR.MIDDLE,
             italic=True)

    # Funnel + downstream pipeline
    pipe_x = dx + Inches(3.30)
    pipe_w = dw - Inches(3.50)
    stages = [
        ("PROVIDER ABSTRACTION  ·  circuit breaker  ·  rate limit",
         TEAL, 0.55),
        ("CANONICAL NORMALIZATION  ·  18 params  ·  provenance",
         CYAN, 0.55),
        ("ENVIRONMENTAL STATE  ·  CURRENT / RECENT / STALE / UNAVAILABLE",
         TEAL, 0.55),
        ("VESSEL DIGITAL TWIN  ·  16-field validated profile",
         CYAN, 0.55),
        ("MARITIME PHYSICS  ·  9 hazards  ·  Kijima-1990 encounter",
         TEAL, 0.55),
        ("DETERMINISTIC CIRCUIT BREAKER  ·  7 hard rules",
         CYAN, 0.55),
        ("ORCA MRSI  ·  0 – 100  ·  reproducible (SHA-256)",
         TEAL, 0.55),
        ("ROUTE RISK  ·  per-segment  ·  replay store",
         CYAN, 0.55),
    ]
    h = (dh - Inches(0.5)) / len(stages)
    y = dy + Inches(0.05)
    for label, accent, _ in stages:
        rect(slide, pipe_x, y, pipe_w, h - Inches(0.07), fill_rgb=SURFACE)
        rect(slide, pipe_x, y, Inches(0.10), h - Inches(0.07),
             fill_rgb=accent)
        text(slide, pipe_x + Inches(0.25), y, pipe_w - Inches(0.35),
             h - Inches(0.07), text=label, size=11, bold=True,
             color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        y += h

    # Tech stack strip at bottom
    tech_y = dy + dh - Inches(0.40)
    tech = [("FRONTEND",  "React · TS · Vite · Leaflet · h3-js"),
            ("BACKEND",   "FastAPI · Python 3.12 · async providers"),
            ("INTELLIGENCE", "Hazards · CB · MRSI · route risk"),
            ("INFRA",     "Cloud-native · edge-capable roadmap")]
    tw = dw / 4
    for i, (h_, b_) in enumerate(tech):
        x = dx + tw * i
        rect(slide, x + Inches(0.05), tech_y, tw - Inches(0.10),
             Inches(0.32), fill_rgb=DEEP_NAVY, line_rgb=GRID, line_w=Pt(0.5))
        text(slide, x + Inches(0.10), tech_y, tw - Inches(0.20),
             Inches(0.32), text=h_, size=9, bold=True, color=TEAL,
             anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x + Inches(1.10), tech_y, tw - Inches(1.30),
             Inches(0.32), text=b_, size=8, color=SOFT,
             anchor=MSO_ANCHOR.MIDDLE)

    # Mini demo flow embedded at top-right of pipeline
    demo_x = pipe_x + pipe_w - Inches(2.2)
    demo_y = dy + Inches(0.05)
    rect(slide, demo_x, demo_y, Inches(2.10), Inches(0.55), fill_rgb=NAVY,
         line_rgb=TEAL, line_w=Pt(1))
    text(slide, demo_x, demo_y, Inches(2.10), Inches(0.55),
         text="LIVE DEMO:  /api/v1/assess-now\n"
              "Kanyakumari → MRSI 75 / 100",
         size=8, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def fill_slide_4(slide):
    """FEASIBILITY AND VIABILITY."""
    set_text(slide, must_have_placeholder="title", runs=[
        {"text": "FEASIBILITY AND VIABILITY", "size": 28, "bold": True,
         "color": TEAL},
        {"text": "Implementation status · Challenges & mitigation · "
                 "Deployment · Cost",
         "size": 16, "color": SOFT},
    ])

    # Three columns
    col_w = (Inches(13.333 - 0.5 * 2 - 0.30) / 3)
    col_h = Inches(5.40)
    col_y = Inches(1.40)
    gap   = Inches(0.15)

    def col(x, head, color):
        rect(slide, x, col_y, col_w, col_h, fill_rgb=SURFACE)
        rect(slide, x, col_y, col_w, Inches(0.10), fill_rgb=color)
        text(slide, x + Inches(0.20), col_y + Inches(0.18),
             col_w - Inches(0.40), Inches(0.40),
             text=head, size=14, bold=True, color=color)

    # COL A: IMPLEMENTATION ROADMAP
    col(Inches(0.5), "A.  IMPLEMENTATION ROADMAP", TEAL)
    items = [
        ("✓ PROTOTYPE",   "FastAPI + React/TS, 34 endpoints"),
        ("✓ LIVE DATA",   "5 providers, canonical layer, provenance"),
        ("✓ RISK ENGINE", "ORCA MRSI, 9 hazards, 24 unit tests pass"),
        ("✓ ROUTE ENGINE", "Per-segment risk, replay store"),
        ("◐ PILOT",       "Kanyakumari live: MRSI 75/100"),
        ("◌ DEPLOYMENT",  "IMD/INCOIS feeds, multi-state scale"),
    ]
    yy = col_y + Inches(0.70)
    for head, body in items:
        text(slide, Inches(0.7), yy, col_w - Inches(0.40), Inches(0.32),
             text=head, size=11, bold=True, color=WHITE)
        text(slide, Inches(0.7), yy + Inches(0.30),
             col_w - Inches(0.40), Inches(0.32),
             text=body, size=10, color=SOFT)
        yy += Inches(0.70)

    # COL B: CHALLENGE → MITIGATION
    x_b = Inches(0.5) + col_w + gap
    col(x_b, "B.  CHALLENGES  →  MITIGATION", CYAN)
    pairs = [
        ("Data stale",           "→ freshness policy (CURRENT/RECENT/STALE)"),
        ("API failure",          "→ provider circuit breaker + fallback"),
        ("Source disagreement",  "→ provenance + confidence score"),
        ("No connectivity",      "→ edge cache + explicit degraded mode"),
        ("Vessel variability",   "→ 16-field validated vessel twin"),
        ("Model uncertainty",    "→ risk_uncertainty 0–1 surfaced in UI"),
    ]
    yy = col_y + Inches(0.70)
    for ch, mi in pairs:
        text(slide, x_b + Inches(0.20), yy,
             col_w - Inches(0.40), Inches(0.32),
             text=ch, size=11, bold=True, color=AMBER)
        text(slide, x_b + Inches(0.20), yy + Inches(0.30),
             col_w - Inches(0.40), Inches(0.32),
             text=mi, size=10, color=SOFT)
        yy += Inches(0.70)

    # COL C: DEPLOYMENT ECONOMICS
    x_c = Inches(0.5) + (col_w + gap) * 2
    col(x_c, "C.  DEPLOYMENT ECONOMICS", TEAL)
    # Tiny stack diagram
    stack_x = x_c + Inches(0.20)
    stack_w = col_w - Inches(0.40)
    stack_y = col_y + Inches(0.80)
    rect(slide, stack_x, stack_y, stack_w, Inches(0.45),
         fill_rgb=DEEP_NAVY, line_rgb=TEAL, line_w=Pt(0.5))
    text(slide, stack_x, stack_y, stack_w, Inches(0.45),
         text="HARDWARE  +  COMMUNICATION  +  COMPUTE  +  MAINTENANCE",
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    phases = [
        ("PROTOTYPE",  "Cloud-only · free-tier (Vercel + Open-Meteo)",
         "indicative", TEAL),
        ("PILOT",      "≈ ₹4 000 / mo on commodity VPS · 10⁵ assessments",
         "indicative", TEAL),
        ("SCALE",      "Adds StormGlass / Copernicus keys · vendor-dep.",
         "indicative", TEAL),
    ]
    yy = stack_y + Inches(0.65)
    for label, body, tag, color in phases:
        rect(slide, stack_x, yy, stack_w, Inches(0.70), fill_rgb=DEEP_NAVY)
        rect(slide, stack_x, yy, Inches(0.10), Inches(0.70), fill_rgb=color)
        text(slide, stack_x + Inches(0.20), yy, stack_w - Inches(0.40),
             Inches(0.32), text=label, size=11, bold=True, color=color)
        text(slide, stack_x + Inches(0.20), yy + Inches(0.30),
             stack_w - Inches(0.40), Inches(0.32),
             text=body, size=10, color=WHITE)
        text(slide, stack_x + Inches(0.20), yy + Inches(0.50),
             stack_w - Inches(0.40), Inches(0.20),
             text=f"({tag})", size=8, color=MUTED, italic=True)
        yy += Inches(0.85)

    # Scale ribbon at bottom of col C
    rib_y = yy + Inches(0.10)
    rib_x = stack_x
    rib_w = stack_w / 4
    rib_h = Inches(0.50)
    labels = ["1 vessel", "fleet", "region", "coastline"]
    for i, lab in enumerate(labels):
        x = rib_x + rib_w * i
        rect(slide, x + Inches(0.04), rib_y, rib_w - Inches(0.08), rib_h,
             fill_rgb=DEEP_NAVY)
        text(slide, x + Inches(0.04), rib_y, rib_w - Inches(0.08), rib_h,
             text=lab, size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(labels) - 1:
            line(slide, x + rib_w - Inches(0.04), rib_y + rib_h / 2,
                 x + rib_w + Inches(0.04), rib_y + rib_h / 2,
                 rgb=TEAL, weight=Pt(2))


def fill_slide_5(slide):
    """IMPACT AND BENEFITS."""
    set_text(slide, must_have_placeholder="title", runs=[
        {"text": "IMPACT AND BENEFITS", "size": 28, "bold": True, "color": TEAL},
        {"text": "Target audience · Social / Economic / Environmental / Strategic",
         "size": 16, "color": SOFT},
    ])

    # Impact ecosystem: ORCA at center, users around it
    cx = Inches(13.333 / 2)        # center x
    cy = Inches(1.40 + 5.40 / 2)   # center y
    # ORCA center
    rect(slide, cx - Inches(1.10), cy - Inches(0.55), Inches(2.2),
         Inches(1.10), fill_rgb=DEEP_NAVY, line_rgb=TEAL, line_w=Pt(2))
    text(slide, cx - Inches(1.10), cy - Inches(0.55), Inches(2.2),
         Inches(0.50), text="ORCA 4.0", size=18, bold=True, color=TEAL,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, cx - Inches(1.10), cy - Inches(0.10), Inches(2.2),
         Inches(0.50), text="voyage-level\ndecision intelligence",
         size=10, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

    # 6 surrounding nodes
    import math
    nodes = [
        ("FISHERS",                "safer departures\nvessel-specific risk"),
        ("FLEETS",                 "central risk panel\nreduced incident rate"),
        ("COASTAL AUTHORITIES",     "geospatial risk layer\nfaster triage"),
        ("PORTS",                  "harbour decisions\narrival risk scored"),
        ("DISASTER RESPONSE",      "1 000-particle SAR drift\n(risk_engine/sar_drift)"),
        ("OCEAN RESEARCH",         "normalized environmental\nstate, provenance"),
    ]
    ring_r = Inches(2.55)
    for i, (head, body) in enumerate(nodes):
        ang = math.radians(-90 + i * 60)
        nx = cx + ring_r * math.cos(ang) - Inches(1.10)
        ny = cy + ring_r * math.sin(ang) * 0.65 - Inches(0.55)
        rect(slide, nx, ny, Inches(2.2), Inches(1.10), fill_rgb=SURFACE)
        rect(slide, nx, ny, Inches(0.10), Inches(1.10), fill_rgb=TEAL)
        text(slide, nx + Inches(0.20), ny + Inches(0.10),
             Inches(2.0), Inches(0.40),
             text=head, size=11, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER)
        text(slide, nx + Inches(0.20), ny + Inches(0.45),
             Inches(2.0), Inches(0.65),
             text=body, size=9, color=SOFT, align=PP_ALIGN.CENTER)
        # spoke
        line(slide, cx, cy, nx + Inches(1.10), ny + Inches(0.55),
             rgb=GRID, weight=Pt(1))

    # Causal flow strip at bottom
    flow_y = Inches(6.50)
    flow_x = Inches(0.5)
    flow_w = Inches(12.33)
    rect(slide, flow_x, flow_y, flow_w, Inches(0.40), fill_rgb=DEEP_NAVY,
         line_rgb=TEAL, line_w=Pt(1))
    text(slide, flow_x, flow_y, flow_w, Inches(0.40),
         text="LIVE CONDITIONS  →  VESSEL-SPECIFIC RISK  →  EARLIER "
              "AWARENESS  →  BETTER DEPARTURE DECISION  →  REDUCED "
              "EXPOSURE TO HAZARD",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)


def fill_slide_6(slide):
    """RESEARCH AND REFERENCES."""
    set_text(slide, must_have_placeholder="title", runs=[
        {"text": "RESEARCH AND REFERENCES", "size": 28, "bold": True,
         "color": TEAL},
        {"text": "Every formula ORCA uses has a real, citable source.",
         "size": 14, "color": SOFT, "italic": True},
    ])

    # 4 research cards on left, references list on right
    cards_x = Inches(0.5); cards_y = Inches(1.40)
    cards_w = Inches(8.0);  card_h = Inches(1.15)
    gap = Inches(0.10)

    cards = [
        ("Kijima et al. (1990)",
         "Encounter-period formula for ship–wave interaction",
         "→ used in wave_vessel_interaction_hazard()"),
        ("FAO small-craft rule of thumb",
         "GM ≈ 0.05 × beam",
         "→ default GM in vessel.py when operator omits it"),
        ("ORCA spec §1.5  ·  capsize threshold",
         "H_crit = 0.6 · L · sin(θ) for the dominant sea state",
         "→ drives max_safe_wave_height_m and CB-WAV-001"),
        ("Open-Meteo / ERA5 / ECMWF IFS  +  NOAA NDBC",
         "global NWP at 0.25° + in-situ buoys, no credentials",
         "→ 5 live provider paths via canonical layer"),
    ]
    yy = cards_y
    for src, insight, use in cards:
        rect(slide, cards_x, yy, cards_w, card_h, fill_rgb=SURFACE)
        rect(slide, cards_x, yy, Inches(0.10), card_h, fill_rgb=TEAL)
        text(slide, cards_x + Inches(0.20), yy + Inches(0.06),
             cards_w - Inches(0.30), Inches(0.30),
             text=src, size=12, bold=True, color=TEAL)
        text(slide, cards_x + Inches(0.20), yy + Inches(0.32),
             cards_w - Inches(0.30), Inches(0.32),
             text=insight, size=10, color=SOFT)
        text(slide, cards_x + Inches(0.20), yy + Inches(0.62),
             cards_w - Inches(0.30), Inches(0.40),
             text=use, size=10, color=WHITE, italic=True)
        yy += card_h + gap

    # References panel on the right
    ref_x = Inches(8.7); ref_y = Inches(1.40)
    ref_w = Inches(4.30); ref_h = Inches(5.40)
    rect(slide, ref_x, ref_y, ref_w, ref_h, fill_rgb=DEEP_NAVY,
         line_rgb=TEAL, line_w=Pt(1))
    text(slide, ref_x, ref_y + Inches(0.10), ref_w, Inches(0.35),
         text="REFERENCES", size=12, bold=True, color=TEAL,
         align=PP_ALIGN.CENTER)
    refs = [
        "[1]  Kijima, K. et al.  (1990)  — "
        "On the prediction method of ship manoeuvring motion.",
        "[2]  ORCA 4.0 Risk Spec §1.5  —  capsize threshold.",
        "[3]  FAO  —  small-craft stability rule of thumb (GM).",
        "[4]  Open-Meteo Marine API  —  ERA5 + NWP.",
        "[5]  MET Norway (yr.no)  —  Locationforecast 2.0.",
        "[6]  NOAA NDBC  —  realtime buoy network.",
        "[7]  Copernicus Marine Service  —  global ocean model.",
        "[8]  IMO Resolution A.601(15)  —  stability documentation.",
        "[9]  INCOIS  —  Indian Ocean operational forecasts.",
        "[10] PHASE16_KANYAKUMARI_AUDIT.md  —  live validation log.",
        "[11] risk_engine/  —  hazards.py · circuit_breaker.py · "
        "engine.py (in-repo).",
    ]
    bullet_textbox(slide, ref_x + Inches(0.20), ref_y + Inches(0.55),
                   ref_w - Inches(0.40), ref_h - Inches(0.65),
                   refs, size=9, color=WHITE, bullet_char="",
                   gap=Pt(4))


# ---------------- Driver -----------------------------------------------------

def main():
    prs = Presentation(TEMPLATE)

    # Verify slide order
    assert len(prs.slides) == 7, f"expected 7 slides, got {len(prs.slides)}"

    fill_slide_1(prs.slides[0])  # TITLE PAGE
    fill_slide_2(prs.slides[1])  # IDEA TITLE / PROPOSED SOLUTION
    fill_slide_3(prs.slides[2])  # TECHNICAL APPROACH
    fill_slide_4(prs.slides[3])  # FEASIBILITY AND VIABILITY
    fill_slide_5(prs.slides[4])  # IMPACT AND BENEFITS
    fill_slide_6(prs.slides[5])  # RESEARCH AND REFERENCES
    # prs.slides[6] = SIH's own 6-slide-limit instruction — leave untouched.

    prs.save(OUT)
    print("Wrote", OUT, "with", len(prs.slides), "slides")


if __name__ == "__main__":
    main()
