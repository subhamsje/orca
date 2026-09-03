"""
Build ORCA 4.0 SIH National Finals deck from real repo data.

Every claim is grounded in the ORCA repository:
- 24/24 risk-engine unit tests pass (backend/tests/test_risk_engine.py)
- 18 canonical parameters (backend/data_providers/canonical.py)
- 9 weighted hazards (backend/risk_engine/hazards.py)
- 7 deterministic circuit-breaker rules (backend/risk_engine/circuit_breaker.py)
- 34 FastAPI endpoints (backend/main.py)
- 5 live providers without credentials, 3 with credentials
- Real Kanyakumari MRSI = 75/100 (PHASE16_KANYAKUMARI_AUDIT.md)

NO FABRICATED STATISTICS. NO INVENTED FISHERMEN COUNTS. NO FAKE RUPEE FIGURES.
"""

import io
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- Design system ---------------------------------------------------

NAVY        = RGBColor(0x08, 0x1B, 0x33)   # deep ocean
DEEP_NAVY   = RGBColor(0x04, 0x10, 0x20)   # near-black navy
TEAL        = RGBColor(0x14, 0xB8, 0xA6)   # primary accent
TEAL_DARK   = RGBColor(0x0E, 0x8C, 0x7F)
CYAN        = RGBColor(0x67, 0xE8, 0xF9)   # data accent
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # caution
RED         = RGBColor(0xEF, 0x44, 0x44)   # extreme
GREEN       = RGBColor(0x10, 0xB9, 0x81)   # safe
WHITE       = RGBColor(0xF8, 0xFA, 0xFC)
SOFT        = RGBColor(0xCB, 0xD5, 0xE1)   # soft body text
MUTED       = RGBColor(0x94, 0xA3, 0xB8)   # captions
GRID        = RGBColor(0x1E, 0x3A, 0x5F)   # diagram lines
SURFACE     = RGBColor(0x0F, 0x24, 0x40)   # card surface


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill(slide, rgb: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def rect(slide, x, y, w, h, *, fill_rgb=None, line_rgb=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill_rgb is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
        if line_w is not None:
            shp.line.width = line_w
    return shp


def text(slide, x, y, w, h, *, text="", size=14, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
         italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def bullet(slide, x, y, w, h, items, *, size=14, color=WHITE, gap=Pt(6),
           bullet_char="•", bullet_color=None):
    """items: list of (heading, body) tuples or plain strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    bc = bullet_color or TEAL
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = gap
        r = p.add_run()
        r.text = bullet_char + "  "
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = bc
        if isinstance(item, tuple):
            head, body = item
            r2 = p.add_run()
            r2.text = head
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.bold = True
            r2.font.color.rgb = color
            if body:
                r3 = p.add_run()
                r3.text = "  " + body
                r3.font.name = "Calibri"
                r3.font.size = Pt(size)
                r3.font.bold = False
                r3.font.color.rgb = SOFT
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
    return tb


def chip(slide, x, y, w, h, label, *, fill_rgb=TEAL, text_color=NAVY, size=11):
    """Pill / chip."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Emu(40000)
    tf.margin_right = Emu(40000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shp


def line(slide, x1, y1, x2, y2, rgb=TEAL, weight=Pt(1.5)):
    cx = slide.shapes.add_connector(1, x1, y1, x2, y2)
    cx.line.color.rgb = rgb
    cx.line.width = weight
    return cx


def page_header(slide, eyebrow, title, slide_num, total=10):
    """Standardised slide chrome."""
    # Top accent bar
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_rgb=TEAL)
    # Eyebrow
    text(slide, Inches(0.5), Inches(0.18), Inches(10), Inches(0.3),
         text=eyebrow, size=10, color=TEAL, bold=True)
    # Title
    text(slide, Inches(0.5), Inches(0.42), Inches(11), Inches(0.7),
         text=title, size=28, color=WHITE, bold=True)
    # Slide counter
    text(slide, Inches(12.0), Inches(0.18), Inches(1.0), Inches(0.3),
         text=f"{slide_num:02d} / {total:02d}", size=10, color=MUTED,
         align=PP_ALIGN.RIGHT, bold=True)
    # Subtle divider
    rect(slide, Inches(0.5), Inches(1.15), Inches(12.33), Emu(9525),
         fill_rgb=GRID)


def footer(slide, left="ORCA 4.0  ·  SIH 2026  ·  Problem Statement ID: SIH26176"):
    text(slide, Inches(0.5), Inches(7.18), Inches(9), Inches(0.25),
         text=left, size=9, color=MUTED)
    text(slide, Inches(10.5), Inches(7.18), Inches(2.5), Inches(0.25),
         text="All claims verified against repo or audit log.",
         size=9, color=MUTED, align=PP_ALIGN.RIGHT, italic=True)


# ---------- Slide-specific helpers ------------------------------------------

def s01_hook(prs):
    s = add_blank(prs)
    fill(s, NAVY)

    # Background wave
    sw = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-2.5),
                            Inches(7), Inches(7))
    sw.shadow.inherit = False
    sw.fill.solid(); sw.fill.fore_color.rgb = TEAL_DARK
    sw.line.fill.background()
    swp = sw.fill
    # Use solid (transparency not always honoured); keep low chroma
    sw.fill.fore_color.rgb = RGBColor(0x10, 0x3A, 0x52)

    # ORCA brand
    text(s, Inches(0.6), Inches(0.4), Inches(7), Inches(0.6),
         text="ORCA 4.0", size=24, bold=True, color=TEAL)
    text(s, Inches(0.6), Inches(0.95), Inches(7), Inches(0.3),
         text="SMART INDIA HACKATHON 2026 · NATIONAL FINALS",
         size=10, color=MUTED, bold=True)
    text(s, Inches(0.6), Inches(1.20), Inches(7), Inches(0.3),
         text="Problem Statement ID: SIH26176 · ISRO / INCOIS",
         size=10, color=MUTED)

    # Eyebrow
    text(s, Inches(0.6), Inches(2.2), Inches(8), Inches(0.4),
         text="THE HOOK", size=11, bold=True, color=TEAL)

    # Headline
    text(s, Inches(0.6), Inches(2.55), Inches(12), Inches(1.7),
         text="The ocean is dynamic.\nThe decision cannot be blind.",
         size=44, bold=True, color=WHITE)

    # Sub
    text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.5),
         text="Existing marine information products deliver forecasts. None tell the fisherman",
         size=16, color=SOFT)
    text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
         text="what those forecasts mean for THIS vessel on THIS route at THIS time.",
         size=16, color=SOFT)

    # Three-pillar strip
    y0 = Inches(5.7)
    h = Inches(0.85)
    gap = Inches(0.15)
    x = Inches(0.6)
    w_each = Inches((13.333 - 0.6*2 - 0.15*2) / 3)
    for label, sub in [
        ("FORECAST",   "Open-Meteo · MET Norway · NDBC buoys"),
        ("VESSEL",     "16-field validated digital twin"),
        ("DECISION",   "ORCA MRSI · deterministic circuit breaker"),
    ]:
        rect(s, x, y0, w_each, h, fill_rgb=SURFACE)
        rect(s, x, y0, w_each, Inches(0.04), fill_rgb=TEAL)
        text(s, x + Inches(0.25), y0 + Inches(0.12), w_each - Inches(0.5),
             Inches(0.4), text=label, size=14, bold=True, color=TEAL)
        text(s, x + Inches(0.25), y0 + Inches(0.42), w_each - Inches(0.5),
             Inches(0.4), text=sub, size=11, color=SOFT)
        x += w_each + gap

    footer(s)


def s02_solution(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "THE SOLUTION", "ORCA: a digital twin of the voyage.", 2)

    # Top row: input blocks flowing into engine, then outputs
    y_inputs = Inches(1.55)
    inp_h = Inches(1.0)
    inp_w = Inches(1.95)
    gap = Inches(0.10)
    inputs = [
        ("LIVE OCEAN",   "waves · swell · current\nSST · salinity"),
        ("WEATHER",      "wind · gust\npressure · visibility"),
        ("VESSEL",       "length · beam · draft\nGM · max_wave"),
        ("ALERTS",       "IMD CWD hook\nINCOIS EWC hook"),
        ("ROUTE",        "up to 6 waypoints\nper-segment state"),
        ("GEOFENCE",     "IMBL buffer\nNaval Restricted Zones"),
    ]
    x = Inches(0.5)
    for label, body in inputs:
        rect(s, x, y_inputs, inp_w, inp_h, fill_rgb=SURFACE)
        rect(s, x, y_inputs, inp_w, Inches(0.04), fill_rgb=TEAL)
        text(s, x + Inches(0.15), y_inputs + Inches(0.08),
             inp_w - Inches(0.3), Inches(0.32),
             text=label, size=11, bold=True, color=TEAL)
        text(s, x + Inches(0.15), y_inputs + Inches(0.40),
             inp_w - Inches(0.3), Inches(0.55),
             text=body, size=10, color=SOFT)
        x += inp_w + gap

    # Engine block
    eng_y = Inches(3.05)
    eng_h = Inches(2.1)
    rect(s, Inches(0.5), eng_y, Inches(12.33), eng_h, fill_rgb=DEEP_NAVY,
         line_rgb=TEAL, line_w=Pt(1.5))
    text(s, Inches(0.5), eng_y + Inches(0.05), Inches(12), Inches(0.4),
         text="ORCA INTELLIGENCE ENGINE", size=12, bold=True, color=TEAL,
         align=PP_ALIGN.CENTER)
    # 5 stage pipeline
    stages = ["Environmental\nState",
              "Vessel Digital\nTwin",
              "Maritime\nPhysics",
              "Deterministic\nCircuit Breaker",
              "ORCA MRSI\n0–100"]
    sx = Inches(0.7)
    sw = Inches((12.33 - 0.4) / 5)
    for i, st in enumerate(stages):
        rect(s, sx + Inches(0.05), eng_y + Inches(0.7),
             sw - Inches(0.1), Inches(0.95), fill_rgb=SURFACE)
        text(s, sx + Inches(0.05), eng_y + Inches(0.95),
             sw - Inches(0.1), Inches(0.5),
             text=st, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            line(s, sx + sw - Inches(0.05), eng_y + Inches(1.18),
                 sx + sw + Inches(0.05), eng_y + Inches(1.18),
                 rgb=TEAL, weight=Pt(2.0))
        sx += sw

    # Outputs
    out_y = Inches(5.55)
    out_h = Inches(1.10)
    outputs = [
        ("SAFE TO PROCEED",      GREEN,  "< 20 / 100"),
        ("PROCEED WITH CAUTION", TEAL,   "20 – 40"),
        ("HIGH RISK",            AMBER,  "40 – 70"),
        ("EXTREME DANGER",       RED,    "≥ 70 (or CB hit)"),
    ]
    x = Inches(0.5)
    out_w = Inches((13.333 - 0.5*2 - 0.15*3) / 4)
    for label, color, range_ in outputs:
        rect(s, x, out_y, out_w, out_h, fill_rgb=SURFACE)
        rect(s, x, out_y, out_w, Inches(0.06), fill_rgb=color)
        text(s, x + Inches(0.15), out_y + Inches(0.18), out_w - Inches(0.3),
             Inches(0.4), text=label, size=13, bold=True, color=color)
        text(s, x + Inches(0.15), out_y + Inches(0.62), out_w - Inches(0.3),
             Inches(0.4), text=range_, size=11, color=SOFT)
        x += out_w + Inches(0.15)

    footer(s)


def s03_architecture(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "TECHNICAL APPROACH",
                "One pipeline. Ocean → vessel → risk → route.", 3)

    # Left column: text bullets (verified)
    bullet(s, Inches(0.5), Inches(1.40), Inches(4.6), Inches(5.5), [
        ("5 live providers",
         "MET Norway · Open-Meteo Marine · Open-Meteo ECMWF · "
         "Open-Meteo Forecast · NOAA NDBC buoys — all production, "
         "no credentials required."),
        ("3 provider slots wired, gated on env vars",
         "StormGlass · Copernicus Marine · NASA Earthdata — slot "
         "exists in provider registry; activated when keys are set."),
        ("Canonical record per parameter",
         "18 parameters, each with source, dataset, observation_time, "
         "distance_km, freshness state (CURRENT / RECENT / STALE / UNAVAILABLE)."),
        ("Deterministic safety circuit breaker",
         "7 hard rules (CB-DQ, CB-CYC×2, CB-GEO×2, CB-WND, CB-WAV) "
         "run BEFORE the continuous risk engine."),
        ("ORCA MRSI 0–100",
         "9 weighted hazards · monotonic · reproducible (SHA-256 of "
         "input snapshot = assessment_id)."),
    ], size=11)

    # Right column: architecture diagram (drawn)
    dx = Inches(5.4)
    dy = Inches(1.40)
    dw = Inches(7.5)
    dh = Inches(5.4)

    layers = [
        ("DATA SOURCES", "MET Norway · Open-Meteo · NDBC · Sentinel-1 (roadmap) · IMD/INCOIS (roadmap)",
         TEAL, SURFACE),
        ("PROVIDER ABSTRACTION + CIRCUIT BREAKER + RATE LIMIT",
         "providers/base.py · providers/registry.py · 6 providers registered",
         CYAN, SURFACE),
        ("CANONICAL NORMALIZATION + FRESHNESS",
         "data_providers/canonical.py · 18 params · CURRENT / RECENT / STALE / UNAVAILABLE",
         TEAL, SURFACE),
        ("VESSEL DIGITAL TWIN + ENVIRONMENTAL STATE",
         "risk_engine/vessel.py · 16-field validated profile · risk_engine/state.py",
         CYAN, SURFACE),
        ("MARITIME PHYSICS + 9 HAZARD FUNCTIONS",
         "risk_engine/hazards.py · Kijima-1990 encounter period · vessel-relative thresholds",
         TEAL, SURFACE),
        ("DETERMINISTIC CIRCUIT BREAKER → ORCA MRSI",
         "risk_engine/circuit_breaker.py · risk_engine/engine.py · 24 unit tests pass",
         CYAN, SURFACE),
        ("ROUTE RISK + REPLAY STORE + 34 API ENDPOINTS",
         "risk_engine/route_risk.py · risk_engine/replay.py · backend/main.py",
         TEAL, SURFACE),
    ]
    h = (dh - Inches(0.6)) / len(layers)
    y = dy
    for i, (head, body, accent, bg) in enumerate(layers):
        rect(s, dx, y, dw, h - Inches(0.07), fill_rgb=bg)
        rect(s, dx, y, Inches(0.10), h - Inches(0.07), fill_rgb=accent)
        text(s, dx + Inches(0.25), y + Inches(0.04),
             dw - Inches(0.35), Inches(0.30),
             text=head, size=10, bold=True, color=accent)
        text(s, dx + Inches(0.25), y + Inches(0.32),
             dw - Inches(0.35), Inches(0.45),
             text=body, size=9, color=SOFT)
        y += h

    footer(s)


def s04_differentiation(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "DIFFERENTIATION",
                "Forecast ≠ operational risk.", 4)

    # Lead line
    text(s, Inches(0.5), Inches(1.30), Inches(12.33), Inches(0.5),
         text="ORCA does not aggregate marine data. It produces a "
              "vessel-conditioned, route-aware decision.",
         size=15, color=SOFT, italic=True)

    # Comparison table
    headers = [
        "Capability", "Generic\nweather app", "Marine\nforecast sites",
        "Generic\nroute planner", "AIS\nplatforms", "ORCA 4.0"
    ]
    rows = [
        ["Live ocean (wave · SST · current)",
         "Limited", "Yes", "No", "No", "Yes (5 providers, provenance)"],
        ["Vessel-aware thresholds (H_crit · GM)",
         "No", "No", "No", "No", "Yes (16-field twin)"],
        ["Wave–vessel interaction (encounter period, angle)",
         "No", "No", "No", "No", "Yes (Kijima 1990)"],
        ["Deterministic circuit breaker",
         "No", "Limited", "No", "No", "Yes (7 hard rules)"],
        ["Per-segment route risk",
         "No", "No", "Distance only", "No", "Yes (route_risk.py)"],
        ["Data freshness + provenance per value",
         "No", "Partial", "No", "Partial", "Yes (CURRENT/STALE/UNAVAILABLE)"],
        ["Source-aware confidence + uncertainty",
         "No", "No", "No", "No", "Yes (risk_uncertainty 0–1)"],
        ["Honest INSUFFICIENT_DATA verdict",
         "No", "No", "No", "No", "Yes (no fake SAFE)"],
    ]

    table_x = Inches(0.5)
    table_y = Inches(1.95)
    col_w = [Inches(2.6), Inches(1.55), Inches(1.55), Inches(1.55),
             Inches(1.45), Inches(2.63)]
    row_h = Inches(0.42)

    # Header row
    x = table_x
    for i, h in enumerate(headers):
        bg = TEAL if i == len(headers) - 1 else GRID
        rect(s, x, table_y, col_w[i], row_h, fill_rgb=bg)
        text(s, x, table_y, col_w[i], row_h,
             text=h, size=10, bold=True,
             color=NAVY if i == len(headers) - 1 else WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    # Body
    for r, row in enumerate(rows):
        x = table_x
        y = table_y + row_h + r * row_h
        bg = DEEP_NAVY if r % 2 == 0 else SURFACE
        for c, val in enumerate(row):
            rect(s, x, y, col_w[c], row_h, fill_rgb=bg)
            color = TEAL if c == len(row) - 1 and val.startswith("Yes") \
                else (SOFT if c == 0 else MUTED)
            bold = (c == 0) or (c == len(row) - 1)
            text(s, x + Inches(0.10), y, col_w[c] - Inches(0.20), row_h,
                 text=val, size=9, bold=bold, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[c]

    # Footnote
    text(s, Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.4),
         text="ORCA’s existing multi-objective route planner is a "
              "v0 placeholder; the production Pareto solver is on the "
              "public roadmap. Per-segment ORCA MRSI is shipped today.",
         size=10, color=AMBER, italic=True)

    footer(s)


def s05_feasibility(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "FEASIBILITY + IMPACT",
                "Built to run fast. Designed to scale.", 5)

    # 4 quadrants
    qx0 = Inches(0.5)
    qy0 = Inches(1.40)
    qw  = Inches((13.333 - 0.5*2 - 0.25) / 2)
    qh  = Inches((7.5 - 1.40 - 0.7 - 0.25) / 2)

    def quad(slide, x, y, w, h, head, body_lines, accent):
        rect(slide, x, y, w, h, fill_rgb=SURFACE)
        rect(slide, x, y, Inches(0.10), h, fill_rgb=accent)
        text(slide, x + Inches(0.30), y + Inches(0.12),
             w - Inches(0.40), Inches(0.4),
             text=head, size=14, bold=True, color=accent)
        for i, line in enumerate(body_lines):
            text(slide, x + Inches(0.30), y + Inches(0.50) + Inches(i * 0.32),
                 w - Inches(0.40), Inches(0.32),
                 text=line, size=10, color=SOFT)

    # A. Technical feasibility
    quad(s, qx0, qy0, qw, qh, "A.  TECHNICAL FEASIBILITY", [
        "BUILT          FastAPI + React/TS · 34 endpoints · 9 services",
        "BUILT          Canonical data layer (18 params, provenance)",
        "BUILT          Deterministic circuit breaker (7 rules)",
        "BUILT          ORCA MRSI engine (9 hazards, version-controlled)",
        "BUILT          Per-segment route risk · replay store",
        "BUILT          Kanyakumari live validation: MRSI = 75 / 100",
        "ROADMAP        IMD CWD parser · INCOIS ERDDAP dataset mapping",
        "ROADMAP        Sentinel-1 SAR · LLM explainability layer",
    ], TEAL)

    # B. Cost model
    quad(s, qx0 + qw + Inches(0.25), qy0, qw, qh, "B.  DEPLOYMENT COST MODEL", [
        "PROTOTYPE  Cloud-only (Vercel/Render free tier + Open-Meteo free)",
        "            Indicative — depends on cloud vendor",
        "PILOT      Backend on commodity VPS · ~50 USD/mo for 10⁵ assessments",
        "            Open-Meteo free tier caps at 10k req/day",
        "PRODUCTION Adds StormGlass/Copernicus keys for higher fidelity",
        "            Indicative — vendor-dependent",
        "EDGE       Optional onboard node: ESP32 + GPS + LoRa",
        "            BOM marked indicative — not yet procured",
    ], CYAN)

    # C. Impact framework (KPI-led)
    quad(s, qx0, qy0 + qh + Inches(0.25), qw, qh,
         "C.  IMPACT FRAMEWORK — pilot KPIs to measure", [
        "USER           Problem                  → ORCA intervention   → KPI",
        "Fisher         unsafe departure         → vessel-conditioned   → safer",
        "                                       →   risk score          →   departures",
        "Fleet op.      fragmented monitoring    → central risk panel  → reduced",
        "                                       →                      →   incident rate",
        "Authority      distributed hazards      → geospatial risk layer→ faster",
        "                                       →                      →   triage",
        "All KPIs       validated against pilot data — not extrapolated",
    ], TEAL)

    # D. Scale
    quad(s, qx0 + qw + Inches(0.25), qy0 + qh + Inches(0.25), qw, qh,
         "D.  NATIONAL SCALE — phased roadmap", [
        "PHASE 1   Prototype live (current)",
        "PHASE 2   One coastal-state pilot",
        "PHASE 3   Multi-vessel fleet deployment",
        "PHASE 4   Multi-state coastal integration",
        "PHASE 5   National maritime intelligence layer",
        "Each phase adds a verifiable capability, not a slide.",
    ], CYAN)

    footer(s)


def s06_research(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "RESEARCH + VALIDATION + VISION",
                "From one vessel to India’s coastline.", 6)

    # Left: scientific foundation table
    text(s, Inches(0.5), Inches(1.30), Inches(7.2), Inches(0.4),
         text="SCIENTIFIC FOUNDATION — what we actually used", size=12,
         bold=True, color=TEAL)

    src = [
        ("Kijima et al. 1990",
         "encounter-period formula for ship–wave interaction",
         "ORCA uses it in wave_vessel_interaction_hazard()"),
        ("FAO small-craft rule of thumb",
         "GM ≈ 0.05 × beam",
         "default GM estimate in vessel.py when operator does not supply it"),
        ("Capsize threshold, ORCA spec §1.5",
         "H_crit = 0.6 · L · sin(θ)",
         "drives vessel.max_safe_wave_height_m and CB-WAV-001"),
        ("Open-Meteo / ERA5 / ECMWF IFS",
         "global NWP at 0.25° resolution, no credentials",
         "5 live provider paths; canonical layer wraps them"),
        ("NOAA NDBC realtime buoys",
         "in-situ wave / wind verification",
         "used for nearest-station validation when buoy is < 1200 km"),
    ]
    y = Inches(1.70)
    for head, body, use in src:
        rect(s, Inches(0.5), y, Inches(7.2), Inches(0.78), fill_rgb=SURFACE)
        rect(s, Inches(0.5), y, Inches(0.08), Inches(0.78), fill_rgb=TEAL)
        text(s, Inches(0.65), y + Inches(0.06), Inches(7.0), Inches(0.3),
             text=head, size=11, bold=True, color=TEAL)
        text(s, Inches(0.65), y + Inches(0.32), Inches(7.0), Inches(0.25),
             text=body, size=10, color=SOFT)
        text(s, Inches(0.65), y + Inches(0.55), Inches(7.0), Inches(0.25),
             text="→ " + use, size=9, color=MUTED, italic=True)
        y += Inches(0.86)

    # Right: validation + scale
    text(s, Inches(8.0), Inches(1.30), Inches(4.83), Inches(0.4),
         text="VALIDATION — what we have actually run", size=12,
         bold=True, color=TEAL)

    bullet(s, Inches(8.0), Inches(1.70), Inches(4.83), Inches(2.5), [
        ("24 / 24 risk-engine unit tests",
         "monotonicity, reproducibility, vessel-specificity, "
         "circuit breakers, freshness, wave-vessel interaction."),
        ("68 backend tests in repo",
         "62 pass; 6 legacy test_backend.py tests still reference the "
         "old hardcoded service layer (see PHASE1_AUDIT)."),
        ("Kanyakumari live assessment",
         "8.084°N, 77.55°E · vessel length 8.5 m · "
         "MRSI = 75 / 100, HIGH_RISK_IMBL — reproducible."),
        ("Real providers, no fixtures",
         "every value in the response is traceable to source_id "
         "in CanonicalRecord."),
    ], size=10)

    # Scale block
    text(s, Inches(8.0), Inches(4.50), Inches(4.83), Inches(0.4),
         text="SCALE", size=12, bold=True, color=TEAL)
    stages = ["Vessel", "Coastal\ncommunity", "Fleet", "Port", "Authority",
              "National intelligence"]
    sx = Inches(8.0)
    sw = Inches(4.83 / 6)
    sy = Inches(4.90)
    for i, st in enumerate(stages):
        rect(s, sx + Inches(0.05), sy, sw - Inches(0.10), Inches(0.85),
             fill_rgb=SURFACE)
        text(s, sx + Inches(0.05), sy, sw - Inches(0.10), Inches(0.85),
             text=st, size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            line(s, sx + sw - Inches(0.05), sy + Inches(0.42),
                 sx + sw + Inches(0.05), sy + Inches(0.42),
                 rgb=TEAL, weight=Pt(2.0))
        sx += sw

    text(s, Inches(8.0), Inches(5.95), Inches(4.83), Inches(1.0),
         text="Each step widens the user base without changing the "
              "decision model — the same MRSI engine scales.",
         size=10, color=SOFT, italic=True)

    footer(s)


def s07_demo(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "LIVE DEMO  ·  60 seconds",
                "Let the product speak.", 7)

    # Left: 6-step flow
    text(s, Inches(0.5), Inches(1.30), Inches(7.2), Inches(0.4),
         text="60-SECOND DEMO FLOW", size=12, bold=True, color=TEAL)

    flow = [
        ("Select location",       "Mumbai Sassoon Dock · 18.94°N, 72.83°E"),
        ("Fetch live data",       "5 providers · provenance per value"),
        ("Vessel digital twin",   "8.5 m fishing craft · LADEN · heading 270°"),
        ("Compute risk",          "9 hazards · 7 CB rules · MRSI 0–100"),
        ("Pick route",            "Origin → 1 mid → destination · "
                                  "per-segment state"),
        ("Identify hazardous segment", "Worst segment highlighted · "
                                       "reason shown in plain language"),
    ]
    y = Inches(1.75)
    for i, (head, body) in enumerate(flow):
        rect(s, Inches(0.5), y, Inches(0.55), Inches(0.55), fill_rgb=TEAL)
        text(s, Inches(0.5), y, Inches(0.55), Inches(0.55),
             text=str(i + 1), size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, Inches(1.15), y, Inches(6.55), Inches(0.55), fill_rgb=SURFACE)
        text(s, Inches(1.30), y + Inches(0.05), Inches(6.4), Inches(0.3),
             text=head, size=12, bold=True, color=WHITE)
        text(s, Inches(1.30), y + Inches(0.30), Inches(6.4), Inches(0.25),
             text=body, size=10, color=SOFT)
        y += Inches(0.70)

    # Right: Kanyakumari screenshot block (synthetic, but values are real)
    text(s, Inches(8.0), Inches(1.30), Inches(4.83), Inches(0.4),
         text="KANYAKUMARI · LIVE RESULT", size=12, bold=True, color=TEAL)

    box = (Inches(8.0), Inches(1.75), Inches(4.83), Inches(5.2))
    rect(s, *box, fill_rgb=DEEP_NAVY, line_rgb=TEAL, line_w=Pt(1.5))
    # header
    rect(s, box[0], box[1], box[2], Inches(0.45), fill_rgb=TEAL)
    text(s, box[0] + Inches(0.15), box[1] + Inches(0.08),
         box[2] - Inches(0.3), Inches(0.3),
         text="8.084°N, 77.5505°E  ·  8.5 m craft  ·  W heading",
         size=11, bold=True, color=NAVY)
    # MRSI score
    rect(s, box[0] + Inches(0.25), box[1] + Inches(0.65),
         box[2] - Inches(0.5), Inches(1.05), fill_rgb=AMBER)
    text(s, box[0] + Inches(0.25), box[1] + Inches(0.68),
         box[2] - Inches(0.5), Inches(0.4),
         text="ORCA MRSI", size=11, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
    text(s, box[0] + Inches(0.25), box[1] + Inches(0.95),
         box[2] - Inches(0.5), Inches(0.7),
         text="75 / 100", size=40, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
    text(s, box[0] + Inches(0.25), box[1] + Inches(1.45),
         box[2] - Inches(0.5), Inches(0.3),
         text="HIGH_RISK_IMBL  ·  forced by CB-GEO-002",
         size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Components table (real numbers)
    rows = [
        ("wave_height",           "+16.8"),
        ("wave_vessel_interaction", "+10.0"),
        ("gust",                   "+5.7"),
        ("current",                "+4.3"),
        ("pressure",               "+1.1"),
        ("visibility",             "+0.7"),
        ("wind",                   "+0.4"),
        ("precipitation",          "+0.0"),
        ("RAW",                    "38.9"),
        ("CB override",            "→ 75 (IMBL buffer hit)"),
    ]
    ry = box[1] + Inches(1.95)
    for i, (k, v) in enumerate(rows):
        bg = SURFACE if i % 2 == 0 else DEEP_NAVY
        rect(s, box[0] + Inches(0.15), ry, box[2] - Inches(0.3),
             Inches(0.27), fill_rgb=bg)
        is_total = k in ("RAW", "CB override")
        text(s, box[0] + Inches(0.25), ry, box[2] - Inches(0.9),
             Inches(0.27), text=k, size=9, bold=is_total,
             color=TEAL if is_total else SOFT, anchor=MSO_ANCHOR.MIDDLE)
        text(s, box[0] + box[2] - Inches(1.6), ry, Inches(1.3),
             Inches(0.27), text=v, size=9, bold=is_total,
             color=TEAL if is_total else WHITE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
        ry += Inches(0.27)

    text(s, box[0] + Inches(0.25), box[1] + box[3] - Inches(0.45),
         box[2] - Inches(0.5), Inches(0.3),
         text="Source: PHASE16_KANYAKUMARI_AUDIT.md · 2026-09-03",
         size=8, color=MUTED, italic=True)

    footer(s)


def s08_ai(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "WHERE AI ACTUALLY ADDS VALUE",
                "Deterministic safety ≠ buzzword ML.", 8)

    # Top: 4-layer model
    text(s, Inches(0.5), Inches(1.30), Inches(12.33), Inches(0.4),
         text="ORCA deliberately separates four computation layers.",
         size=14, color=SOFT, italic=True)

    layers = [
        ("DETERMINISTIC SAFETY",
         "Official warnings · H_crit · IMBL buffer · naval zones · "
         "data quality gate",
         "Code, not learned.\nA cyclone warning is a cyclone warning — "
         "the model cannot override the breaker.",
         TEAL, SURFACE),
        ("MARITIME PHYSICS",
         "Kijima 1990 encounter period · vessel-relative thresholds · "
         "steepness H/T",
         "Tables in hazards.py.\nVersion-controlled, auditable, "
         "explainable to a domain expert.",
         CYAN, SURFACE),
        ("ROUTE OPTIMIZATION",
         "Per-segment ORCA MRSI → worst segment, mean risk, "
         "hazardous distance",
         "Multi-objective Pareto is on the roadmap (current implementation "
         "is a v0 placeholder).",
         TEAL, SURFACE),
        ("EXPLAINABILITY",
         "Per-component breakdown that reconciles to total · "
         "future: LLM-plain-language rephrasing",
         "The LLM can only rephrase a forced verdict — it cannot "
         "change the score.",
         CYAN, SURFACE),
    ]
    y = Inches(1.85)
    h = Inches(1.20)
    for head, what, how, accent, bg in layers:
        rect(s, Inches(0.5), y, Inches(12.33), h, fill_rgb=bg)
        rect(s, Inches(0.5), y, Inches(0.10), h, fill_rgb=accent)
        text(s, Inches(0.70), y + Inches(0.08), Inches(2.4), Inches(0.4),
             text=head, size=12, bold=True, color=accent)
        text(s, Inches(3.20), y + Inches(0.10), Inches(4.4), Inches(1.0),
             text=what, size=10, color=WHITE)
        text(s, Inches(7.80), y + Inches(0.10), Inches(5.0), Inches(1.0),
             text=how, size=10, color=SOFT, italic=True)
        y += h + Inches(0.08)

    # Bottom: ML candidate table (honest)
    text(s, Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3),
         text="Where ML is on the roadmap — not shipped today. "
              "Validation in progress against baseline.",
         size=10, color=AMBER, italic=True)

    footer(s)


def s09_hardware(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "FIELD DEPLOYMENT  ·  BOM",
                "Cloud-first today, edge-capable roadmap.", 9)

    # Left: BOM table
    text(s, Inches(0.5), Inches(1.30), Inches(7.5), Inches(0.4),
         text="INDICATIVE BOM  ·  prototype, vendor-dependent",
         size=12, bold=True, color=TEAL)

    headers = ["Component", "Purpose", "Indicative"]
    rows = [
        ["GPS module (u-blox NEO-6M)",
         "Position · speed · plausibility check",
         "≈ ₹600"],
        ["MCU (ESP32-S3)",
         "Edge aggregation · 16-byte packet pack",
         "≈ ₹450"],
        ["LoRa 868 MHz transceiver",
         "Sub-GHz fleet / shore mesh",
         "≈ ₹700"],
        ["IMU (MPU6050)",
         "Heave / pitch / roll — future capsize proxy",
         "≈ ₹250"],
        ["Barometer + temp + RH",
         "Local sensor sanity-check vs provider",
         "≈ ₹300"],
        ["Marine enclosure + cabling",
         "IP65 mount + power harness",
         "≈ ₹900"],
        ["Phone (vessel-side UI)",
         "Operator UI · voice input · offline map",
         "already owned"],
        ["TOTAL (prototype, indicative)", "", "≈ ₹3,200"],
    ]
    col_w = [Inches(3.0), Inches(3.0), Inches(1.5)]
    row_h = Inches(0.34)
    table_x = Inches(0.5)
    table_y = Inches(1.78)
    x = table_x
    for i, h in enumerate(headers):
        rect(s, x, table_y, col_w[i], row_h, fill_rgb=GRID)
        text(s, x + Inches(0.1), table_y, col_w[i] - Inches(0.2), row_h,
             text=h, size=10, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        x = table_x
        y = table_y + row_h * (r + 1)
        bg = DEEP_NAVY if r % 2 == 0 else SURFACE
        is_total = r == len(rows) - 1
        for c, val in enumerate(row):
            rect(s, x, y, col_w[c], row_h, fill_rgb=bg)
            text(s, x + Inches(0.10), y, col_w[c] - Inches(0.20), row_h,
                 text=val, size=9, bold=is_total,
                 color=TEAL if is_total else WHITE if c == 0 else SOFT,
                 anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[c]

    text(s, Inches(0.5), Inches(5.20), Inches(7.5), Inches(0.4),
         text="Production BOM will change once a vendor is selected. "
              "Numbers above are indicative.",
         size=9, color=MUTED, italic=True)

    # Right: deployment topology
    text(s, Inches(8.2), Inches(1.30), Inches(4.83), Inches(0.4),
         text="DEPLOYMENT TOPOLOGY", size=12, bold=True, color=TEAL)

    # VESSEL
    rect(s, Inches(8.6), Inches(1.85), Inches(4.0), Inches(0.85),
         fill_rgb=SURFACE)
    rect(s, Inches(8.6), Inches(1.85), Inches(0.10), Inches(0.85),
         fill_rgb=TEAL)
    text(s, Inches(8.75), Inches(1.90), Inches(3.8), Inches(0.3),
         text="VESSEL", size=11, bold=True, color=TEAL)
    text(s, Inches(8.75), Inches(2.20), Inches(3.8), Inches(0.5),
         text="GPS · IMU · barometer · ESP32 + LoRa",
         size=10, color=SOFT)

    arrow_y = Inches(2.85)
    line(s, Inches(10.6), Inches(2.70), Inches(10.6), Inches(3.10),
         rgb=TEAL, weight=Pt(2.0))

    # EDGE / PHONE
    rect(s, Inches(8.6), Inches(3.15), Inches(4.0), Inches(0.85),
         fill_rgb=SURFACE)
    rect(s, Inches(8.6), Inches(3.15), Inches(0.10), Inches(0.85),
         fill_rgb=TEAL)
    text(s, Inches(8.75), Inches(3.20), Inches(3.8), Inches(0.3),
         text="EDGE / PHONE", size=11, bold=True, color=TEAL)
    text(s, Inches(8.75), Inches(3.50), Inches(3.8), Inches(0.5),
         text="Operator UI · offline map · cached forecast",
         size=10, color=SOFT)

    line(s, Inches(10.6), Inches(4.00), Inches(10.6), Inches(4.40),
         rgb=TEAL, weight=Pt(2.0))

    # CLOUD
    rect(s, Inches(8.6), Inches(4.45), Inches(4.0), Inches(0.85),
         fill_rgb=SURFACE)
    rect(s, Inches(8.6), Inches(4.45), Inches(0.10), Inches(0.85),
         fill_rgb=TEAL)
    text(s, Inches(8.75), Inches(4.50), Inches(3.8), Inches(0.3),
         text="ORCA CLOUD", size=11, bold=True, color=TEAL)
    text(s, Inches(8.75), Inches(4.80), Inches(3.8), Inches(0.5),
         text="5 live providers · canonical · risk engine",
         size=10, color=SOFT)

    line(s, Inches(10.6), Inches(5.30), Inches(10.6), Inches(5.70),
         rgb=TEAL, weight=Pt(2.0))

    # COMMAND
    rect(s, Inches(8.6), Inches(5.75), Inches(4.0), Inches(0.85),
         fill_rgb=SURFACE)
    rect(s, Inches(8.6), Inches(5.75), Inches(0.10), Inches(0.85),
         fill_rgb=TEAL)
    text(s, Inches(8.75), Inches(5.80), Inches(3.8), Inches(0.3),
         text="COMMAND DASHBOARD", size=11, bold=True, color=TEAL)
    text(s, Inches(8.75), Inches(6.10), Inches(3.8), Inches(0.5),
         text="Authority view · fleet risk panel · replay",
         size=10, color=SOFT)

    text(s, Inches(8.6), Inches(6.80), Inches(4.0), Inches(0.3),
         text="Connectivity loss → edge keeps last-known state + "
              "cached forecast + local safety rules.",
         size=9, color=AMBER, italic=True)

    footer(s)


def s10_team(prs):
    s = add_blank(prs)
    fill(s, NAVY)
    page_header(s, "WHY THIS TEAM",
                "We are not pitching an idea. We are building the system.",
                10)

    # Top row: capability matrix
    text(s, Inches(0.5), Inches(1.30), Inches(12.33), Inches(0.4),
         text="CAPABILITY MATRIX  ·  what is actually in the repo",
         size=12, bold=True, color=TEAL)

    rows = [
        ("Backend / API",         "FastAPI · 34 endpoints · async providers"),
        ("Maritime physics",      "Kijima-1990 encounter · capsize threshold · 9 hazards"),
        ("Geospatial",            "haversine · polygon geofence · IMBL buffer · naval zones"),
        ("Deterministic safety",  "7 circuit-breaker rules · CB-WAV/CYC/GEO/WND/DQ"),
        ("Data engineering",      "Canonical layer · freshness policy · "
                                  "provider abstraction + circuit breaker"),
        ("Frontend / maps",       "React + TS + Vite · Leaflet · OSM base · "
                                  "provenance badges"),
        ("Testing",               "24 risk-engine tests · 68 backend tests in repo"),
        ("Embedded / IoT (planned)", "ESP32 + LoRa · indicative BOM on slide 9"),
    ]
    col_w = [Inches(3.5), Inches(8.83)]
    row_h = Inches(0.40)
    x = Inches(0.5)
    y = Inches(1.80)
    for r, (a, b) in enumerate(rows):
        bg = DEEP_NAVY if r % 2 == 0 else SURFACE
        rect(s, x, y, col_w[0], row_h, fill_rgb=bg)
        rect(s, x + col_w[0], y, col_w[1], row_h, fill_rgb=bg)
        text(s, x + Inches(0.15), y, col_w[0] - Inches(0.3), row_h,
             text=a, size=11, bold=True, color=TEAL,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + col_w[0] + Inches(0.15), y, col_w[1] - Inches(0.3), row_h,
             text=b, size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

    # Three callouts at the bottom
    cy = Inches(5.80)
    cw = Inches(4.05)
    ch = Inches(1.20)
    callouts = [
        ("24 / 24 risk-engine unit tests pass",
         "Reproducibility, monotonicity, vessel-specificity, "
         "circuit breakers — all green."),
        ("5 live providers, no fixtures",
         "Every value in the response carries a source_id and "
         "an observation timestamp."),
        ("Kanyakumari live run",
         "8.084°N, 77.55°E → MRSI 75/100, "
         "HIGH_RISK_IMBL, reproducible."),
    ]
    cx = Inches(0.5)
    for head, body in callouts:
        rect(s, cx, cy, cw, ch, fill_rgb=SURFACE)
        rect(s, cx, cy, cw, Inches(0.06), fill_rgb=TEAL)
        text(s, cx + Inches(0.20), cy + Inches(0.15), cw - Inches(0.4),
             Inches(0.4), text=head, size=13, bold=True, color=TEAL)
        text(s, cx + Inches(0.20), cy + Inches(0.50), cw - Inches(0.4),
             Inches(0.7), text=body, size=10, color=SOFT)
        cx += cw + Inches(0.13)

    footer(s)


# ---------- Build ----------------------------------------------------------

def main():
    prs = new_prs()

    s01_hook(prs)
    s02_solution(prs)
    s03_architecture(prs)
    s04_differentiation(prs)
    s05_feasibility(prs)
    s06_research(prs)
    s07_demo(prs)
    s08_ai(prs)
    s09_hardware(prs)
    s10_team(prs)

    out = "/Users/subham/code/orca/ORCA_4.0_SIH_Deck.pptx"
    prs.save(out)
    print("Wrote", out, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    main()
