import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = "/Users/subham/Downloads/test_6slides.pptx"
OUTPUT_PATH = "/Users/subham/Downloads/ORCA_4.0_SIH2026_Extreme_Master_Deck.pptx"
LOGO_PATH = "/Users/subham/code/orca/sih_logo_extracted.png"

# ==============================================================================
# PALETTE: PALANTIR x MARITIME COMMAND x MISSION CONTROL (LIGHT & SCIENTIFIC)
# ==============================================================================
C_WHITE        = RGBColor(255, 255, 255)
C_SLATE_50     = RGBColor(248, 250, 252)   # #f8fafc - Base canvas
C_SLATE_100    = RGBColor(241, 245, 249)   # #f1f5f9 - Card fill
C_SLATE_200    = RGBColor(226, 232, 240)   # #e2e8f0 - Crisp border
C_SLATE_700    = RGBColor(51, 65, 85)      # #334155 - Body text
C_SLATE_900    = RGBColor(15, 23, 42)      # #0f172a - Dark bold text
C_MUTED        = RGBColor(100, 116, 139)   # #64748b - Subtitle / labels

C_NAVY_HERO    = RGBColor(10, 37, 64)      # #0a2540 - Deep maritime hero panels
C_NAVY_CARD    = RGBColor(15, 34, 64)      # #0f2240 - Navy inner cards
C_BLUE_PRIMARY = RGBColor(2, 132, 199)     # #0284c7 - Ocean blue accent
C_BLUE_LIGHT   = RGBColor(240, 249, 255)   # #f0f9ff - Ice blue fill
C_BLUE_BORDER  = RGBColor(186, 230, 253)   # #bae6fd - Ice blue border
C_CYAN_ACCENT  = RGBColor(6, 182, 212)     # #06b6d4 - Vibrant cyan

C_EMERALD      = RGBColor(16, 185, 129)    # #10b981 - Green Safe
C_EMERALD_BG   = RGBColor(236, 253, 245)   # #ecfdf5
C_EMERALD_BORDER = RGBColor(167, 243, 208)

C_AMBER        = RGBColor(217, 119, 6)     # #d97706 - Warning Amber
C_AMBER_BG     = RGBColor(254, 243, 199)   # #fef3c7
C_AMBER_BORDER = RGBColor(253, 230, 138)

C_ROSE         = RGBColor(225, 29, 72)     # #e11d48 - Danger Rose
C_ROSE_BG      = RGBColor(255, 241, 242)   # #fff1f2
C_ROSE_BORDER  = RGBColor(254, 205, 211)

NAV_STAGES = [
    "01 PROBLEM", "02 GAP", "03 SOLUTION", "04 DIFFERENTIATION", 
    "05 ARCHITECTURE", "06 INTELLIGENCE", "07 PROTOTYPE", 
    "08 SAFETY", "09 SCALE", "10 IMPACT"
]

def set_shape_flat(shape, fill_color, border_color=None, border_width=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

def create_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    set_shape_flat(shape, bg_color, border_color, border_width)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return shape, tf

def apply_official_header_footer(slide, slide_num, total_slides, section_badge, title_text, subtitle_text=""):
    # Top Team Oval Badge (Left: 0.36, top: 0.18)
    oval_shape, oval_tf = create_card(slide, Inches(0.36), Inches(0.18), Inches(1.35), Inches(0.72),
                                      bg_color=C_NAVY_HERO, border_color=None, shape_type=MSO_SHAPE.OVAL)
    p_ov = oval_tf.paragraphs[0]
    p_ov.text = "ORCA 4.0"
    p_ov.font.name = "Arial"
    p_ov.font.size = Pt(11)
    p_ov.font.bold = True
    p_ov.font.color.rgb = C_WHITE
    p_ov.alignment = PP_ALIGN.CENTER

    # Official SIH Top Right Logo Picture
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(10.70), Inches(0.02), Inches(2.46), Inches(1.10))

    # Center Title & Subtitle (Left: 1.85, width: 8.75)
    t_box = slide.shapes.add_textbox(Inches(1.85), Inches(0.10), Inches(8.75), Inches(0.85))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
    
    p_title = tf_t.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(17)
    p_title.font.bold = True
    p_title.font.color.rgb = C_NAVY_HERO

    if subtitle_text:
        p_sub = tf_t.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(9)
        p_sub.font.color.rgb = C_MUTED

    # Global 10-Stage Progress Ribbon with Active Highlight (top: 0.96, height: 0.28)
    ribbon_w = Inches(12.13)
    _, tf_rb = create_card(slide, Inches(0.6), Inches(0.96), ribbon_w, Inches(0.28),
                           bg_color=C_SLATE_100, border_color=C_SLATE_200, border_width=1.0)
    p_rb = tf_rb.paragraphs[0]
    p_rb.alignment = PP_ALIGN.CENTER
    
    active_idx = slide_num - 1
    for s_i, s_label in enumerate(NAV_STAGES):
        r_step = p_rb.add_run()
        r_step.text = s_label
        r_step.font.name = "Arial"
        r_step.font.size = Pt(7)
        r_step.font.bold = (s_i == active_idx)
        r_step.font.color.rgb = C_BLUE_PRIMARY if (s_i == active_idx) else C_MUTED
        if s_i < len(NAV_STAGES) - 1:
            r_sep = p_rb.add_run()
            r_sep.text = " ─ "
            r_sep.font.name = "Arial"
            r_sep.font.size = Pt(6.5)
            r_sep.font.color.rgb = C_SLATE_200

    # Official Bottom Footer Bar (top: 6.95, height: 0.40)
    btm_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), Inches(13.333), Inches(0.55))
    set_shape_flat(btm_bar, C_WHITE, C_SLATE_200, 0.5)

    # Footer Text
    ft_box = slide.shapes.add_textbox(Inches(4.5), Inches(7.02), Inches(4.33), Inches(0.35))
    tf_ft = ft_box.text_frame
    tf_ft.margin_left = tf_ft.margin_right = tf_ft.margin_top = tf_ft.margin_bottom = 0
    p_ft = tf_ft.paragraphs[0]
    p_ft.text = "@SIH Idea submission- Template  ·  ORCA 4.0"
    p_ft.font.name = "Arial"
    p_ft.font.size = Pt(8.5)
    p_ft.font.color.rgb = C_MUTED
    p_ft.alignment = PP_ALIGN.CENTER

    # Slide Number
    num_box = slide.shapes.add_textbox(Inches(11.0), Inches(7.02), Inches(1.7), Inches(0.35))
    tf_num = num_box.text_frame
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    p_num = tf_num.paragraphs[0]
    p_num.text = f"{slide_num:02d} / {total_slides:02d}"
    p_num.font.name = "Arial"
    p_num.font.size = Pt(9)
    p_num.font.bold = True
    p_num.font.color.rgb = C_NAVY_HERO
    p_num.alignment = PP_ALIGN.RIGHT

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    print("Building Extreme Master Deck (85% Visual / 15% Text) for ORCA 4.0...")

    # =========================================================================
    # SLIDE 1: HERO / TITLE (WHO?)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    print("Formatting Slide 1: Hero & Title...")

    # Header Bar
    top_bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.15))
    set_shape_flat(top_bar1, C_WHITE, C_SLATE_200, 0.5)

    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(10.70), Inches(0.02), Inches(2.46), Inches(1.10))

    t_box1 = slide1.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(9.8), Inches(0.85))
    tf_t1 = t_box1.text_frame
    tf_t1.margin_left = tf_t1.margin_right = tf_t1.margin_top = tf_t1.margin_bottom = 0
    p_sih_title = tf_t1.paragraphs[0]
    p_sih_title.text = "SMART INDIA HACKATHON 2026"
    p_sih_title.font.name = "Arial"
    p_sih_title.font.size = Pt(20)
    p_sih_title.font.bold = True
    p_sih_title.font.color.rgb = C_NAVY_HERO
    
    p_sih_sub = tf_t1.add_paragraph()
    p_sih_sub.text = "OFFICIAL IDEA PRESENTATION  ·  NATIONAL FINALS"
    p_sih_sub.font.name = "Arial"
    p_sih_sub.font.size = Pt(10)
    p_sih_sub.font.bold = True
    p_sih_sub.font.color.rgb = C_BLUE_PRIMARY

    # Left: Official Form Details Card (width: 5.4)
    _, tf_sih_form = create_card(slide1, Inches(0.6), Inches(1.30), Inches(5.4), Inches(5.50),
                                 bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_fh = tf_sih_form.paragraphs[0]
    p_fh.text = "PROPOSAL METADATA & TEAM REGISTRATION"
    p_fh.font.name = "Arial"
    p_fh.font.bold = True
    p_fh.font.size = Pt(10)
    p_fh.font.color.rgb = C_NAVY_HERO
    p_fh.space_after = Pt(8)

    sih_fields = [
        ("Problem Statement ID", "SIH26176"),
        ("Problem Statement Title", "AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform for Artisanal Fishermen & Coastal Fleets"),
        ("Sponsoring Ministry", "ISRO / Department of Space & INCOIS"),
        ("Theme", "Space Technology / Blue Economy / Disaster Management"),
        ("PS Category", "Software (with Edge Hardware compatibility)"),
        ("Team ID / Registered Name", "ORCA 4.0")
    ]
    for label, val in sih_fields:
        p_f = tf_sih_form.add_paragraph()
        p_f.space_after = Pt(4.5)
        r_lbl = p_f.add_run()
        r_lbl.text = f"{label}:\n"
        r_lbl.font.name = "Arial"
        r_lbl.font.bold = True
        r_lbl.font.size = Pt(8)
        r_lbl.font.color.rgb = C_NAVY_HERO
        r_val = p_f.add_run()
        r_val.text = val
        r_val.font.name = "Arial"
        r_val.font.size = Pt(9)
        r_val.font.color.rgb = C_SLATE_700

    # Right: Dark Navy ORCA 4.0 Showcase with 7 Connected Visual Nodes (width: 6.53)
    _, tf_hero = create_card(slide1, Inches(6.2), Inches(1.30), Inches(6.53), Inches(5.50),
                             bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.5)
    
    p_h_badge = tf_hero.paragraphs[0]
    p_h_badge.text = "MARITIME DECISION OPERATING SYSTEM"
    p_h_badge.font.name = "Arial"
    p_h_badge.font.size = Pt(9.5)
    p_h_badge.font.bold = True
    p_h_badge.font.color.rgb = C_CYAN_ACCENT
    p_h_badge.space_after = Pt(2)

    p_h_title = tf_hero.add_paragraph()
    p_h_title.text = "ORCA 4.0"
    p_h_title.font.name = "Arial"
    p_h_title.font.size = Pt(36)
    p_h_title.font.bold = True
    p_h_title.font.color.rgb = C_WHITE
    p_h_title.space_after = Pt(1)

    p_h_tag = tf_hero.add_paragraph()
    p_h_tag.text = "“FROM OCEAN DATA TO SAFER DECISIONS”"
    p_h_tag.font.name = "Arial"
    p_h_tag.font.size = Pt(12)
    p_h_tag.font.bold = True
    p_h_tag.font.color.rgb = RGBColor(56, 189, 248)
    p_h_tag.space_after = Pt(6)

    # 7 Connected Nodes Visual Grid inside Hero
    nodes_7 = [
        ("1. OCEAN", "Sea State & Waves", C_BLUE_PRIMARY),
        ("2. WEATHER", "Wind & Storms", C_CYAN_ACCENT),
        ("3. VESSEL", "Hull Stability Twin", RGBColor(56, 189, 248)),
        ("4. ROUTE", "Waypoints & Track", C_AMBER),
        ("5. ORCA", "Bio-Physical Core", C_WHITE),
        ("6. RISK", "Dynamic Risk Field", C_ROSE),
        ("7. DECISION", "GO / REROUTE / RETURN", C_EMERALD)
    ]
    
    p_n_hdr = tf_hero.add_paragraph()
    p_n_hdr.text = "7-NODE SYSTEM ARCHITECTURE PIPELINE:"
    p_n_hdr.font.name = "Arial"
    p_n_hdr.font.bold = True
    p_n_hdr.font.size = Pt(8.5)
    p_n_hdr.font.color.rgb = C_EMERALD
    p_n_hdr.space_after = Pt(4)

    for n_title, n_sub, n_col in nodes_7:
        p_node = tf_hero.add_paragraph()
        p_node.space_after = Pt(1.5)
        r_nt = p_node.add_run()
        r_nt.text = f"  [{n_title}]  "
        r_nt.font.name = "Arial"
        r_nt.font.bold = True
        r_nt.font.size = Pt(8)
        r_nt.font.color.rgb = n_col
        r_ns = p_node.add_run()
        r_ns.text = f"➔  {n_sub}"
        r_ns.font.name = "Arial"
        r_ns.font.size = Pt(7.5)
        r_ns.font.color.rgb = RGBColor(203, 213, 225)

    p_val = tf_hero.add_paragraph()
    p_val.space_after = Pt(2)
    r_v = p_val.add_run()
    r_v.text = "\n[VESSEL-AWARE]  ·  [ROUTE-AWARE]  ·  [SAFETY-FIRST]"
    r_v.font.name = "Arial"
    r_v.font.bold = True
    r_v.font.size = Pt(8.5)
    r_v.font.color.rgb = C_CYAN_ACCENT

    # Footer on Slide 1
    btm_bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), Inches(13.333), Inches(0.55))
    set_shape_flat(btm_bar1, C_WHITE, C_SLATE_200, 0.5)
    ft_box1 = slide1.shapes.add_textbox(Inches(4.5), Inches(7.02), Inches(4.33), Inches(0.35))
    tf_ft1 = ft_box1.text_frame
    p_ft1 = tf_ft1.paragraphs[0]
    p_ft1.text = "@SIH Idea submission- Template  ·  ORCA 4.0"
    p_ft1.font.name = "Arial"
    p_ft1.font.size = Pt(8.5)
    p_ft1.font.color.rgb = C_MUTED
    p_ft1.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: THE PROBLEM & GAP (WHY?)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide2, 2, 10, "Problem Discovery", 
                                 "“THE OCEAN IS FORECAST. THE VOYAGE IS NOT.”", 
                                 "Visualizing the fatal decision gap: Why raw forecast numbers fail to answer what this vessel should do now.")
    print("Formatting Slide 2: The Problem & Gap...")

    # Hero Thesis Callout Box (Huge statement)
    _, tf_s2_th = create_card(slide2, Inches(0.6), Inches(1.30), Inches(12.13), Inches(0.52),
                              bg_color=C_ROSE_BG, border_color=C_ROSE, border_width=1.0)
    p_s2_th = tf_s2_th.paragraphs[0]
    r_th1 = p_s2_th.add_run()
    r_th1.text = "CORE PARADIGM SHIFT:  FORECAST ≠ OPERATIONAL RISK.  "
    r_th1.font.name = "Arial"
    r_th1.font.bold = True
    r_th1.font.size = Pt(9.5)
    r_th1.font.color.rgb = C_ROSE
    r_th2 = p_s2_th.add_run()
    r_th2.text = "Weather apps show what the sea is doing. ORCA computes what it means for THIS specific hull on THIS route."
    r_th2.font.name = "Arial"
    r_th2.font.size = Pt(8.5)
    r_th2.font.color.rgb = C_SLATE_900

    # 3-Zone Visual Composition: Left (8 Silos), Center (Broadcast Wall), Right (Decision Gap)
    # Zone 1: 8 Fragmented Telemetry Nodes (width: 4.2)
    _, tf_s2_z1 = create_card(slide2, Inches(0.6), Inches(1.88), Inches(4.2), Inches(3.75),
                              bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_z1_h = tf_s2_z1.paragraphs[0]
    p_z1_h.text = "8 DISCONNECTED DATA STREAMS"
    p_z1_h.font.name = "Arial"
    p_z1_h.font.bold = True
    p_z1_h.font.size = Pt(9.5)
    p_z1_h.font.color.rgb = C_NAVY_HERO
    p_z1_h.space_after = Pt(3)

    s2_silos = [
        ("🌊 WAVES (Hs, Swell Tp)", "INCOIS spectral wave grids"),
        ("💨 WIND (Vel, Gusts)", "IMD surface atmospheric vectors"),
        ("🌀 CYCLONES (IMD Alerts)", "Active storm coordinates"),
        ("🌊 CURRENTS (u, v)", "Ocean drift velocity fields"),
        ("🛰️ SATELLITES (SST)", "Thermal front boundaries"),
        ("🚢 VESSEL PROFILE", "Lwl, Draft, Beam geometry"),
        ("📡 AIS TELEMETRY", "Fleet positions & collisions"),
        ("🗺️ VOYAGE ROUTE", "Target dock & waypoints")
    ]
    for s_title, s_desc in s2_silos:
        p_s = tf_s2_z1.add_paragraph()
        p_s.space_after = Pt(1.5)
        r1 = p_s.add_run()
        r1.text = f"• {s_title}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p_s.add_run()
        r2.text = s_desc
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_MUTED

    # Zone 2: The Generic Forecast Wall (width: 3.5)
    _, tf_s2_z2 = create_card(slide2, Inches(4.95), Inches(1.88), Inches(3.5), Inches(3.75),
                              bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_z2_h = tf_s2_z2.paragraphs[0]
    p_z2_h.text = "TODAY: GENERIC FORECAST"
    p_z2_h.font.name = "Arial"
    p_z2_h.font.bold = True
    p_z2_h.font.size = Pt(9.5)
    p_z2_h.font.color.rgb = C_ROSE
    p_z2_h.space_after = Pt(6)

    p_z2_msg = tf_s2_z2.add_paragraph()
    p_z2_msg.text = "⚠️ BROADCAST SMS:\n“Waves 2.5m in South Arabian Sea. Wind 22 kts.”\n"
    p_z2_msg.font.name = "Arial"
    p_z2_msg.font.bold = True
    p_z2_msg.font.size = Pt(8.5)
    p_z2_msg.font.color.rgb = C_SLATE_900
    p_z2_msg.space_after = Pt(6)

    p_z2_flaws = tf_s2_z2.add_paragraph()
    p_z2_flaws.text = "❌ ZERO VESSEL CONTEXT:\nIdentical SMS sent to a 40m steel trawler and a 7m wooden canoe.\n\n❌ ZERO ROUTE AWARENESS:\nCannot evaluate encounter wave angle or shallow breaker reef."
    p_z2_flaws.font.name = "Arial"
    p_z2_flaws.font.size = Pt(7.5)
    p_z2_flaws.font.color.rgb = C_SLATE_700

    # Zone 3: The Confused Decision State (width: 4.13)
    _, tf_s2_z3 = create_card(slide2, Inches(8.6), Inches(1.88), Inches(4.13), Inches(3.75),
                              bg_color=C_ROSE_BG, border_color=C_ROSE, border_width=1.0)
    p_z3_h = tf_s2_z3.paragraphs[0]
    p_z3_h.text = "OPERATIONAL DECISION GAP"
    p_z3_h.font.name = "Arial"
    p_z3_h.font.bold = True
    p_z3_h.font.size = Pt(9.5)
    p_z3_h.font.color.rgb = C_ROSE
    p_z3_h.space_after = Pt(4)

    p_z3_q = tf_s2_z3.add_paragraph()
    p_z3_q.text = "THE FISHERMAN'S DILEMMA:\n"
    p_z3_q.font.name = "Arial"
    p_z3_q.font.bold = True
    p_z3_q.font.size = Pt(8.5)
    p_z3_q.font.color.rgb = C_NAVY_HERO

    q_items = [
        ("• Should I sail today?", "Livelihood vs. Drowning Risk"),
        ("• Will my hull capsize?", "Unknown beam-sea resonance"),
        ("• Which route is safe?", "Blind navigation beyond 12 NM"),
        ("• Where are the fish?", "Blind fuel burn hunting for PFZ")
    ]
    for qn, qd in q_items:
        p_qi = tf_s2_z3.add_paragraph()
        p_qi.space_after = Pt(2)
        r1 = p_qi.add_run()
        r1.text = f"{qn} "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_ROSE
        r2 = p_qi.add_run()
        r2.text = f"➔ {qd}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Bottom: The ORCA Paradigm Equation (top: 5.72, height: 1.05)
    _, tf_s2_btm = create_card(slide2, Inches(0.6), Inches(5.72), Inches(12.13), Inches(1.05),
                               bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.5)
    p_s2_bh = tf_s2_btm.paragraphs[0]
    p_s2_bh.text = "THE ORCA SOLUTION EQUATION: DATA + CONTEXT + PHYSICS + ROUTE = ACTIONABLE DECISION"
    p_s2_bh.font.name = "Arial"
    p_s2_bh.font.bold = True
    p_s2_bh.font.size = Pt(9)
    p_s2_bh.font.color.rgb = C_NAVY_HERO
    
    p_s2_bb = tf_s2_btm.add_paragraph()
    p_s2_bb.text = "ORCA fuses 8 data streams into a 0.083° spatial grid, instantiates a parametric vessel digital twin, evaluates deterministic wave encounter physics (<10ms), and computes an optimized A* route."
    p_s2_bb.font.name = "Arial"
    p_s2_bb.font.size = Pt(7.5)
    p_s2_bb.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 3: THE ORCA CONVERGENCE ENGINE (WHAT?)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide3, 3, 10, "Solution Overview", 
                                 "“ORCA TURNS SEA CONDITIONS INTO A VOYAGE DECISION”", 
                                 "A central computational engine fusing raw environmental data, vessel physics, and dynamic route risk.")
    print("Formatting Slide 3: The ORCA Engine...")

    # Left: 8 Converging Telemetry Feeds (width: 3.0)
    _, tf_s3_in = create_card(slide3, Inches(0.6), Inches(1.30), Inches(3.0), Inches(4.35),
                              bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_s3_ih = tf_s3_in.paragraphs[0]
    p_s3_ih.text = "8 CONVERGING INPUTS"
    p_s3_ih.font.name = "Arial"
    p_s3_ih.font.bold = True
    p_s3_ih.font.size = Pt(9)
    p_s3_ih.font.color.rgb = C_BLUE_PRIMARY
    p_s3_ih.space_after = Pt(2)

    in_8 = [
        ("Waves (Hs, Tp)", "INCOIS ERDDAP"),
        ("Wind & Gusts", "IMD Station Grid"),
        ("Cyclone Tracks", "IMD Storm Center"),
        ("Ocean Currents", "Copernicus CMEMS"),
        ("SST Fronts", "ISRO MOSDAC"),
        ("AIS Transponders", "Fleet Positions"),
        ("Vessel Hull Profile", "Registry Twin"),
        ("Target Waypoints", "Port Dock Plan")
    ]
    for in_name, in_src in in_8:
        p_in = tf_s3_in.add_paragraph()
        p_in.space_after = Pt(1.5)
        r1 = p_in.add_run()
        r1.text = f"• {in_name}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p_in.add_run()
        r2.text = f"   Source: {in_src}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_MUTED

    # Center: Central ORCA Engine (4 Computational Layers: width: 5.6)
    _, tf_s3_eng = create_card(slide3, Inches(3.75), Inches(1.30), Inches(5.6), Inches(4.35),
                               bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.5)
    p_s3_eh = tf_s3_eng.paragraphs[0]
    p_s3_eh.text = "⚡ ORCA 4.0 COMPUTATIONAL ENGINE"
    p_s3_eh.font.name = "Arial"
    p_s3_eh.font.bold = True
    p_s3_eh.font.size = Pt(10.5)
    p_s3_eh.font.color.rgb = C_CYAN_ACCENT
    p_s3_eh.space_after = Pt(4)

    layers_4 = [
        ("LAYER 1: ENVIRONMENTAL DIGITAL STATE", "Spatial H3 Normalization (0.083° / ~9km) + Multi-Source Temporal Alignment.", C_BLUE_LIGHT),
        ("LAYER 2: PARAMETRIC VESSEL DIGITAL TWIN", "Hydrodynamic hull modeling: Length (Lwl), Beam (B), Draft (Td), Freeboard.", RGBColor(56, 189, 248)),
        ("LAYER 3: DETERMINISTIC PHYSICS & HSI ML", "Critical wave threshold (Hcrit = 0.6·L) + XGBoost multi-species habitat index.", C_EMERALD),
        ("LAYER 4: A* DYNAMIC ROUTE RISK MESH", "Time-varying 2D cost surface pathfinder avoiding high-risk breaker pockets.", C_AMBER)
    ]
    for l_t, l_d, l_c in layers_4:
        p_l = tf_s3_eng.add_paragraph()
        p_l.space_after = Pt(3.5)
        r1 = p_l.add_run()
        r1.text = f"{l_t}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = l_c
        r2 = p_l.add_run()
        r2.text = l_d
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = RGBColor(203, 213, 225)

    # Right: 3 Actionable Verdicts (width: 3.33)
    _, tf_s3_out = create_card(slide3, Inches(9.5), Inches(1.30), Inches(3.23), Inches(4.35),
                               bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_s3_oh = tf_s3_out.paragraphs[0]
    p_s3_oh.text = "3 ACTIONABLE OUTPUTS"
    p_s3_oh.font.name = "Arial"
    p_s3_oh.font.bold = True
    p_s3_oh.font.size = Pt(9)
    p_s3_oh.font.color.rgb = C_EMERALD
    p_s3_oh.space_after = Pt(4)

    out_3 = [
        ("🟢 GO / SAFE DEPARTURE", "Sea state within vessel limits. Direct optimal route clear.", C_EMERALD_BG, C_EMERALD),
        ("🟡 CAUTION / REROUTE", "Wave danger detected. A* dynamic diversion generated.", C_AMBER_BG, C_AMBER),
        ("🔴 DANGER / STAY ASHORE", "Waves > Hcrit or cyclone alert. Deterministic hard stop.", C_ROSE_BG, C_ROSE)
    ]
    for ot, od, obg, oc in out_3:
        p_o = tf_s3_out.add_paragraph()
        p_o.space_after = Pt(4)
        r1 = p_o.add_run()
        r1.text = f"{ot}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = oc
        r2 = p_o.add_run()
        r2.text = od
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Bottom Positioning Banner: "ORCA IS THE DECISION LAYER"
    _, tf_s3_pos = create_card(slide3, Inches(0.6), Inches(5.75), Inches(12.13), Inches(1.05),
                               bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER, border_width=1.0)
    p_s3_ph = tf_s3_pos.paragraphs[0]
    p_s3_ph.text = "POSITIONING: “ORCA IS THE DECISION LAYER BETWEEN DATA AND ACTION”"
    p_s3_ph.font.name = "Arial"
    p_s3_ph.font.bold = True
    p_s3_ph.font.size = Pt(10)
    p_s3_ph.font.color.rgb = C_NAVY_HERO
    
    p_s3_pb = tf_s3_pos.add_paragraph()
    p_s3_pb.text = "Existing platforms stop at displaying raw maps and forecasts. ORCA provides the automated computational bridge that converts multi-agency oceanographic data into safe, vessel-specific voyage clearances."
    p_s3_pb.font.name = "Arial"
    p_s3_pb.font.size = Pt(7.5)
    p_s3_pb.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 4: VISUAL DIFFERENTIATION & MATRIX (DIFFERENT FROM WHAT?)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide4, 4, 10, "Competitive Landscape", 
                                 "“OTHERS SHOW THE SEA. ORCA UNDERSTANDS THE VOYAGE.”", 
                                 "Side-by-side comparison: Why weather portals stop at data, while ORCA delivers context and deterministic physics.")
    print("Formatting Slide 4: Visual Differentiation...")

    # Top: Two-Part Architectural Shift Comparison (top: 1.30, height: 1.65)
    _, tf_s4_top = create_card(slide4, Inches(0.6), Inches(1.30), Inches(12.13), Inches(1.65),
                               bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_s4_th = tf_s4_top.paragraphs[0]
    p_s4_th.text = "ARCHITECTURAL COMPARISON: GENERIC FORECAST vs. ORCA DECISION INTELLIGENCE"
    p_s4_th.font.name = "Arial"
    p_s4_th.font.bold = True
    p_s4_th.font.size = Pt(9.5)
    p_s4_th.font.color.rgb = C_NAVY_HERO
    p_s4_th.space_after = Pt(2)

    p_s4_cmp = tf_s4_top.add_paragraph()
    r1 = p_s4_cmp.add_run()
    r1.text = "• EXISTING APPROACHES:  "
    r1.font.name = "Arial"
    r1.font.bold = True
    r1.font.size = Pt(8)
    r1.font.color.rgb = C_ROSE
    r2 = p_s4_cmp.add_run()
    r2.text = "Raw Sea State  ➔  Regional Broadcast SMS  ➔  Human Guesswork (Blind to vessel size & route)\n"
    r2.font.name = "Arial"
    r2.font.size = Pt(8)
    r2.font.color.rgb = C_SLATE_700

    r3 = p_s4_cmp.add_run()
    r3.text = "• ORCA 4.0 PLATFORM:   "
    r3.font.name = "Arial"
    r3.font.bold = True
    r3.font.size = Pt(8)
    r3.font.color.rgb = C_EMERALD
    r4 = p_s4_cmp.add_run()
    r4.text = "Sea State  +  Vessel Hull Twin  +  Route Track  +  Deterministic Physics  ➔  Actionable Verdict"
    r4.font.name = "Arial"
    r4.font.bold = True
    r4.font.size = Pt(8)
    r4.font.color.rgb = C_NAVY_HERO

    # Bottom: Compact Comparison Matrix (top: 3.05, height: 3.75)
    _, tf_s4_m = create_card(slide4, Inches(0.6), Inches(3.05), Inches(12.13), Inches(3.75),
                             bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_s4_mh = tf_s4_m.paragraphs[0]
    p_s4_mh.text = "CAPABILITY EVALUATION MATRIX"
    p_s4_mh.font.name = "Arial"
    p_s4_mh.font.bold = True
    p_s4_mh.font.size = Pt(9.5)
    p_s4_mh.font.color.rgb = C_NAVY_HERO
    p_s4_mh.space_after = Pt(2)

    comp_rows_s4 = [
        ("Capability / Dimension", "Weather Forecast Apps", "Marine Port Portals", "ORCA 4.0 Platform"),
        ("Weather & Wave Forecast", "✓ Regional Broadcast", "✓ Coarse Marine Grid", "✓ High-Res (0.083° H3 Spatial Grid)"),
        ("Cyclone & Hazard Alerts", "✓ Generic SMS Alert", "✓ Port Warning Bulletin", "✓ Dynamic Real-Time Geo-Fencing"),
        ("Vessel Hull & Draft Profile", "— Ignored", "— Ignored", "✓ Parametric Vessel Digital Twin"),
        ("Wave-Vessel Dynamic Physics", "— None", "— None", "✓ Hcrit Capsizing & Roll Resonance Math"),
        ("Dynamic Route Pathfinder", "— None (Straight Line)", "△ Fixed Channel Lanes", "✓ A* Risk-Cost Avoidance Routing"),
        ("Deterministic Safety Guard", "— None", "△ Manual Port Officer", "✓ Non-Bypassable Rule Shield (<10ms)"),
        ("Multi-Species Habitat (PFZ)", "— None", "△ Static PFZ PDF Maps", "✓ XGBoost Habitat Suitability Index"),
        ("Offshore Reach Beyond 12 NM", "— 0% (GSM Cellular Blind)", "△ Expensive Satellite Satcom", "✓ Multi-Hop LoRa Mesh (up to 50 km)")
    ]
    for idx, (c1, c2, c3, c4) in enumerate(comp_rows_s4):
        p = tf_s4_m.add_paragraph()
        p.space_after = Pt(1.5)
        r1 = p.add_run()
        r1.text = f"{c1.ljust(30)} "
        r1.font.name = "Arial"
        r1.font.bold = True if idx == 0 else False
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_NAVY_HERO if idx == 0 else C_SLATE_900

        r2 = p.add_run()
        r2.text = f"{c2.ljust(25)} "
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_ROSE if "—" in c2 else (C_AMBER if "△" in c2 else C_MUTED)

        r3 = p.add_run()
        r3.text = f"{c3.ljust(25)} "
        r3.font.name = "Arial"
        r3.font.size = Pt(7)
        r3.font.color.rgb = C_ROSE if "—" in c3 else (C_AMBER if "△" in c3 else C_MUTED)

        r4 = p.add_run()
        r4.text = c4
        r4.font.name = "Arial"
        r4.font.bold = True
        r4.font.size = Pt(7.5)
        r4.font.color.rgb = C_EMERALD if "✓" in c4 else C_NAVY_HERO

    # =========================================================================
    # SLIDE 5: CTO-LEVEL SYSTEM ARCHITECTURE (HOW?)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide5, 5, 10, "System Architecture", 
                                 "“ONE PIPELINE. ONE WORLD MODEL.”", 
                                 "Layered CTO-grade architecture: Telemetry Ingestion ➔ Normalization ➔ Dual Digital State ➔ Tripartite Engine ➔ Safety Shield.")
    print("Formatting Slide 5: CTO-Level Architecture...")

    # Tier 1: Data Ingestion (top: 1.30, height: 0.82)
    _, tf_l1 = create_card(slide5, Inches(0.6), Inches(1.30), Inches(12.13), Inches(0.82),
                           bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_l1_h = tf_l1.paragraphs[0]
    p_l1_h.text = "TIER 1: MULTI-SOURCE TELEMETRY INGESTION LAYER"
    p_l1_h.font.name = "Arial"
    p_l1_h.font.bold = True
    p_l1_h.font.size = Pt(8.5)
    p_l1_h.font.color.rgb = C_BLUE_PRIMARY
    
    p_l1_b = tf_l1.add_paragraph()
    p_l1_b.text = "📡 INCOIS ERDDAP (Wave Spectra)  |  📡 IMD Alerts (Cyclones)  |  📡 Copernicus CMEMS (0.083° Physical Grids)  |  📡 ISRO MOSDAC (SST)  |  📡 AIS Feeds  |  📡 NMEA Sensors"
    p_l1_b.font.name = "Arial"
    p_l1_b.font.size = Pt(8)
    p_l1_b.font.color.rgb = C_SLATE_900

    # Tier 2: Validation, Provenance & Normalization (top: 2.20, height: 0.88)
    _, tf_l2 = create_card(slide5, Inches(0.6), Inches(2.20), Inches(12.13), Inches(0.88),
                           bg_color=C_SLATE_50, border_color=C_BLUE_BORDER, border_width=1.0)
    p_l2_h = tf_l2.paragraphs[0]
    p_l2_h.text = "TIER 2: DATA VALIDATION, PROVENANCE & DUAL DIGITAL STATE"
    p_l2_h.font.name = "Arial"
    p_l2_h.font.bold = True
    p_l2_h.font.size = Pt(8.5)
    p_l2_h.font.color.rgb = C_NAVY_HERO
    
    p_l2_b = tf_l2.add_paragraph()
    p_l2_b.text = "• Environmental Digital State: Spatial H3 Indexing (0.083° / ~9km) + Temporal Alignment + Sensor Calibration (<50ms)\n• Parametric Vessel Digital Twin: Dynamic hydrodynamics incorporating Hull Length (Lwl), Draft (Td), Beam (B), Freeboard & Engine Power"
    p_l2_b.font.name = "Arial"
    p_l2_b.font.size = Pt(7.5)
    p_l2_b.font.color.rgb = C_SLATE_700

    # Tier 3: Tripartite Computing Core (top: 3.16, height: 1.40)
    _, tf_l3 = create_card(slide5, Inches(0.6), Inches(3.16), Inches(12.13), Inches(1.40),
                           bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.0)
    p_l3_h = tf_l3.paragraphs[0]
    p_l3_h.text = "TIER 3: TRIPARTITE COMPUTING CORE (STRICT ARCHITECTURAL SEPARATION)"
    p_l3_h.font.name = "Arial"
    p_l3_h.font.bold = True
    p_l3_h.font.size = Pt(9)
    p_l3_h.font.color.rgb = C_CYAN_ACCENT
    p_l3_h.space_after = Pt(1.5)

    p_l3_b1 = tf_l3.add_paragraph()
    p_l3_b1.text = "⚙️ DETERMINISTIC PHYSICS: Wave encounter angle math, Critical wave height (Hcrit = 0.6 · L), wave steepness ratio & roll resonance."
    p_l3_b1.font.name = "Arial"
    p_l3_b1.font.bold = True
    p_l3_b1.font.size = Pt(8)
    p_l3_b1.font.color.rgb = C_EMERALD

    p_l3_b2 = tf_l3.add_paragraph()
    p_l3_b2.text = "🧠 PREDICTIVE ML LAYER: Multi-Species Habitat Suitability Index (HSI XGBoost), 2D Sobel thermal fronts & 1,000-particle SAR drift."
    p_l3_b2.font.name = "Arial"
    p_l3_b2.font.bold = True
    p_l3_b2.font.size = Pt(8)
    p_l3_b2.font.color.rgb = RGBColor(56, 189, 248)

    p_l3_b3 = tf_l3.add_paragraph()
    p_l3_b3.text = "📍 DYNAMIC OPTIMIZATION: Modified A* navigational pathfinder navigating through time-varying 2D risk-cost surface meshes."
    p_l3_b3.font.name = "Arial"
    p_l3_b3.font.bold = True
    p_l3_b3.font.size = Pt(8)
    p_l3_b3.font.color.rgb = C_AMBER

    # Tier 4: Non-Bypassable Safety Circuit Breaker Shield (top: 4.64, height: 0.85)
    _, tf_l4 = create_card(slide5, Inches(0.6), Inches(4.64), Inches(12.13), Inches(0.85),
                           bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.5)
    p_l4_h = tf_l4.paragraphs[0]
    p_l4_h.text = "TIER 4: NON-BYPASSABLE SAFETY CIRCUIT BREAKER (<10ms Pure Python, Zero Hallucination)"
    p_l4_h.font.name = "Arial"
    p_l4_h.font.bold = True
    p_l4_h.font.size = Pt(8.5)
    p_l4_h.font.color.rgb = C_EMERALD
    
    p_l4_b = tf_l4.add_paragraph()
    p_l4_b.text = "All physical safety checks and active cyclone overrides execute as deterministic compiled code. LLM models sit strictly downstream for plain-language vernacular voice translation only."
    p_l4_b.font.name = "Arial"
    p_l4_b.font.size = Pt(7.5)
    p_l4_b.font.color.rgb = C_SLATE_900

    # Tier 5: Operational Dispatch (top: 5.57, height: 1.22)
    _, tf_l5 = create_card(slide5, Inches(0.6), Inches(5.57), Inches(12.13), Inches(1.22),
                           bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_l5_h = tf_l5.paragraphs[0]
    p_l5_h.text = "TIER 5: MULTI-CHANNEL OPERATIONAL DISPATCH & SIDEBAR TELEMETRY"
    p_l5_h.font.name = "Arial"
    p_l5_h.font.bold = True
    p_l5_h.font.size = Pt(8.5)
    p_l5_h.font.color.rgb = C_NAVY_HERO
    
    p_l5_b = tf_l5.add_paragraph()
    p_l5_b.text = "• Fisherman PWA & LoRa Voice Device: Actionable safety verdict dial, optimized 4-point waypoints, offline voice guidance.\n• Fleet Web & Coast Guard Command: Real-time dark-fleet AIS anomaly detection, harbor fleet congestion & SAR drift tracking.\n• Production Resilience Sidebar: Automated TTL Freshness Checks (<30m), Provider Failover Cascade & Provenance Weighting."
    p_l5_b.font.name = "Arial"
    p_l5_b.font.size = Pt(7.5)
    p_l5_b.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 6: INTELLIGENCE ENGINE & AI GOVERNANCE (HOW INTELLIGENT?)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide6, 6, 10, "Intelligence Separation", 
                                 "“PHYSICS + AI + OPTIMIZATION. EACH DOES WHAT IT IS GOOD AT.”", 
                                 "Strict architectural separation: Physics constrains, ML predicts, A* optimizes, and deterministic safety governs.")
    print("Formatting Slide 6: Intelligence Engine...")

    tri_data = [
        ("⚙️ DETERMINISTIC PHYSICS", 
         "“KNOWS THE RULES (CONSTRAIN)”", 
         [("Critical Wave Limit", "Hcrit ≈ 0.6 · Lhull threshold prevents capsizing"),
          ("Wave Steepness Ratio", "S = Hs / λ calculated dynamically from wave spectra"),
          ("Beam-Sea Encounter", "Identifies roll resonance & parametric broaching risk"),
          ("Geofenced Boundaries", "Hard stop on Marine Protected Areas & international borders")],
         C_EMERALD_BG, C_EMERALD),
        ("🧠 PREDICTIVE MACHINE LEARNING", 
         "“FINDS PATTERNS (PREDICT)”", 
         [("Multi-Species HSI", "XGBoost model predicts Bangda, Surmai & Pomfret zones"),
          ("Sobel Thermal Fronts", "2D spatial filtering identifies productive SST boundaries"),
          ("SAR Anomaly Matching", "Flags dark vessels turning off AIS transponders"),
          ("Monte Carlo SAR Drift", "1,000-particle stochastic simulation for rescue targets")],
         C_BLUE_LIGHT, C_BLUE_PRIMARY),
        ("📍 GRAPH OPTIMIZATION (A*)", 
         "“FINDS THE PATH (OPTIMIZE)”", 
         [("2D Risk Cost Surface", "Dynamic mesh where hazardous sea cells have infinite cost"),
          ("4-Point Waypoints", "Generates fuel-optimal path circumnavigating storm pockets"),
          ("Current Assist Optimization", "Leverages ocean surface drift velocity to save diesel"),
          ("Harbor Dock Allocation", "Optimizes arrival harbor based on catch ROI & congestion")],
         C_AMBER_BG, C_AMBER)
    ]

    for idx, (t_title, t_sub, t_items, bg_c, b_col) in enumerate(tri_data):
        _, tf_tri = create_card(slide6, Inches(0.6 + idx * 4.14), Inches(1.30), Inches(3.85), Inches(4.35),
                                bg_color=bg_c, border_color=b_col, border_width=1.0)
        p_th = tf_tri.paragraphs[0]
        p_th.text = t_title
        p_th.font.name = "Arial"
        p_th.font.bold = True
        p_th.font.size = Pt(9.5)
        p_th.font.color.rgb = b_col
        p_th.space_after = Pt(1.5)

        p_ts = tf_tri.add_paragraph()
        p_ts.text = t_sub
        p_ts.font.name = "Arial"
        p_ts.font.bold = True
        p_ts.font.size = Pt(7.5)
        p_ts.font.color.rgb = C_NAVY_HERO
        p_ts.space_after = Pt(5)

        for iname, idesc in t_items:
            p_i = tf_tri.add_paragraph()
            p_i.space_after = Pt(3)
            r1 = p_i.add_run()
            r1.text = f"• {iname}: "
            r1.font.name = "Arial"
            r1.font.bold = True
            r1.font.size = Pt(8)
            r1.font.color.rgb = C_SLATE_900
            r2 = p_i.add_run()
            r2.text = idesc
            r2.font.name = "Arial"
            r2.font.size = Pt(7.5)
            r2.font.color.rgb = C_SLATE_700

    # Bottom Banner: AI Governance Directive (Red Boundary)
    _, tf_llm = create_card(slide6, Inches(0.6), Inches(5.72), Inches(12.13), Inches(1.10),
                            bg_color=C_NAVY_HERO, border_color=C_ROSE, border_width=1.5)
    p_llm_h = tf_llm.paragraphs[0]
    p_llm_h.text = "AI SAFETY GOVERNANCE: “INTELLIGENCE IS ALLOWED. SAFETY OVERRIDE IS NOT.”"
    p_llm_h.font.name = "Arial"
    p_llm_h.font.bold = True
    p_llm_h.font.size = Pt(9)
    p_llm_h.font.color.rgb = C_CYAN_ACCENT
    p_llm_h.space_after = Pt(2)

    p_llm_b = tf_llm.add_paragraph()
    r_l1 = p_llm_b.add_run()
    r_l1.text = "STRICT SAFETY BOUNDARIES:  "
    r_l1.font.name = "Arial"
    r_l1.font.bold = True
    r_l1.font.size = Pt(7.5)
    r_l1.font.color.rgb = C_ROSE
    r_l2 = p_llm_b.add_run()
    r_l2.text = "LLM ──X──> SAFETY  |  ML ──X──> OVERRIDE  |  A* ──X──> PHYSICAL LIMIT\n"
    r_l2.font.name = "Arial"
    r_l2.font.bold = True
    r_l2.font.size = Pt(7.5)
    r_l2.font.color.rgb = C_WHITE

    r_l3 = p_llm_b.add_run()
    r_l3.text = "Downstream Role: "
    r_l3.font.name = "Arial"
    r_l3.font.bold = True
    r_l3.font.size = Pt(7.5)
    r_l3.font.color.rgb = C_EMERALD
    r_l4 = p_llm_b.add_run()
    r_l4.text = "LLM is strictly constrained to voice intent parsing (Speech ➔ JSON) and multilingual explanation synthesis (JSON ➔ Vernacular Audio)."
    r_l4.font.name = "Arial"
    r_l4.font.size = Pt(7.5)
    r_l4.font.color.rgb = RGBColor(203, 213, 225)

    # =========================================================================
    # SLIDE 7: LIVE PROTOTYPE & DYNAMIC DECISION LOOP (DOES IT WORK?)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide7, 7, 10, "Prototype Validation", 
                                 "“FROM LOCATION TO DECISION: LIVE COMMAND LOOP”", 
                                 "Live simulation loop demonstrating how ORCA processes a 9m FRP craft navigating from Kanyakumari.")
    print("Formatting Slide 7: Live Prototype & Command Center...")

    # Left: Telemetry Input & Vessel Profile Panel (width: 3.5)
    _, tf_s7_l = create_card(slide7, Inches(0.6), Inches(1.30), Inches(3.5), Inches(5.55),
                             bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_s7_lh = tf_s7_l.paragraphs[0]
    p_s7_lh.text = "1. LIVE TELEMETRY & VESSEL PROFILE"
    p_s7_lh.font.name = "Arial"
    p_s7_lh.font.bold = True
    p_s7_lh.font.size = Pt(9)
    p_s7_lh.font.color.rgb = C_NAVY_HERO
    p_s7_lh.space_after = Pt(4)

    tel_items = [
        ("Location", "Kanyakumari Harbor (8.08°N, 77.55°E)"),
        ("Vessel Hull", "9m FRP Artisanal Gillnetter"),
        ("Wave Height (Hs)", "2.8 m [Open-Meteo Feed]"),
        ("Swell Period (Tp)", "8.2 s [INCOIS Buoy]"),
        ("Wind Velocity", "24 kts (Gusts: 31 kts)"),
        ("Critical Wave (Hcrit)", "5.4 m (Hcrit = 0.6 · 9m)"),
        ("Physics Result", "Hcrit > Hs (Safe from direct capsize)")
    ]
    for tn, tv in tel_items:
        p_t = tf_s7_l.add_paragraph()
        p_t.space_after = Pt(2.5)
        r1 = p_t.add_run()
        r1.text = f"• {tn}:\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p_t.add_run()
        r2.text = f"  {tv}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Center: Command Map Visualization (width: 4.8)
    _, tf_s7_map = create_card(slide7, Inches(4.3), Inches(1.30), Inches(4.8), Inches(3.50),
                               bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER, border_width=1.0)
    p_s7_mh = tf_s7_map.paragraphs[0]
    p_s7_mh.text = "2. MARITIME COMMAND MAP: DYNAMIC ROUTE"
    p_s7_mh.font.name = "Arial"
    p_s7_mh.font.bold = True
    p_s7_mh.font.size = Pt(9)
    p_s7_mh.font.color.rgb = C_BLUE_PRIMARY
    p_s7_mh.space_after = Pt(4)

    p_map_v = tf_s7_map.add_paragraph()
    p_map_v.text = "📍 KANYAKUMARI HARBOR (Departure Point)\n"
    p_map_v.font.name = "Arial"
    p_map_v.font.bold = True
    p_map_v.font.size = Pt(8)
    p_map_v.font.color.rgb = C_NAVY_HERO

    p_map_d = tf_s7_map.add_paragraph()
    p_map_d.text = "   |─── ❌ Direct Route: Intersects Shallow Reef Breaker Zone (Risk Score: 0.88)\n   └─── ✅ ORCA A* Diversion: +3.2 NM East (Risk Score: 0.24, Fuel Est: 14.2 L)\n"
    p_map_d.font.name = "Arial"
    p_map_d.font.size = Pt(7.5)
    p_map_d.font.color.rgb = C_SLATE_700

    p_map_e = tf_s7_map.add_paragraph()
    p_map_e.text = "🏁 PELAGIC TARGET ZONE (HSI Score: 0.82 High Yield)"
    p_map_e.font.name = "Arial"
    p_map_e.font.bold = True
    p_map_e.font.size = Pt(8)
    p_map_e.font.color.rgb = C_EMERALD

    # Center Bottom: Benchmark Performance Metrics (width: 4.8)
    _, tf_s7_bm = create_card(slide7, Inches(4.3), Inches(4.90), Inches(4.8), Inches(1.95),
                              bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_bm_h = tf_s7_bm.paragraphs[0]
    p_bm_h.text = "BENCHMARK PERFORMANCE [SIMULATION]"
    p_bm_h.font.name = "Arial"
    p_bm_h.font.bold = True
    p_bm_h.font.size = Pt(8.5)
    p_bm_h.font.color.rgb = C_NAVY_HERO
    p_bm_h.space_after = Pt(2)

    bms = [
        ("End-to-End DAG Latency", "< 85 ms (FastAPI + SQLite WAL)"),
        ("Safety Circuit Breaker", "< 8 ms (Pure Python Deterministic)"),
        ("A* Pathfinder Compute", "< 35 ms (100x100 Spatial Cost Grid)")
    ]
    for bmn, bmv in bms:
        p_bm = tf_s7_bm.add_paragraph()
        p_bm.space_after = Pt(1.5)
        r1 = p_bm.add_run()
        r1.text = f"• {bmn}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p_bm.add_run()
        r2.text = bmv
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_SLATE_700

    # Right: Actionable Decision Output Panel (width: 3.43)
    _, tf_s7_dec = create_card(slide7, Inches(9.3), Inches(1.30), Inches(3.43), Inches(5.55),
                               bg_color=C_AMBER_BG, border_color=C_AMBER, border_width=1.5)
    p_s7_dh = tf_s7_dec.paragraphs[0]
    p_s7_dh.text = "3. OPERATIONAL DECISION"
    p_s7_dh.font.name = "Arial"
    p_s7_dh.font.bold = True
    p_s7_dh.font.size = Pt(9.5)
    p_s7_dh.font.color.rgb = C_AMBER
    p_s7_dh.space_after = Pt(4)

    p_dial = tf_s7_dec.add_paragraph()
    p_dial.text = "STATUS: PROCEED WITH CAUTION\nRisk Score: 0.24 / 1.00 (Safe with Diversion)\n"
    p_dial.font.name = "Arial"
    p_dial.font.bold = True
    p_dial.font.size = Pt(8.5)
    p_dial.font.color.rgb = C_NAVY_HERO
    p_dial.space_after = Pt(4)

    p_voice_s7 = tf_s7_dec.add_paragraph()
    p_voice_s7.text = "📢 TAMIL AUDIO ADVISORY:\n“கவனமாக செல்லவும். ஆழமற்ற பவளப்பாறை அலைகளை தவிர்க்க 3.2 கடல் மைல் கிழக்கு நோக்கி செல்லவும்.”\n"
    p_voice_s7.font.name = "Arial"
    p_voice_s7.font.size = Pt(7.5)
    p_voice_s7.font.color.rgb = C_SLATE_900
    p_voice_s7.space_after = Pt(4)

    p_loop = tf_s7_dec.add_paragraph()
    p_loop.text = "🔄 LIVE DECISION LOOP:\nChanging Wave Height, Vessel Size, or Heading triggers instant recalculation (<85ms)."
    p_loop.font.name = "Arial"
    p_loop.font.bold = True
    p_loop.font.size = Pt(7.5)
    p_loop.font.color.rgb = C_NAVY_HERO

    # =========================================================================
    # SLIDE 8: SAFETY & FAILURE RESILIENCE FIREWALL (WHAT IF THINGS FAIL?)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide8, 8, 10, "Resilience & Failsafes", 
                                 "“DESIGNED FOR WHEN THE SYSTEM IS WRONG”", 
                                 "Hostile maritime failure modes converge into the ORCA Resilience Core: When data becomes uncertain, the system becomes more conservative.")
    print("Formatting Slide 8: Safety & Failure Resilience...")

    # Left: 6 Hostile Failure Modes (width: 5.6)
    _, tf_s8_l = create_card(slide8, Inches(0.6), Inches(1.30), Inches(5.6), Inches(4.35),
                             bg_color=C_ROSE_BG, border_color=C_ROSE, border_width=1.0)
    p_s8_lh = tf_s8_l.paragraphs[0]
    p_s8_lh.text = "HOSTILE SEA FAILURE CONDITIONS"
    p_s8_lh.font.name = "Arial"
    p_s8_lh.font.bold = True
    p_s8_lh.font.size = Pt(9.5)
    p_s8_lh.font.color.rgb = C_ROSE
    p_s8_lh.space_after = Pt(3)

    fails_6 = [
        ("1. Stale Ocean Telemetry", "Forecast timestamp older than 3 hours or satellite pass delayed."),
        ("2. API Provider Outage", "INCOIS / Open-Meteo REST endpoints return HTTP 500 or timeout."),
        ("3. Model Disagreement", "IMD predicts 18 kt wind while ECMWF model indicates 32 kt gusts."),
        ("4. Offshore Signal Blackout", "Vessel sails past 12 NM territorial limit; GSM cellular signal dies."),
        ("5. Onboard Sensor Failure", "GPS or pitch-roll IMU disconnects or outputs erratic 5-sigma noise."),
        ("6. ML Out-of-Distribution", "Extreme unpredicted atmospheric event triggers high ML uncertainty.")
    ]
    for fn, fd in fails_6:
        p_f = tf_s8_l.add_paragraph()
        p_f.space_after = Pt(2.5)
        r1 = p_f.add_run()
        r1.text = f"⚠️ {fn}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_SLATE_900
        r2 = p_f.add_run()
        r2.text = fd
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Right: ORCA Engineering Mitigations (width: 6.33)
    _, tf_s8_r = create_card(slide8, Inches(6.4), Inches(1.30), Inches(6.33), Inches(4.35),
                             bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.5)
    p_s8_rh = tf_s8_r.paragraphs[0]
    p_s8_rh.text = "ORCA RESILIENCE FIREWALL"
    p_s8_rh.font.name = "Arial"
    p_s8_rh.font.bold = True
    p_s8_rh.font.size = Pt(9.5)
    p_s8_rh.font.color.rgb = C_EMERALD
    p_s8_rh.space_after = Pt(3)

    mits_6 = [
        ("TTL Freshness Check (<30m)", "Automatically flags stale data; applies conservative wave safety buffers."),
        ("Provider Failover Cascade", "Automated failover: INCOIS ERDDAP ➔ Open-Meteo Marine ➔ MOSDAC Backup."),
        ("Provenance & Confidence Bounds", "Calculates ensemble disagreement; defaults to highest-risk scenario."),
        ("Multi-Hop LoRa Fleet Mesh", "Vessels relay 16-byte emergency packets across 3-8 km hops to lighthouse."),
        ("Kalman Filter Sensor Sanity", "Rejects IMU spikes > 3 sigma; falls back to static vessel stability limits."),
        ("100% Deterministic Override", "Zero AI in safety gate; pure compiled physical logic enforces non-negotiable stop.")
    ]
    for mn, md in mits_6:
        p_m = tf_s8_r.add_paragraph()
        p_m.space_after = Pt(2.5)
        r1 = p_m.add_run()
        r1.text = f"🛡️ {mn}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p_m.add_run()
        r2.text = md
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Bottom Hero Statement: "WHEN DATA BECOMES UNCERTAIN..."
    _, tf_s8_hero = create_card(slide8, Inches(0.6), Inches(5.75), Inches(12.13), Inches(1.05),
                                bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.0)
    p_s8_hh = tf_s8_hero.paragraphs[0]
    p_s8_hh.text = "“WHEN DATA BECOMES UNCERTAIN, THE SYSTEM BECOMES MORE CONSERVATIVE.”"
    p_s8_hh.font.name = "Arial"
    p_s8_hh.font.bold = True
    p_s8_hh.font.size = Pt(11)
    p_s8_hh.font.color.rgb = C_WHITE
    p_s8_hh.alignment = PP_ALIGN.CENTER
    p_s8_hh.space_after = Pt(2)

    p_s8_hb = tf_s8_hero.add_paragraph()
    p_s8_hb.text = "BAD DATA ≠ FALSE CONFIDENCE. ORCA degrades safely to physical upper bounds rather than risking mariner lives on unverified AI inferences."
    p_s8_hb.font.name = "Arial"
    p_s8_hb.font.size = Pt(8)
    p_s8_hb.font.color.rgb = RGBColor(203, 213, 225)
    p_s8_hb.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 9: FEASIBILITY, DEPLOYMENT & ECONOMICS (CAN IT SCALE?)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide9, 9, 10, "Execution & Scale", 
                                 "“FROM PROTOTYPE TO COASTAL INFRASTRUCTURE”", 
                                 "Three-pillar scaling model: Verified software milestones, low-cost LoRa mesh deployment, and ₹3,500 per-boat BOM economics.")
    print("Formatting Slide 9: Feasibility, Deployment & Economics...")

    # Top: Scale Progression Journey (top: 1.30, height: 0.65)
    _, tf_s9_top = create_card(slide9, Inches(0.6), Inches(1.30), Inches(12.13), Inches(0.65),
                               bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.0)
    p_s9_th = tf_s9_top.paragraphs[0]
    p_s9_th.text = "NATIONAL SCALE PROGRESSION JOURNEY"
    p_s9_th.font.name = "Arial"
    p_s9_th.font.bold = True
    p_s9_th.font.size = Pt(8.5)
    p_s9_th.font.color.rgb = C_CYAN_ACCENT
    
    p_s9_tb = tf_s9_top.add_paragraph()
    p_s9_tb.text = "1 ARTISANAL VESSEL  ➔  HARBOR FLEET (500+)  ➔  COASTAL STATE  ➔  PAN-INDIA MARITIME DIGITAL TWIN LAYER (7,516 KM)"
    p_s9_tb.font.name = "Arial"
    p_s9_tb.font.bold = True
    p_s9_tb.font.size = Pt(8)
    p_s9_tb.font.color.rgb = C_WHITE

    # 3 Columns: Roadmap, Deployment Architecture, BOM Economics (top: 2.05, height: 4.75)
    # Col 1: Verified Roadmap (width: 3.6)
    _, tf_s9_c1 = create_card(slide9, Inches(0.6), Inches(2.05), Inches(3.6), Inches(4.75),
                              bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_c1_h = tf_s9_c1.paragraphs[0]
    p_c1_h.text = "ZONE 1: VERIFIED ROADMAP"
    p_c1_h.font.name = "Arial"
    p_c1_h.font.bold = True
    p_c1_h.font.size = Pt(9)
    p_c1_h.font.color.rgb = C_NAVY_HERO
    p_c1_h.space_after = Pt(3)

    milestones_s9 = [
        ("✓ END-TO-END PROTOTYPE", "FastAPI + React MapLibre working architecture [BUILT]"),
        ("✓ MULTI-SOURCE INGESTION", "Open-Meteo, IMD, MOSDAC pipelines active [BUILT]"),
        ("✓ DETERMINISTIC RISK ENGINE", "Physical capsizing & Hcrit circuit breaker [BUILT]"),
        ("✓ DYNAMIC A* PATHFINDER", "Time-varying 2D risk surface routing [BUILT]"),
        ("◐ HARBOR PILOT TRIAL", "Kanyakumari & Munambam field trials [IN PROGRESS]"),
        ("◌ PAN-INDIA EXPANSION", "Integration across 9 coastal states & Satcom [PLANNED]")
    ]
    for mt, md in milestones_s9:
        p_m = tf_s9_c1.add_paragraph()
        p_m.space_after = Pt(3)
        r1 = p_m.add_run()
        r1.text = f"{mt}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_EMERALD if "✓" in mt else (C_AMBER if "◐" in mt else C_MUTED)
        r2 = p_m.add_run()
        r2.text = md
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_SLATE_700

    # Col 2: Mesh Deployment Architecture (width: 4.2)
    _, tf_s9_c2 = create_card(slide9, Inches(4.4), Inches(2.05), Inches(4.2), Inches(4.75),
                              bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_c2_h = tf_s9_c2.paragraphs[0]
    p_c2_h.text = "ZONE 2: DEPLOYMENT ARCHITECTURE"
    p_c2_h.font.name = "Arial"
    p_c2_h.font.bold = True
    p_c2_h.font.size = Pt(9)
    p_c2_h.font.color.rgb = C_BLUE_PRIMARY
    p_c2_h.space_after = Pt(3)

    dep_layers = [
        ("1. Vessel Edge Node", "Low-power ESP32 MCU + LoRa SX1262 transceiver + GPS NEO-6M."),
        ("2. Fleet Ad-Hoc Relay", "Boats spaced 3–8 km apart form a dynamic peer-to-peer radio mesh."),
        ("3. Coastal Lighthouse Mast", "60m lighthouse receiver with line-of-sight reach of 33.8 km to boats."),
        ("4. ORCA Fast Cloud Engine", "Ingests 16-byte binary packets; executes A* & HSI in <85ms."),
        ("5. Dual User Interface", "Offline PWA for fishers; real-time dashboard for Coast Guard.")
    ]
    for dt, dd in dep_layers:
        p_d = tf_s9_c2.add_paragraph()
        p_d.space_after = Pt(3)
        r1 = p_d.add_run()
        r1.text = f"{dt}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p_d.add_run()
        r2.text = dd
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_SLATE_700

    # Col 3: BOM Economics & Triple Pillar (width: 3.93)
    _, tf_s9_c3 = create_card(slide9, Inches(8.8), Inches(2.05), Inches(3.93), Inches(4.75),
                              bg_color=C_EMERALD_BG, border_color=C_EMERALD, border_width=1.0)
    p_c3_h = tf_s9_c3.paragraphs[0]
    p_c3_h.text = "ZONE 3: ECONOMICS & BOM"
    p_c3_h.font.name = "Arial"
    p_c3_h.font.bold = True
    p_c3_h.font.size = Pt(9)
    p_c3_h.font.color.rgb = C_EMERALD
    p_c3_h.space_after = Pt(3)

    bom_items = [
        ("ESP32 MCU + LoRa SX1262", "₹1,800 (~$22)"),
        ("GPS NEO-6M Receiver", "₹600 (~$7)"),
        ("IP67 Enclosure + Solar Panel", "₹1,100 (~$13)"),
        ("Total Hardware Cost / Boat", "≈ ₹3,500 ($42) [INDICATIVE]"),
        ("Public Data Ingestion Cost", "₹0 (Open INCOIS & MOSDAC)"),
        ("Serverless Compute / Trip", "₹0.04 per assessment [ESTIMATE]")
    ]
    for bn, bv in bom_items:
        p_b = tf_s9_c3.add_paragraph()
        p_b.space_after = Pt(2)
        r1 = p_b.add_run()
        r1.text = f"• {bn}: "
        r1.font.name = "Arial"
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_SLATE_900
        r2 = p_b.add_run()
        r2.text = bv
        r2.font.name = "Arial"
        r2.font.bold = True
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_EMERALD if "Total" in bn or "₹0" in bv else C_BLUE_PRIMARY

    p_c3_pil = tf_s9_c3.add_paragraph()
    p_c3_pil.space_after = Pt(2)
    r_pil_h = p_c3_pil.add_run()
    r_pil_h.text = "\nTHREE-PILLAR SUMMARY:\n"
    r_pil_h.font.name = "Arial"
    r_pil_h.font.bold = True
    r_pil_h.font.size = Pt(8)
    r_pil_h.font.color.rgb = C_NAVY_HERO
    r_pil_b = p_c3_pil.add_run()
    r_pil_b.text = "FAST (<85ms Engine) | CHEAP (~₹3,500 BOM) | DEMOABLE (Live Loop)"
    r_pil_b.font.name = "Arial"
    r_pil_b.font.size = Pt(7)
    r_pil_b.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 10: IMPACT, RESEARCH & FINAL VISION (WHY DOES IT MATTER?)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_official_header_footer(slide10, 10, 10, "Impact & Vision", 
                                 "“ONE VESSEL. ONE HARBOR. ONE NATIONAL MARITIME INTELLIGENCE LAYER.”", 
                                 "Transforming academic oceanographic research and space telemetry into sovereign digital infrastructure for India's Blue Economy.")
    print("Formatting Slide 10: Impact, Research & Final Vision...")

    # Top: 4 Grounded Research-to-Code Blueprint Cards (top: 1.30, height: 1.75)
    research_s10 = [
        ("IMO Code on Intact Stability (2008)", "Res. MSC.267(85)", "Hcrit ≈ 0.6·Lhull capsizing threshold", C_BLUE_LIGHT, C_BLUE_PRIMARY),
        ("Wave Hydrodynamics (Faltinsen)", "Encounter Physics", "Wave steepness (S = Hs/λ) & roll resonance", C_EMERALD_BG, C_EMERALD),
        ("Satellite Oceanography (INCOIS)", "MOSDAC Telemetry", "Sobel thermal fronts & Multi-species HSI", C_AMBER_BG, C_AMBER),
        ("Dynamic Pathfinding (Modified A*)", "Risk Cost Surfaces", "4-Point waypoint storm avoidance routing", C_BLUE_LIGHT, C_BLUE_PRIMARY)
    ]
    for idx, (rt, rsrc, rapp, bg_c, b_col) in enumerate(research_s10):
        _, tf_r = create_card(slide10, Inches(0.6 + idx * 3.08), Inches(1.30), Inches(2.9), Inches(1.75),
                              bg_color=bg_c, border_color=b_col, border_width=1.0)
        p1 = tf_r.paragraphs[0]
        p1.text = f"📚 {rt}"
        p1.font.name = "Arial"
        p1.font.bold = True
        p1.font.size = Pt(8)
        p1.font.color.rgb = b_col
        
        p2 = tf_r.add_paragraph()
        p2.text = f"Source: {rsrc}"
        p2.font.name = "Arial"
        p2.font.italic = True
        p2.font.size = Pt(7)
        p2.font.color.rgb = C_MUTED
        p2.space_after = Pt(1.5)

        p3 = tf_r.add_paragraph()
        r1 = p3.add_run()
        r1.text = "ORCA Implementation: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_NAVY_HERO
        r2 = p3.add_run()
        r2.text = rapp
        r2.font.name = "Arial"
        r2.font.size = Pt(7)
        r2.font.color.rgb = C_SLATE_700

    # Center: 5-Stakeholder Multiplier Ring (top: 3.15, height: 1.65)
    _, tf_stk = create_card(slide10, Inches(0.6), Inches(3.15), Inches(12.13), Inches(1.65),
                            bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_stk_h = tf_stk.paragraphs[0]
    p_stk_h.text = "NATIONAL BLUE ECONOMY STAKEHOLDER MULTIPLIER"
    p_stk_h.font.name = "Arial"
    p_stk_h.font.bold = True
    p_stk_h.font.size = Pt(9.5)
    p_stk_h.font.color.rgb = C_NAVY_HERO
    p_stk_h.space_after = Pt(2)

    stk_items = [
        ("🎣 4.5M Artisanal Fishers", "Zero preventable drownings, timely voice warnings in native dialect, precise fishing spots."),
        ("⚓ Harbors & Port Authorities", "Automated safety clearance, harbor fleet congestion management & departure scheduling."),
        ("🚢 Commercial Trawlers", "20–30% fuel savings via current-assisted A* routing; higher pelagic catch ROI."),
        ("🛡️ Coast Guard & NDRF", "Instant Dark-Fleet anomaly matching & 1,000-particle Monte Carlo SAR drift tracking."),
        ("🔬 Oceanographic Scientists", "Crowdsourced ground-truth sea state telemetry across India's 2.37M sq km EEZ.")
    ]
    for s_title, s_desc in stk_items:
        p_s = tf_stk.add_paragraph()
        p_s.space_after = Pt(1.5)
        r1 = p_s.add_run()
        r1.text = f"{s_title}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p_s.add_run()
        r2.text = s_desc
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Bottom: Final Heroic Vision Statement (top: 4.90, height: 1.95)
    _, tf_vis = create_card(slide10, Inches(0.6), Inches(4.90), Inches(12.13), Inches(1.95),
                            bg_color=C_NAVY_HERO, border_color=C_CYAN_ACCENT, border_width=1.5)
    p_vis_h = tf_vis.paragraphs[0]
    p_vis_h.text = "“ORCA DOESN’T JUST PREDICT THE OCEAN. IT HELPS PEOPLE DECIDE WHAT TO DO NEXT.”"
    p_vis_h.font.name = "Arial"
    p_vis_h.font.size = Pt(12)
    p_vis_h.font.bold = True
    p_vis_h.font.color.rgb = C_WHITE
    p_vis_h.alignment = PP_ALIGN.CENTER
    p_vis_h.space_after = Pt(3)

    p_vis_sub = tf_vis.add_paragraph()
    p_vis_sub.text = "From a single artisanal fishing boat in Kanyakumari to a sovereign digital intelligence layer for India’s entire 7,516 km coastline."
    p_vis_sub.font.name = "Arial"
    p_vis_sub.font.size = Pt(9)
    p_vis_sub.font.color.rgb = RGBColor(203, 213, 225)
    p_vis_sub.alignment = PP_ALIGN.CENTER
    p_vis_sub.space_after = Pt(4)

    p_vis_flow = tf_vis.add_paragraph()
    p_vis_flow.text = "OCEAN  ➔  WEATHER  ➔  VESSEL  ➔  ROUTE  ➔  ORCA  ➔  RISK  ➔  DECISION"
    p_vis_flow.font.name = "Arial"
    p_vis_flow.font.size = Pt(8.5)
    p_vis_flow.font.bold = True
    p_vis_flow.font.color.rgb = C_CYAN_ACCENT
    p_vis_flow.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Extreme Master Deck generated successfully!")
    print(f"Saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_presentation()
