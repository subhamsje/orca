import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_PATH = "/Users/subham/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx"
OUTPUT_PATH = "/Users/subham/Downloads/ORCA_4.0_SIH2026_Official_Deck.pptx"

# ==============================================================================
# PREMIUM HUMAN-CRAFTED DESIGN SYSTEM (CLEAN, MODERN, SOPHISTICATED)
# ==============================================================================
# Base Theme
C_WHITE        = RGBColor(255, 255, 255)
C_SLATE_50     = RGBColor(248, 250, 252)   # #f8fafc - Lightest background
C_SLATE_100    = RGBColor(241, 245, 249)   # #f1f5f9 - Card fill
C_SLATE_200    = RGBColor(226, 232, 240)   # #e2e8f0 - Border subtle
C_SLATE_700    = RGBColor(51, 65, 85)      # #334155 - Body text
C_SLATE_900    = RGBColor(15, 23, 42)      # #0f172a - Dark bold text
C_MUTED        = RGBColor(100, 116, 139)   # #64748b - Subtitle text

# Brand & Accent Colors (Maritime Intelligence)
C_NAVY_DEEP    = RGBColor(10, 37, 64)      # #0a2540 - Stripe/Hero Navy
C_BLUE_PRIMARY = RGBColor(14, 116, 144)    # #0e7490 - Ocean Cyan/Teal
C_BLUE_LIGHT   = RGBColor(240, 249, 255)   # #f0f9ff - Ice Blue Tint
C_BLUE_BORDER  = RGBColor(186, 230, 253)   # #bae6fd - Blue border

C_EMERALD      = RGBColor(16, 185, 129)    # #10b981 - Green Safe
C_EMERALD_BG   = RGBColor(236, 253, 245)   # #ecfdf5 - Green Tint
C_EMERALD_BORDER = RGBColor(167, 243, 208) # #a7f3d0

C_AMBER        = RGBColor(217, 119, 6)     # #d97706 - Warning Amber
C_AMBER_BG     = RGBColor(254, 243, 199)   # #fef3c7 - Amber Tint
C_AMBER_BORDER = RGBColor(253, 230, 138)   # #fde68a

C_ROSE         = RGBColor(225, 29, 72)     # #e11d48 - Rose Danger
C_ROSE_BG      = RGBColor(255, 241, 242)   # #fff1f2 - Rose Tint

def style_box(shape, bg_color, border_color=None, border_width=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

def create_box(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    style_box(shape, bg_color, border_color, border_width)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return shape, tf

def setup_header(slide, title_text, pointer_badges):
    for s in slide.shapes:
        if s.name == "Title 1":
            s.left = Inches(1.85)
            s.top = Inches(0.12)
            s.width = Inches(8.75)
            s.height = Inches(0.80)
            tf = s.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = "Arial"
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = C_NAVY_DEEP
            p.alignment = PP_ALIGN.LEFT
        elif "Oval" in s.name:
            s.left = Inches(0.36)
            s.top = Inches(0.18)
            s.width = Inches(1.35)
            s.height = Inches(0.72)
            tf = s.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "ORCA 4.0"
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.alignment = PP_ALIGN.CENTER
            s.fill.solid()
            s.fill.fore_color.rgb = C_NAVY_DEEP
            s.line.fill.background()
        elif "TextBox" in s.name:
            s.text_frame.clear()

    # Add pointer badges bar
    badge_w = Inches(2.95)
    badge_h = Inches(0.28)
    for idx, (b_text, b_bg, b_fg) in enumerate(pointer_badges):
        bx = Inches(0.6 + idx * 3.06)
        by = Inches(0.98)
        b_shape, b_tf = create_box(slide, bx, by, badge_w, badge_h, bg_color=b_bg, border_color=C_SLATE_200, border_width=1.0)
        p = b_tf.paragraphs[0]
        p.text = b_text
        p.font.name = "Arial"
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = b_fg
        p.alignment = PP_ALIGN.CENTER

def build_deck():
    print(f"Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)

    if len(prs.slides) >= 7:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        print("Removed instruction slide 7. Exact 6 slides active.")

    # =========================================================================
    # SLIDE 1: TITLE PAGE (PREMIUM HERO)
    # =========================================================================
    slide1 = prs.slides[0]
    print("Formatting Slide 1: TITLE PAGE...")

    for shape in slide1.shapes:
        if shape.name == "TextBox 9":
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fields = [
                ("Problem Statement ID:", " SIH26176"),
                ("Problem Statement Title:", " AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform for Artisanal Fishermen & Coastal Fleets"),
                ("Theme:", " Space Technology / Blue Economy / Disaster Management"),
                ("PS Category:", " Software (with Edge Hardware compatibility)"),
                ("Team ID:", " [Team ID]"),
                ("Team Name:", " ORCA 4.0")
            ]
            for idx, (label, val) in enumerate(fields):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.space_after = Pt(7)
                r1 = p.add_run()
                r1.text = label
                r1.font.name = "Arial"
                r1.font.bold = True
                r1.font.size = Pt(12)
                r1.font.color.rgb = C_NAVY_DEEP
                r2 = p.add_run()
                r2.text = val
                r2.font.name = "Arial"
                r2.font.bold = False
                r2.font.size = Pt(12)
                r2.font.color.rgb = C_SLATE_900
        elif shape.name == "Subtitle 3":
            shape.text_frame.text = "TITLE PAGE"

    # Right Hero Showcase Card (Clean Luxury Dark Navy with Cyan Accents)
    _, tf_h1 = create_box(slide1, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5),
                          bg_color=C_NAVY_DEEP, border_color=C_BLUE_PRIMARY, border_width=1.5)
    
    p0 = tf_h1.paragraphs[0]
    p0.text = "SMART INDIA HACKATHON 2026 · NATIONAL FINALS"
    p0.font.name = "Arial"
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = C_BLUE_PRIMARY
    p0.space_after = Pt(4)

    p1 = tf_h1.add_paragraph()
    p1.text = "ORCA 4.0"
    p1.font.name = "Arial"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.space_after = Pt(2)

    p2 = tf_h1.add_paragraph()
    p2.text = "“FROM OCEAN DATA TO SAFER DECISIONS”"
    p2.font.name = "Arial"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(56, 189, 248)
    p2.space_after = Pt(6)

    p3 = tf_h1.add_paragraph()
    p3.text = "Universal Bio-Physical Decision Operating System for India's 4.5M Marine Workforce"
    p3.font.name = "Arial"
    p3.font.size = Pt(10)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(12)

    p4 = tf_h1.add_paragraph()
    p4.text = "CENTRAL OPERATIONAL CHAIN:"
    p4.font.name = "Arial"
    p4.font.size = Pt(9)
    p4.font.bold = True
    p4.font.color.rgb = C_EMERALD
    
    p5 = tf_h1.add_paragraph()
    p5.text = "OCEAN  ➔  WEATHER  ➔  VESSEL  ➔  ROUTE  ➔  ORCA  ➔  RISK  ➔  DECISION"
    p5.font.name = "Arial"
    p5.font.size = Pt(8.5)
    p5.font.bold = True
    p5.font.color.rgb = C_WHITE
    p5.space_after = Pt(14)

    s1_cards = [
        ("🛡️ Deterministic Safety Shield", "Non-bypassable capsizing physics (<10ms) with zero AI hallucination", C_EMERALD),
        ("🚢 Parametric Vessel Twin", "Dynamic wave encounter math tailored to hull length, draft & beam", RGBColor(56, 189, 248)),
        ("📻 Multi-Hop LoRa Fleet Mesh", "Overcomes 12 NM 4G blackout for zero-cost telemetry up to 50 km", C_AMBER)
    ]
    for c_title, c_desc, c_col in s1_cards:
        p_c = tf_h1.add_paragraph()
        p_c.space_after = Pt(5)
        r_t = p_c.add_run()
        r_t.text = f"{c_title}: "
        r_t.font.name = "Arial"
        r_t.font.bold = True
        r_t.font.size = Pt(9.5)
        r_t.font.color.rgb = c_col
        r_d = p_c.add_run()
        r_d.text = c_desc
        r_d.font.name = "Arial"
        r_d.font.size = Pt(8.5)
        r_d.font.color.rgb = RGBColor(203, 213, 225)

    # =========================================================================
    # SLIDE 2: IDEA TITLE & PROPOSED SOLUTION
    # =========================================================================
    slide2 = prs.slides[1]
    print("Formatting Slide 2: IDEA TITLE / PROPOSED SOLUTION...")
    setup_header(slide2, "IDEA TITLE: ORCA 4.0 — Maritime Decision Intelligence",
                 [("PROPOSED SOLUTION", C_BLUE_LIGHT, C_BLUE_PRIMARY),
                  ("DETAILED EXPLANATION", C_SLATE_100, C_SLATE_700),
                  ("HOW IT ADDRESSES PROBLEM", C_SLATE_100, C_SLATE_700),
                  ("INNOVATION & UNIQUENESS", C_EMERALD_BG, C_EMERALD)])

    # Top Hero Callout (Paradigm Shift)
    _, tf_th = create_box(slide2, Inches(0.6), Inches(1.34), Inches(12.13), Inches(0.50),
                          bg_color=C_BLUE_LIGHT, border_color=C_BLUE_BORDER, border_width=1.0)
    p_th = tf_th.paragraphs[0]
    r1 = p_th.add_run()
    r1.text = "CORE PARADIGM SHIFT:  FORECAST ≠ OPERATIONAL RISK.  "
    r1.font.name = "Arial"
    r1.font.bold = True
    r1.font.size = Pt(9.5)
    r1.font.color.rgb = C_BLUE_PRIMARY
    r2 = p_th.add_run()
    r2.text = "Existing apps ask 'What is happening at sea?'. ORCA computes 'What does it mean for THIS vessel, on THIS route, at THIS time?'"
    r2.font.name = "Arial"
    r2.font.size = Pt(9)
    r2.font.color.rgb = C_SLATE_900

    # Left Container: Data Transformation (width: 7.2)
    _, tf_t_hdr = create_box(slide2, Inches(0.6), Inches(1.90), Inches(7.2), Inches(4.30),
                             bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_thd = tf_t_hdr.paragraphs[0]
    p_thd.text = "VISUAL TRANSFORMATION: FROM DATA CHAOS TO ACTIONABLE DECISION"
    p_thd.font.name = "Arial"
    p_thd.font.bold = True
    p_thd.font.size = Pt(9.5)
    p_thd.font.color.rgb = C_NAVY_DEEP

    inputs_1 = [
        ("🌊 WAVES (Hs, Swell Tp)", "Sea & swell wave energy spectra"),
        ("💨 WIND (Vel, Gusts)", "Surface boundary wind vectors"),
        ("🌀 CYCLONES (IMD Alerts)", "Active storm center coordinates"),
        ("🌊 CURRENTS (Surface u,v)", "Ocean surface drift velocity")
    ]
    inputs_2 = [
        ("🛰️ SATELLITES (SST / Chlor)", "Thermal front boundaries"),
        ("🚢 VESSEL (Lwl, Draft, Beam)", "Displacement hull profile"),
        ("📡 AIS FEEDS (MMSI Data)", "Real-time vessel position"),
        ("🗺️ ROUTE (Waypoints)", "Departure & target docks")
    ]

    for idx, (title, sub) in enumerate(inputs_1):
        _, tf_in = create_box(slide2, Inches(0.72), Inches(2.26 + idx * 0.92), Inches(1.95), Inches(0.82),
                              bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
        p = tf_in.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = C_BLUE_PRIMARY
        p2 = tf_in.add_paragraph()
        p2.text = sub
        p2.font.name = "Arial"
        p2.font.size = Pt(7)
        p2.font.color.rgb = C_MUTED

    for idx, (title, sub) in enumerate(inputs_2):
        _, tf_in = create_box(slide2, Inches(2.76), Inches(2.26 + idx * 0.92), Inches(1.95), Inches(0.82),
                              bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
        p = tf_in.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = C_BLUE_PRIMARY
        p2 = tf_in.add_paragraph()
        p2.text = sub
        p2.font.name = "Arial"
        p2.font.size = Pt(7)
        p2.font.color.rgb = C_MUTED

    # Central Engine Box
    _, tf_ce = create_box(slide2, Inches(4.82), Inches(2.26), Inches(2.86), Inches(1.70),
                          bg_color=C_NAVY_DEEP, border_color=C_BLUE_PRIMARY, border_width=1.5)
    p_ce = tf_ce.paragraphs[0]
    p_ce.text = "⚡ ORCA 4.0 ENGINE"
    p_ce.font.name = "Arial"
    p_ce.font.bold = True
    p_ce.font.size = Pt(11)
    p_ce.font.color.rgb = C_WHITE
    p_ce.alignment = PP_ALIGN.CENTER
    p_ce2 = tf_ce.add_paragraph()
    p_ce2.text = "• 0.083° Spatial Normalization\n• Parametric Vessel Digital Twin\n• Multi-Species Habitat Index"
    p_ce2.font.name = "Arial"
    p_ce2.font.size = Pt(8)
    p_ce2.font.color.rgb = RGBColor(203, 213, 225)

    # Downstream Decision Boxes
    out_cards = [
        ("🛡️ DETERMINISTIC RISK", "Hcrit capsizing check (<10ms)", C_EMERALD_BG, C_EMERALD),
        ("🗺️ ROUTE RISK FIELD", "Dynamic 2D risk surface A*", C_BLUE_LIGHT, C_BLUE_PRIMARY),
        ("📢 ACTIONABLE VERDICT", "Clear Go/No-Go + Audio Voice", C_AMBER_BG, C_AMBER)
    ]
    for idx, (title, sub, bg_c, b_col) in enumerate(out_cards):
        _, tf_oc = create_box(slide2, Inches(4.82), Inches(4.06 + idx * 0.70), Inches(2.86), Inches(0.62),
                              bg_color=bg_c, border_color=b_col, border_width=1.0)
        p = tf_oc.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{title}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = b_col
        r2 = p.add_run()
        r2.text = sub
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Right Box: Competitive Differentiation Matrix
    _, tf_diff = create_box(slide2, Inches(7.95), Inches(1.90), Inches(4.78), Inches(4.30),
                            bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_dfh = tf_diff.paragraphs[0]
    p_dfh.text = "COMPETITIVE DIFFERENTIATION MATRIX"
    p_dfh.font.name = "Arial"
    p_dfh.font.bold = True
    p_dfh.font.size = Pt(10)
    p_dfh.font.color.rgb = C_NAVY_DEEP
    p_dfh.space_after = Pt(4)

    comp_rows = [
        ("Capability", "Existing Systems", "ORCA 4.0"),
        ("Weather Forecast", "● Regional Only", "✓ High-Res (0.083°)"),
        ("Cyclone Warnings", "● Generic SMS", "✓ Dynamic Geofenced"),
        ("Vessel Hull Context", "— Ignored", "✓ Vessel Twin (Lwl, Draft)"),
        ("Wave Physics Math", "— None", "✓ Hcrit & Resonance Math"),
        ("Safe Route Pathfinder", "— Static Path", "✓ Dynamic A* Avoidance"),
        ("Deterministic Safety", "— Unverified", "✓ Non-bypassable Rule"),
        ("Offshore Reach", "— Dead >12 NM", "✓ LoRa Mesh (50 km)")
    ]
    for idx, (c1, c2, c3) in enumerate(comp_rows):
        p_r = tf_diff.add_paragraph()
        p_r.space_after = Pt(2.5)
        r1 = p_r.add_run()
        r1.text = f"{c1.ljust(20)}: "
        r1.font.name = "Arial"
        r1.font.bold = True if idx == 0 else False
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_DEEP if idx == 0 else C_SLATE_900
        
        r2 = p_r.add_run()
        r2.text = f"{c2}  ➔  "
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_ROSE if "—" in c2 else C_MUTED
        
        r3 = p_r.add_run()
        r3.text = c3
        r3.font.name = "Arial"
        r3.font.bold = True
        r3.font.size = Pt(8)
        r3.font.color.rgb = C_EMERALD if "✓" in c3 else C_SLATE_900

    # Bottom Innovation Strip
    _, tf_btm2 = create_box(slide2, Inches(0.6), Inches(6.28), Inches(12.13), Inches(0.42),
                            bg_color=C_SLATE_100, border_color=C_SLATE_200, border_width=1.0)
    p_b2 = tf_btm2.paragraphs[0]
    p_b2.text = "① MULTI-SOURCE HARMONIZATION   ➔   ② VESSEL-AWARE PHYSICS TWIN   ➔   ③ ROUTE RISK OPTIMIZATION   ➔   ④ EXPLAINABLE SAFETY SHIELD"
    p_b2.font.name = "Arial"
    p_b2.font.bold = True
    p_b2.font.size = Pt(9)
    p_b2.font.color.rgb = C_NAVY_DEEP
    p_b2.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides[2]
    print("Formatting Slide 3: TECHNICAL APPROACH...")
    setup_header(slide3, "TECHNICAL APPROACH: End-to-End Architecture",
                 [("TECHNOLOGIES TO BE USED", C_BLUE_LIGHT, C_BLUE_PRIMARY),
                  ("METHODOLOGY & PROCESS", C_SLATE_100, C_SLATE_700),
                  ("SYSTEM FLOWCHARTS", C_SLATE_100, C_SLATE_700),
                  ("PROTOTYPE SIMULATION", C_EMERALD_BG, C_EMERALD)])

    # Top Architecture Box
    _, tf_a_bg = create_box(slide3, Inches(0.6), Inches(1.34), Inches(12.13), Inches(3.30),
                            bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)

    # 6 Ingestion Nodes
    sources = [
        ("INCOIS ERDDAP", "Wave spectra & currents"),
        ("IMD Alerts", "Real-time cyclone tracks"),
        ("Copernicus Marine", "0.083° global ocean grids"),
        ("ISRO MOSDAC", "INSAT-3DR SST & Oceansat-3"),
        ("AIS Streams", "Vessel transponder data"),
        ("NMEA Sensors", "Onboard GPS / IMU pitch")
    ]
    for idx, (s_name, s_sub) in enumerate(sources):
        _, tf_s = create_box(slide3, Inches(0.72 + idx * 1.96), Inches(1.44), Inches(1.88), Inches(0.64),
                             bg_color=C_WHITE, border_color=C_BLUE_BORDER, border_width=1.0)
        p = tf_s.paragraphs[0]
        p.text = f"📡 {s_name}"
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8)
        p.font.color.rgb = C_BLUE_PRIMARY
        p2 = tf_s.add_paragraph()
        p2.text = s_sub
        p2.font.name = "Arial"
        p2.font.size = Pt(7)
        p2.font.color.rgb = C_MUTED

    # Normalization Funnel Bar
    _, tf_fn = create_box(slide3, Inches(0.72), Inches(2.16), Inches(11.89), Inches(0.42),
                          bg_color=C_BLUE_PRIMARY, border_color=None)
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = "⚡ DATA ENGINE & DUAL DIGITAL STATE: Spatial H3 Indexing (0.083°) + Environmental State  ⟷  Parametric Vessel Digital Twin (Lwl, Draft, Beam)"
    p_fn.font.name = "Arial"
    p_fn.font.bold = True
    p_fn.font.size = Pt(8.5)
    p_fn.font.color.rgb = C_WHITE
    p_fn.alignment = PP_ALIGN.CENTER

    # Tripartite Engines
    engines = [
        ("⚙️ DETERMINISTIC PHYSICS (<10ms)", 
         "• Critical Wave Height (Hcrit = 0.6 · Lhull)\n• Wave Steepness (S = Hs / λ) & Period\n• Beam-Sea Encounter Angle & Roll Resonance",
         C_EMERALD_BG, C_EMERALD),
        ("🧠 PREDICTIVE ML & ANALYTICS", 
         "• Multi-Species Habitat Index (HSI XGBoost)\n• 2D Sobel Filter SST Thermal Fronts\n• 1,000-Particle Monte Carlo SAR Drift",
         C_BLUE_LIGHT, C_BLUE_PRIMARY),
        ("📍 DYNAMIC A* OPTIMIZATION", 
         "• 2D Dynamic Risk Cost Surface Mesh\n• A* Navigational Pathfinding & Avoidance\n• Optimal Fuel Burn & Dock Allocation",
         C_AMBER_BG, C_AMBER)
    ]
    for idx, (e_title, e_desc, bg_c, b_col) in enumerate(engines):
        _, tf_e = create_box(slide3, Inches(0.72 + idx * 4.02), Inches(2.66), Inches(3.85), Inches(1.24),
                             bg_color=bg_c, border_color=b_col, border_width=1.0)
        p = tf_e.paragraphs[0]
        p.text = e_title
        p.font.name = "Arial"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = b_col
        p2 = tf_e.add_paragraph()
        p2.text = e_desc
        p2.font.name = "Arial"
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = C_SLATE_900

    # Safety Shield Banner
    _, tf_scb = create_box(slide3, Inches(0.72), Inches(3.98), Inches(11.89), Inches(0.54),
                           bg_color=C_NAVY_DEEP, border_color=C_EMERALD, border_width=1.0)
    p_scb = tf_scb.paragraphs[0]
    r1 = p_scb.add_run()
    r1.text = "🛡️ NON-BYPASSABLE SAFETY CIRCUIT BREAKER: "
    r1.font.name = "Arial"
    r1.font.bold = True
    r1.font.size = Pt(9)
    r1.font.color.rgb = C_EMERALD
    r2 = p_scb.add_run()
    r2.text = "Pure deterministic compiled logic evaluates physical capsizing risk (<10ms). LLM is strictly constrained downstream for vernacular audio synthesis only."
    r2.font.name = "Arial"
    r2.font.size = Pt(8)
    r2.font.color.rgb = C_WHITE

    # Bottom Left: Live Prototype Simulation
    _, tf_dm = create_box(slide3, Inches(0.6), Inches(4.72), Inches(5.9), Inches(1.98),
                          bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_dm_h = tf_dm.paragraphs[0]
    p_dm_h.text = "LIVE WORKING PROTOTYPE SIMULATION (KANYAKUMARI CASE)"
    p_dm_h.font.name = "Arial"
    p_dm_h.font.bold = True
    p_dm_h.font.size = Pt(9.5)
    p_dm_h.font.color.rgb = C_NAVY_DEEP
    p_dm_h.space_after = Pt(3)

    k_steps = [
        ("1. Input", "Kanyakumari Port (8.08°N, 77.55°E), 9m FRP Craft"),
        ("2. Live State", "Wave Hs = 2.8m, Swell Tp = 8.2s, Wind = 24 kts"),
        ("3. Physics Check", "Hcrit = 5.4m > 2.8m (Safe from Immediate Capsize)"),
        ("4. Route Engine", "Diverts 3.2 NM East to avoid shallow reef breaker"),
        ("5. Output Verdict", "PROCEED WITH CAUTION (Audio Advisory Generated)")
    ]
    for sn, sv in k_steps:
        p = tf_dm.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"{sn}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_NAVY_DEEP
        r2 = p.add_run()
        r2.text = sv
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Bottom Right: Production Tech Stack
    _, tf_tc = create_box(slide3, Inches(6.65), Inches(4.72), Inches(6.08), Inches(1.98),
                          bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_tc_h = tf_tc.paragraphs[0]
    p_tc_h.text = "PRODUCTION TECHNOLOGY STACK"
    p_tc_h.font.name = "Arial"
    p_tc_h.font.bold = True
    p_tc_h.font.size = Pt(9.5)
    p_tc_h.font.color.rgb = C_NAVY_DEEP
    p_tc_h.space_after = Pt(3)

    t_stack = [
        ("Frontend / UI", "React 18, TypeScript, Vite, MapLibre GL, WebGL, TailwindCSS, PWA"),
        ("Backend / API", "FastAPI, Python 3.10+, Asynchronous DAG Orchestrator, SQLite/WAL"),
        ("Intelligence", "NumPy, SciPy, XGBoost (HSI Matrix), Modified A* Graph Pathfinder"),
        ("IoT / Telecom", "ESP32 LoRa Gateway (868/433 MHz), Multi-Hop Mesh (OLSR/BATMAN)")
    ]
    for cat, spec in t_stack:
        p = tf_tc.add_paragraph()
        p.space_after = Pt(2.5)
        r1 = p.add_run()
        r1.text = f"• {cat}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p.add_run()
        r2.text = spec
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides[3]
    print("Formatting Slide 4: FEASIBILITY AND VIABILITY...")
    setup_header(slide4, "FEASIBILITY AND VIABILITY: Build, Mitigate & Deploy",
                 [("FEASIBILITY ANALYSIS", C_BLUE_LIGHT, C_BLUE_PRIMARY),
                  ("POTENTIAL RISKS & CHALLENGES", C_AMBER_BG, C_AMBER),
                  ("STRATEGIES FOR OVERCOMING", C_EMERALD_BG, C_EMERALD),
                  ("DEPLOYMENT & ECONOMICS", C_SLATE_100, C_SLATE_700)])

    # Column 1: Implementation Roadmap (width: 3.0)
    _, tf_rd = create_box(slide4, Inches(0.6), Inches(1.34), Inches(3.0), Inches(5.36),
                          bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_rh = tf_rd.paragraphs[0]
    p_rh.text = "IMPLEMENTATION ROADMAP"
    p_rh.font.name = "Arial"
    p_rh.font.bold = True
    p_rh.font.size = Pt(10)
    p_rh.font.color.rgb = C_NAVY_DEEP
    p_rh.space_after = Pt(6)

    milestones = [
        ("✓ PROTOTYPE", "End-to-End FastAPI & React UI working", C_EMERALD),
        ("✓ LIVE DATA", "Open-Meteo & IMD ingestion active", C_EMERALD),
        ("✓ RISK ENGINE", "Deterministic safety breaker active", C_EMERALD),
        ("✓ ROUTE ENGINE", "A* dynamic pathfinder & HSI matrix", C_EMERALD),
        ("◐ HARBOR PILOT", "Kanyakumari & Munambam trial", C_AMBER),
        ("◌ SCALE PHASE", "9 Coastal States & ISRO Satcom", C_MUTED)
    ]
    for title, desc, col in milestones:
        p = tf_rd.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = col
        r2 = p.add_run()
        r2.text = desc
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Column 2: Challenge ➔ Mitigation (width: 5.15)
    _, tf_ch = create_box(slide4, Inches(3.75), Inches(1.34), Inches(5.15), Inches(5.36),
                          bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_chh = tf_ch.paragraphs[0]
    p_chh.text = "CRITICAL RISKS ➔ PROVEN MITIGATION STRATEGIES"
    p_chh.font.name = "Arial"
    p_chh.font.bold = True
    p_chh.font.size = Pt(10)
    p_chh.font.color.rgb = C_NAVY_DEEP
    p_chh.space_after = Pt(6)

    challenges = [
        ("Data Latency / Stale Feeds", "Automated TTL monitoring (<30m) & cached offline fallback"),
        ("API Provider Outages", "Multi-source failover cascade (INCOIS ➔ Open-Meteo ➔ MOSDAC)"),
        ("Source Disagreements", "Provenance-weighted ensemble & explicit uncertainty bounds"),
        ("Offshore Signal Blackout", "Multi-Hop LoRa Fleet Mesh (up to 50 km) + Offline PWA Caching"),
        ("Vessel Hull Diversity", "Parametric Digital Twin model library (Catamaran, FRP, Trawler)"),
        ("Model Hallucination", "100% Deterministic Safety Rules (Zero AI in safety path)")
    ]
    for ct, md in challenges:
        p = tf_ch.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"⚠️ {ct}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_ROSE
        r2 = p.add_run()
        r2.text = f"  ➔ 🛡️ {md}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Column 3 Top: Deployment Architecture
    _, tf_dp = create_box(slide4, Inches(9.05), Inches(1.34), Inches(3.68), Inches(2.62),
                          bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_dph = tf_dp.paragraphs[0]
    p_dph.text = "DEPLOYMENT ARCHITECTURE"
    p_dph.font.name = "Arial"
    p_dph.font.bold = True
    p_dph.font.size = Pt(9.5)
    p_dph.font.color.rgb = C_NAVY_DEEP
    p_dph.space_after = Pt(4)

    dep_steps = [
        ("1. Vessel Edge", "GPS / IMU + ESP32 LoRa Node"),
        ("2. Fleet Relay", "Multi-Hop Boat Mesh (3-8 km hops)"),
        ("3. Coastal Mast", "Lighthouse LoRa Gateway (60m mast)"),
        ("4. Cloud Engine", "FastAPI / SQLite / Orchestrator"),
        ("5. User Tiers", "Fisher PWA | Fleet Web | Coast Guard")
    ]
    for ds, dv in dep_steps:
        p = tf_dp.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"{ds}: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p.add_run()
        r2.text = dv
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Column 3 Bottom: Deployment Economics & BOM
    _, tf_cs = create_box(slide4, Inches(9.05), Inches(4.08), Inches(3.68), Inches(2.62),
                          bg_color=C_EMERALD_BG, border_color=C_EMERALD_BORDER, border_width=1.0)
    p_csh = tf_cs.paragraphs[0]
    p_csh.text = "DEPLOYMENT ECONOMICS & BOM"
    p_csh.font.name = "Arial"
    p_csh.font.bold = True
    p_csh.font.size = Pt(9.5)
    p_csh.font.color.rgb = C_EMERALD
    p_csh.space_after = Pt(4)

    cost_items = [
        ("ESP32 MCU + LoRa SX1262", "₹1,800 (~$22)"),
        ("GPS NEO-6M Receiver", "₹600 (~$7)"),
        ("IP67 Enclosure + Solar", "₹1,100 (~$13)"),
        ("Total Hardware Cost / Boat", "≈ ₹3,500 ($42)"),
        ("Data & Server API Cost", "₹0 (Open Feeds / Serverless)")
    ]
    for cn, cv in cost_items:
        p = tf_cs.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"• {cn}: "
        r1.font.name = "Arial"
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_SLATE_900
        r2 = p.add_run()
        r2.text = cv
        r2.font.name = "Arial"
        r2.font.bold = True
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_EMERALD if "Total" in cn or "₹0" in cv else C_BLUE_PRIMARY

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides[4]
    print("Formatting Slide 5: IMPACT AND BENEFITS...")
    setup_header(slide5, "IMPACT AND BENEFITS: Empowering India's Blue Economy",
                 [("TARGET AUDIENCE IMPACT", C_BLUE_LIGHT, C_BLUE_PRIMARY),
                  ("SOCIAL & HUMAN BENEFITS", C_EMERALD_BG, C_EMERALD),
                  ("ECONOMIC & OPERATIONAL ROI", C_AMBER_BG, C_AMBER),
                  ("NATIONAL STRATEGIC SCALE", C_SLATE_100, C_SLATE_700)])

    # Left Container: ORCA Impact Ecosystem
    _, tf_ec = create_box(slide5, Inches(0.6), Inches(1.34), Inches(5.5), Inches(4.82),
                          bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_ech = tf_ec.paragraphs[0]
    p_ech.text = "ORCA MULTI-STAKEHOLDER ECOSYSTEM"
    p_ech.font.name = "Arial"
    p_ech.font.bold = True
    p_ech.font.size = Pt(10.5)
    p_ech.font.color.rgb = C_NAVY_DEEP
    p_ech.space_after = Pt(6)

    stakeholders = [
        ("🎣 Artisanal Fishermen (4.5M Lives)", "Zero preventable capsizings, localized vernacular audio safety alerts, precision PFZ fishing hotspots."),
        ("⚓ Harbors & Port Authorities", "Live harbor fleet tracking, automated departure safety clearances, harbor congestion mitigation."),
        ("🚢 Commercial Fishing Fleets", "Optimal fuel-saving route planning, dynamic weather avoidance, higher catch yield ROI."),
        ("🛡️ Coast Guard & NDRF", "Real-time Dark-Fleet anomaly alerts, 1,000-particle Monte Carlo SAR drift tracking."),
        ("🔬 Oceanographic Institutes", "Crowdsourced ground-truth sea state telemetry & high-resolution model validation.")
    ]
    for st, sd in stakeholders:
        p = tf_ec.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"{st}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = C_BLUE_PRIMARY
        r2 = p.add_run()
        r2.text = f"  • {sd}"
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_700

    # Right Container: 4 Pillars of National Impact
    _, tf_pl = create_box(slide5, Inches(6.25), Inches(1.34), Inches(6.48), Inches(4.82),
                          bg_color=C_WHITE, border_color=C_SLATE_200, border_width=1.0)
    p_plh = tf_pl.paragraphs[0]
    p_plh.text = "THE FOUR PILLARS OF NATIONAL IMPACT"
    p_plh.font.name = "Arial"
    p_plh.font.bold = True
    p_plh.font.size = Pt(10.5)
    p_plh.font.color.rgb = C_NAVY_DEEP
    p_plh.space_after = Pt(6)

    pillars = [
        ("👥 SOCIAL IMPACT (Saving Lives)", 
         "Eliminates fatal capsizings by calculating real-time wave-vessel physical resonance. Delivers plain-language voice alerts in Tamil, Malayalam, Bengali, Telugu, and Hindi for non-literate crews.",
         C_EMERALD),
        ("💰 ECONOMIC IMPACT (Fuel & Yield)", 
         "Reduces diesel expenditure by 20–30% through current-assisted dynamic A* pathfinding. Maximizes high-value pelagic catch yields via multi-species Habitat Suitability Index (HSI) mapping.",
         C_AMBER),
        ("🌿 ENVIRONMENTAL IMPACT (Ocean Health)", 
         "Substantially lowers maritime diesel emissions per fishing voyage. Provides automated geofencing to prevent accidental intrusion into Marine Protected Areas (MPAs) and international boundaries.",
         C_BLUE_PRIMARY),
        ("🇮🇳 STRATEGIC IMPACT (Digital Sovereignty)", 
         "Delivers end-to-end sovereign maritime domain awareness across India's 7,516 km coastline and 2.37M sq km EEZ, creating an integrated national ocean safety infrastructure.",
         C_NAVY_DEEP)
    ]
    for p_name, p_body, p_col in pillars:
        p = tf_pl.add_paragraph()
        p.space_after = Pt(5)
        r1 = p.add_run()
        r1.text = f"{p_name}\n"
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = p_col
        r2 = p.add_run()
        r2.text = p_body
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900

    # Bottom Progression Banner
    _, tf_btm5 = create_box(slide5, Inches(0.6), Inches(6.28), Inches(12.13), Inches(0.42),
                            bg_color=C_SLATE_100, border_color=C_SLATE_200, border_width=1.0)
    p_b5 = tf_btm5.paragraphs[0]
    p_b5.text = "1 ARTISANAL VESSEL   ➔   HARBOR FLEET (500+)   ➔   COASTAL STATE   ➔   PAN-INDIA MARITIME DIGITAL TWIN LAYER"
    p_b5.font.name = "Arial"
    p_b5.font.bold = True
    p_b5.font.size = Pt(9)
    p_b5.font.color.rgb = C_NAVY_DEEP
    p_b5.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides[5]
    print("Formatting Slide 6: RESEARCH AND REFERENCES...")
    setup_header(slide6, "RESEARCH  AND REFERENCES: Science to Implementation",
                 [("DETAILS OF RESEARCH WORK", C_BLUE_LIGHT, C_BLUE_PRIMARY),
                  ("IMO STABILITY CRITERIA", C_SLATE_100, C_SLATE_700),
                  ("HYDRODYNAMIC WAVE PHYSICS", C_SLATE_100, C_SLATE_700),
                  ("VERIFIED ACADEMIC CITATIONS", C_EMERALD_BG, C_EMERALD)])

    # 4 Core Research Cards (2x2 Grid)
    research_cards = [
        ("IMO Code on Intact Stability (2008)",
         "Resolution MSC.267(85) / Small Craft Capsizing Dynamics",
         "Small vessel dynamic instability under steep wave action, beam-sea broaching, and parametric roll resonance.",
         "Deterministic Safety Engine: Computes critical wave threshold (Hcrit ≈ 0.6 · Lhull) and enforces non-bypassable hazard overrides (<10ms).",
         C_BLUE_BORDER),
        ("Marine Hydrodynamics & Encounter Physics",
         "O.M. Faltinsen (1990) & J.N. Newman (1977)",
         "Wave steepness ratio (S = Hs / λ), deep-water dispersion relations (λ = g·Tp² / 2π), and encounter angle frequency shifts.",
         "Vessel Digital Twin: Dynamic calculation of wave steepness & vessel-wave interaction rather than treating sea states as static point forecasts.",
         C_EMERALD_BORDER),
        ("Satellite Oceanography & Pelagic Habitats",
         "INCOIS PFZ Mission & ISRO MOSDAC Telemetry",
         "Aggregation of pelagic marine species along thermal gradients, SST frontal boundaries, and chlorophyll-a density concentration zones.",
         "Multi-Species Habitat Index (HSI): XGBoost model integrated with 2D Sobel spatial filter for high-precision fishery zone recommendations.",
         C_AMBER_BORDER),
        ("Dynamic Constrained Pathfinding & Graph Optimization",
         "Modified A* over Time-Varying Navigational Cost Surfaces",
         "Optimization of multiobjective vessel routes under dynamic weather constraints, geofenced boundaries, and risk cost penalization.",
         "ORCA Navigational Pathfinder: Computes 4-point optimized waypoints balancing weather safety, fuel consumption, and dock proximity.",
         C_BLUE_BORDER)
    ]

    for idx, (title, src, insight, app, bcol) in enumerate(research_cards):
        col = idx % 2
        row = idx // 2
        x = Inches(0.6 + col * 6.18)
        y = Inches(1.34 + row * 2.40)
        
        _, tf_rc = create_box(slide6, x, y, Inches(5.95), Inches(2.30),
                              bg_color=C_WHITE, border_color=bcol, border_width=1.0)
        p1 = tf_rc.paragraphs[0]
        p1.text = f"📚 {title}"
        p1.font.name = "Arial"
        p1.font.bold = True
        p1.font.size = Pt(9)
        p1.font.color.rgb = C_NAVY_DEEP
        
        p_src = tf_rc.add_paragraph()
        p_src.text = f"Source: {src}"
        p_src.font.name = "Arial"
        p_src.font.italic = True
        p_src.font.size = Pt(7)
        p_src.font.color.rgb = C_MUTED
        p_src.space_after = Pt(2)
        
        p_ins = tf_rc.add_paragraph()
        r1 = p_ins.add_run()
        r1.text = "Key Insight: "
        r1.font.name = "Arial"
        r1.font.bold = True
        r1.font.size = Pt(7.5)
        r1.font.color.rgb = C_AMBER
        r2 = p_ins.add_run()
        r2.text = insight
        r2.font.name = "Arial"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = C_SLATE_900
        p_ins.space_after = Pt(2)

        p_app = tf_rc.add_paragraph()
        r3 = p_app.add_run()
        r3.text = "ORCA Implementation: "
        r3.font.name = "Arial"
        r3.font.bold = True
        r3.font.size = Pt(7.5)
        r3.font.color.rgb = C_EMERALD
        r4 = p_app.add_run()
        r4.text = app
        r4.font.name = "Arial"
        r4.font.size = Pt(7.5)
        r4.font.color.rgb = C_SLATE_700

    # Bottom Citations Container
    _, tf_rf = create_box(slide6, Inches(0.6), Inches(6.22), Inches(12.13), Inches(0.48),
                          bg_color=C_SLATE_50, border_color=C_SLATE_200, border_width=1.0)
    p_rf = tf_rf.paragraphs[0]
    p_rf.text = "AUTHENTIC OPERATIONAL & SCIENTIFIC REFERENCES"
    p_rf.font.name = "Arial"
    p_rf.font.bold = True
    p_rf.font.size = Pt(8)
    p_rf.font.color.rgb = C_NAVY_DEEP
    
    p_cits = tf_rf.add_paragraph()
    p_cits.text = (
        "[1] IMO Intact Stability Code (2008). Res. MSC.267(85). | "
        "[2] INCOIS Ocean State Forecast & PFZ Protocols. | "
        "[3] ISRO MOSDAC Oceansat-3 & INSAT-3DR Products. | "
        "[4] Copernicus Marine CMEMS Analysis. | "
        "[5] WMO-No. 558 Marine Services."
    )
    p_cits.font.name = "Arial"
    p_cits.font.size = Pt(7)
    p_cits.font.color.rgb = C_SLATE_700

    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Clean Human-Crafted Presentation generated successfully!")
    print(f"Saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_deck()
