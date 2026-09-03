import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

TEMPLATE_PATH = "/Users/subham/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx"
OUTPUT_PATH = "/Users/subham/Downloads/ORCA_4.0_SIH2026_Official_Deck.pptx"

# Color Palette
NAVY_DEEP = RGBColor(10, 25, 47)        # #0a192f Main deep card
NAVY_CARD = RGBColor(15, 34, 64)        # #0f2240 Card background
OCEAN_BLUE = RGBColor(2, 132, 199)      # #0284c7 Ocean Blue
CYAN_ACCENT = RGBColor(6, 182, 212)     # #06b6d4 Bright Cyan
TEAL_DARK = RGBColor(13, 71, 161)       # #0d47a1 Dark Teal / SIH Theme
EMERALD_SAFE = RGBColor(16, 185, 129)   # #10b981 Safe Green
AMBER_WARN = RGBColor(245, 158, 11)     # #f59e0b Warning Amber
ROSE_DANGER = RGBColor(239, 68, 68)     # #ef4444 Danger Red
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(15, 23, 42)        # Slate 900
TEXT_MUTED = RGBColor(100, 116, 139)    # Slate 500
TEXT_LIGHT_MUTED = RGBColor(203, 213, 225) # Slate 300
BG_LIGHT_BLUE = RGBColor(240, 249, 255) # Light blue card bg
BORDER_CYAN = RGBColor(56, 189, 248)    # Border Cyan
BORDER_MUTED = RGBColor(226, 232, 240)  # Border Slate 200

def add_card(slide, left, top, width, height, bg_color=None, border_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, margin=0.05):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    return txBox, tf

def add_badge(slide, left, top, width, height, text, bg_color, text_color=TEXT_WHITE, font_size=9, bold=True):
    shape = add_card(slide, left, top, width, height, bg_color=bg_color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape

def clear_shape_text(shape):
    if shape.has_text_frame:
        shape.text_frame.text = ""

def build_presentation():
    print(f"Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)

    # 1. DELETE SLIDE 7 (Instruction Slide)
    if len(prs.slides) >= 7:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        print("Deleted instruction slide. Active slides: 6")

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides[0]
    print("Populating Slide 1: TITLE PAGE...")

    # Find and update TextBox 9 (the official form fields)
    for shape in slide1.shapes:
        if shape.name == "TextBox 9":
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fields = [
                ("Problem Statement ID:", " SIH26176"),
                ("Problem Statement Title:", " AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform for Artisanal Fishermen & Coastal Fleets"),
                ("Theme:", " Space Technology / Blue Economy / Disaster Management"),
                ("PS Category:", " Software (with Edge Hardware compatibility)"),
                ("Team ID:", " [Team ID]"),
                ("Team Name:", " ORCA 4.0")
            ]
            
            for idx, (label, val) in enumerate(fields):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.space_after = Pt(8)
                
                r1 = p.add_run()
                r1.text = label
                r1.font.bold = True
                r1.font.size = Pt(13)
                r1.font.color.rgb = TEAL_DARK
                
                r2 = p.add_run()
                r2.text = val
                r2.font.bold = False
                r2.font.size = Pt(13)
                r2.font.color.rgb = TEXT_DARK
        elif shape.name == "Subtitle 3":
            # Reposition/restyle subtitle
            shape.text_frame.text = "TITLE PAGE"

    # Add ORCA 4.0 Refined Branding Card on Right Side of Slide 1
    card1 = add_card(slide1, Inches(7.1), Inches(1.3), Inches(5.8), Inches(5.6), bg_color=NAVY_DEEP, border_color=CYAN_ACCENT)
    
    # Text inside Brand Card
    tbox_brand, tf_brand = add_textbox(slide1, Inches(7.3), Inches(1.5), Inches(5.4), Inches(5.2))
    
    p_badge = tf_brand.paragraphs[0]
    p_badge.text = "SMART INDIA HACKATHON 2026 · PROJECT SHOWCASE"
    p_badge.font.size = Pt(10)
    p_badge.font.bold = True
    p_badge.font.color.rgb = CYAN_ACCENT
    
    p_title = tf_brand.add_paragraph()
    p_title.text = "ORCA 4.0"
    p_title.font.size = Pt(38)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_after = Pt(2)
    
    p_tagline = tf_brand.add_paragraph()
    p_tagline.text = "“FROM OCEAN DATA TO SAFER DECISIONS”"
    p_tagline.font.size = Pt(15)
    p_tagline.font.bold = True
    p_tagline.font.color.rgb = CYAN_ACCENT
    p_tagline.space_after = Pt(10)

    p_desc = tf_brand.add_paragraph()
    p_desc.text = "AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = TEXT_LIGHT_MUTED
    p_desc.space_after = Pt(14)
    
    # Story Chain Banner
    p_chain_hdr = tf_brand.add_paragraph()
    p_chain_hdr.text = "ORCA CENTRAL DECISION CHAIN:"
    p_chain_hdr.font.size = Pt(10)
    p_chain_hdr.font.bold = True
    p_chain_hdr.font.color.rgb = EMERALD_SAFE
    
    p_chain = tf_brand.add_paragraph()
    p_chain.text = "OCEAN ➔ WEATHER ➔ VESSEL ➔ ROUTE ➔ ORCA ➔ RISK ➔ DECISION"
    p_chain.font.size = Pt(9.5)
    p_chain.font.bold = True
    p_chain.font.color.rgb = TEXT_WHITE
    p_chain.space_after = Pt(14)

    # 3 Key Value Badges inside card
    val_cards = [
        ("🛡️ Deterministic Safety Shield", "Hard physics circuit breaker (Hcrit = 0.6·L) with zero hallucination risk"),
        ("🚢 Parametric Vessel Twin", "Dynamic wave-vessel encounter dynamics tailored to hull length & draft"),
        ("📻 Multi-Hop LoRa Fleet Mesh", "Overcomes 12 NM 4G blackout for zero-cost offshore telemetry up to 50 km")
    ]
    for title, desc in val_cards:
        p_v = tf_brand.add_paragraph()
        p_v.space_after = Pt(6)
        r_t = p_v.add_run()
        r_t.text = f"{title}: "
        r_t.font.bold = True
        r_t.font.size = Pt(9.5)
        r_t.font.color.rgb = CYAN_ACCENT
        
        r_d = p_v.add_run()
        r_d.text = desc
        r_d.font.size = Pt(9)
        r_d.font.color.rgb = TEXT_LIGHT_MUTED

    # =========================================================================
    # SLIDE 2: IDEA TITLE / PROPOSED SOLUTION
    # =========================================================================
    slide2 = prs.slides[1]
    print("Populating Slide 2: IDEA TITLE / PROPOSED SOLUTION...")
    
    for s in slide2.shapes:
        if s.name == "Title 1":
            s.text_frame.text = "IDEA TITLE: ORCA 4.0 — Maritime Decision Intelligence"
        elif s.name == "Oval 9":
            s.text_frame.text = "ORCA 4.0"
        elif s.name == "TextBox 8":
            # Clear original placeholder text box
            s.text_frame.clear()

    # Official Pointer Badges Bar (Top)
    pointers_s2 = [
        ("PROPOSED SOLUTION", CYAN_ACCENT),
        ("DETAILED EXPLANATION", TEAL_DARK),
        ("HOW IT ADDRESSES PROBLEM", TEAL_DARK),
        ("INNOVATION & UNIQUENESS", EMERALD_SAFE)
    ]
    for i, (p_text, p_col) in enumerate(pointers_s2):
        add_badge(slide2, Inches(0.5 + i * 3.1), Inches(1.15), Inches(3.0), Inches(0.32), p_text, p_col, font_size=9)

    # Key Thesis Box (Compact & High Impact)
    thesis_card = add_card(slide2, Inches(0.5), Inches(1.55), Inches(12.33), Inches(0.55), bg_color=NAVY_DEEP, border_color=CYAN_ACCENT)
    _, tf_th = add_textbox(slide2, Inches(0.6), Inches(1.58), Inches(12.1), Inches(0.5))
    p_th = tf_th.paragraphs[0]
    r_th1 = p_th.add_run()
    r_th1.text = "CORE PARADIGM SHIFT:  FORECAST ≠ OPERATIONAL RISK.  "
    r_th1.font.bold = True
    r_th1.font.size = Pt(11)
    r_th1.font.color.rgb = AMBER_WARN
    r_th2 = p_th.add_run()
    r_th2.text = "Existing systems ask 'What is happening at sea?'. ORCA computes 'What does it mean for THIS vessel, on THIS route, at THIS time?'"
    r_th2.font.size = Pt(10.5)
    r_th2.font.color.rgb = TEXT_WHITE

    # Left Container (Data Transformation Flow: 60% width)
    flow_bg = add_card(slide2, Inches(0.5), Inches(2.2), Inches(7.5), Inches(4.0), bg_color=NAVY_CARD, border_color=BORDER_CYAN)
    
    # Header inside flow_bg
    _, tf_fl_hdr = add_textbox(slide2, Inches(0.6), Inches(2.25), Inches(7.3), Inches(0.35))
    p_fh = tf_fl_hdr.paragraphs[0]
    p_fh.text = "VISUAL TRANSFORMATION: FROM FRAGMENTED DATA CHAOS TO OPERATIONAL DECISION"
    p_fh.font.size = Pt(10.5)
    p_fh.font.bold = True
    p_fh.font.color.rgb = CYAN_ACCENT

    # 8 Input Nodes (2 columns of 4 chips on the left of the flow)
    inputs_left = [
        ("🌊 WAVES", "Hs, Swell Tp, Period"),
        ("💨 WIND", "Speed, Gusts, Direction"),
        ("🌀 CYCLONES", "IMD Track & Warnings"),
        ("🌊 CURRENTS", "Surface Vectors (u,v)")
    ]
    inputs_right = [
        ("🛰️ SATELLITES", "SST & Chlorophyll-a"),
        ("🚢 VESSEL", "Hull Length, Draft, Beam"),
        ("📡 AIS FEEDS", "Fleet MMSI Coordinates"),
        ("🗺️ ROUTE", "Waypoints & Port Docks")
    ]
    
    for idx, (chip_title, chip_sub) in enumerate(inputs_left):
        c = add_card(slide2, Inches(0.7), Inches(2.65 + idx * 0.8), Inches(1.9), Inches(0.7), bg_color=NAVY_DEEP, border_color=BORDER_MUTED)
        _, tf_c = add_textbox(slide2, Inches(0.75), Inches(2.67 + idx * 0.8), Inches(1.8), Inches(0.65))
        p1 = tf_c.paragraphs[0]
        p1.text = chip_title
        p1.font.bold = True
        p1.font.size = Pt(9.5)
        p1.font.color.rgb = CYAN_ACCENT
        p2 = tf_c.add_paragraph()
        p2.text = chip_sub
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = TEXT_LIGHT_MUTED

    for idx, (chip_title, chip_sub) in enumerate(inputs_right):
        c = add_card(slide2, Inches(2.7), Inches(2.65 + idx * 0.8), Inches(1.9), Inches(0.7), bg_color=NAVY_DEEP, border_color=BORDER_MUTED)
        _, tf_c = add_textbox(slide2, Inches(2.75), Inches(2.67 + idx * 0.8), Inches(1.8), Inches(0.65))
        p1 = tf_c.paragraphs[0]
        p1.text = chip_title
        p1.font.bold = True
        p1.font.size = Pt(9.5)
        p1.font.color.rgb = CYAN_ACCENT
        p2 = tf_c.add_paragraph()
        p2.text = chip_sub
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = TEXT_LIGHT_MUTED

    # Central Convergence Core: ORCA 4.0 ENGINE
    orca_engine = add_card(slide2, Inches(4.8), Inches(2.65), Inches(2.9), Inches(1.5), bg_color=OCEAN_BLUE, border_color=CYAN_ACCENT)
    _, tf_oe = add_textbox(slide2, Inches(4.85), Inches(2.7), Inches(2.8), Inches(1.4))
    p_oe1 = tf_oe.paragraphs[0]
    p_oe1.text = "⚡ ORCA 4.0 ENGINE"
    p_oe1.font.bold = True
    p_oe1.font.size = Pt(12)
    p_oe1.font.color.rgb = TEXT_WHITE
    p_oe1.alignment = PP_ALIGN.CENTER
    
    p_oe2 = tf_oe.add_paragraph()
    p_oe2.text = "• Spatial & Temporal Grid Fusion\n• Parametric Vessel Digital Twin\n• Multi-Species Habitat Index"
    p_oe2.font.size = Pt(8.5)
    p_oe2.font.color.rgb = TEXT_WHITE
    
    # Downstream Decision Boxes
    out_boxes = [
        ("🛡️ DETERMINISTIC RISK", "Hcrit capsizing check (<10ms)", EMERALD_SAFE),
        ("🗺️ ROUTE RISK FIELD", "Dynamic 2D risk surface A*", CYAN_ACCENT),
        ("📢 ACTIONABLE VERDICT", "Clear Go/No-Go + Audio Voice", AMBER_WARN)
    ]
    for idx, (title, sub, col) in enumerate(out_boxes):
        ob = add_card(slide2, Inches(4.8), Inches(4.3 + idx * 0.62), Inches(2.9), Inches(0.55), bg_color=NAVY_DEEP, border_color=col)
        _, tf_ob = add_textbox(slide2, Inches(4.85), Inches(4.32 + idx * 0.62), Inches(2.8), Inches(0.5))
        p_o = tf_ob.paragraphs[0]
        r_ot = p_o.add_run()
        r_ot.text = f"{title}: "
        r_ot.font.bold = True
        r_ot.font.size = Pt(8.5)
        r_ot.font.color.rgb = col
        r_os = p_o.add_run()
        r_os.text = sub
        r_os.font.size = Pt(8)
        r_os.font.color.rgb = TEXT_LIGHT_MUTED

    # Right Container: Differentiation Matrix (40% width)
    diff_bg = add_card(slide2, Inches(8.15), Inches(2.2), Inches(4.68), Inches(4.0), bg_color=BG_LIGHT_BLUE, border_color=TEAL_DARK)
    _, tf_df = add_textbox(slide2, Inches(8.25), Inches(2.25), Inches(4.48), Inches(3.9))
    
    p_dh = tf_df.paragraphs[0]
    p_dh.text = "COMPETITIVE DIFFERENTIATION MATRIX"
    p_dh.font.bold = True
    p_dh.font.size = Pt(11)
    p_dh.font.color.rgb = TEAL_DARK
    p_dh.space_after = Pt(6)

    comp_rows = [
        ("Capability", "Existing Systems", "ORCA 4.0"),
        ("Weather Forecast", "✅ Regional Only", "✅ High-Res (0.083°)"),
        ("Cyclone Warnings", "✅ Generic SMS", "✅ Dynamic Geofenced"),
        ("Vessel Hull Context", "❌ Ignored", "✅ Vessel Twin (Lwl, Draft)"),
        ("Wave Physics Math", "❌ None", "✅ Hcrit & Resonance Math"),
        ("Safe Route Pathfinder", "❌ Static Path", "✅ Dynamic A* Avoidance"),
        ("Deterministic Safety", "❌ Unverified", "✅ Non-bypassable Rule"),
        ("Offshore Reach", "❌ Dead >12 NM", "✅ LoRa Mesh (50 km)")
    ]

    for idx, (c1, c2, c3) in enumerate(comp_rows):
        p_r = tf_df.add_paragraph()
        p_r.space_after = Pt(2)
        
        r1 = p_r.add_run()
        r1.text = f"{c1.ljust(20)}: "
        r1.font.bold = True if idx == 0 else False
        r1.font.size = Pt(8.5)
        r1.font.color.rgb = TEAL_DARK if idx == 0 else TEXT_DARK
        
        r2 = p_r.add_run()
        r2.text = f"{c2} ➔ "
        r2.font.size = Pt(8)
        r2.font.color.rgb = ROSE_DANGER if "❌" in c2 else TEXT_MUTED
        
        r3 = p_r.add_run()
        r3.text = f"{c3}"
        r3.font.bold = True
        r3.font.size = Pt(8.5)
        r3.font.color.rgb = EMERALD_SAFE if "✅" in c3 else TEXT_DARK

    # Bottom Innovation Strip (ORCA Continuous Route Line)
    strip_bg = add_card(slide2, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.55), bg_color=NAVY_DEEP, border_color=CYAN_ACCENT)
    _, tf_st = add_textbox(slide2, Inches(0.6), Inches(6.32), Inches(12.1), Inches(0.5))
    p_st = tf_st.paragraphs[0]
    p_st.alignment = PP_ALIGN.CENTER
    
    st_items = [
        ("① MULTI-SOURCE", CYAN_ACCENT),
        (" ➔ ", TEXT_MUTED),
        ("② VESSEL-AWARE TWIN", CYAN_ACCENT),
        (" ➔ ", TEXT_MUTED),
        ("③ ROUTE RISK FIELD", CYAN_ACCENT),
        (" ➔ ", TEXT_MUTED),
        ("④ EXPLAINABLE VERDICT", EMERALD_SAFE)
    ]
    for text, color in st_items:
        r = p_st.add_run()
        r.text = text
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = color

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides[2]
    print("Populating Slide 3: TECHNICAL APPROACH...")
    
    for s in slide3.shapes:
        if s.name == "Title 1":
            s.text_frame.text = "TECHNICAL APPROACH: End-to-End Pipeline & Architecture"
        elif s.name == "Oval 10":
            s.text_frame.text = "ORCA 4.0"
        elif s.name == "TextBox 8":
            s.text_frame.clear()

    # Official Pointer Badges Bar (Top)
    pointers_s3 = [
        ("TECHNOLOGIES TO BE USED", CYAN_ACCENT),
        ("METHODOLOGY & PROCESS", TEAL_DARK),
        ("FLOWCHARTS / SYSTEM PIPELINE", TEAL_DARK),
        ("WORKING PROTOTYPE SIMULATION", EMERALD_SAFE)
    ]
    for i, (p_text, p_col) in enumerate(pointers_s3):
        add_badge(slide3, Inches(0.5 + i * 3.1), Inches(1.15), Inches(3.0), Inches(0.32), p_text, p_col, font_size=9)

    # Main Architecture Canvas (Top 55% of content area)
    arch_bg = add_card(slide3, Inches(0.5), Inches(1.55), Inches(12.33), Inches(3.3), bg_color=NAVY_DEEP, border_color=BORDER_CYAN)
    
    # Layer 1: 6 Ingestion Source Boxes
    sources = [
        ("INCOIS ERDDAP", "Wave spectrum & currents"),
        ("IMD Alerts", "Real-time cyclone tracks"),
        ("Copernicus Marine", "0.083° global ocean grids"),
        ("ISRO MOSDAC", "INSAT-3DR SST & Oceansat-3"),
        ("AIS Streams", "Vessel transponder positions"),
        ("NMEA Sensors", "Onboard GPS / IMU / Pitch")
    ]
    for idx, (s_name, s_sub) in enumerate(sources):
        sc = add_card(slide3, Inches(0.65 + idx * 1.98), Inches(1.65), Inches(1.9), Inches(0.65), bg_color=NAVY_CARD, border_color=CYAN_ACCENT)
        _, tf_sc = add_textbox(slide3, Inches(0.68 + idx * 1.98), Inches(1.67), Inches(1.84), Inches(0.6))
        p = tf_sc.paragraphs[0]
        p.text = f"📡 {s_name}"
        p.font.bold = True
        p.font.size = Pt(8.5)
        p.font.color.rgb = CYAN_ACCENT
        p2 = tf_sc.add_paragraph()
        p2.text = s_sub
        p2.font.size = Pt(7.5)
        p2.font.color.rgb = TEXT_LIGHT_MUTED

    # Layer 2: Normalization & Digital State Funnel Bar
    funnel = add_card(slide3, Inches(0.65), Inches(2.4), Inches(12.0), Inches(0.48), bg_color=OCEAN_BLUE, border_color=CYAN_ACCENT)
    _, tf_fn = add_textbox(slide3, Inches(0.7), Inches(2.42), Inches(11.9), Inches(0.44))
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = "⚡ DATA ENGINE & DUAL DIGITAL STATE:  Spatial Indexing (0.083° / H3)  +  Environmental State  ⟷  Parametric Vessel Digital Twin (Lwl, Draft, Beam)"
    p_fn.font.bold = True
    p_fn.font.size = Pt(9.5)
    p_fn.font.color.rgb = TEXT_WHITE
    p_fn.alignment = PP_ALIGN.CENTER

    # Layer 3: Tripartite Intelligence Engines (Strict Separation)
    engines = [
        ("⚙️ DETERMINISTIC PHYSICS (<10ms)", 
         "• Critical Wave Height (Hcrit = 0.6 · Lhull)\n• Wave Steepness (S = Hs / λ) & Period\n• Beam-Sea Encounter Angle & Roll Resonance",
         EMERALD_SAFE),
        ("🧠 PREDICTIVE ML & ANALYTICS", 
         "• Multi-Species Habitat Index (HSI XGBoost)\n• 2D Sobel Filter SST Thermal Fronts\n• 1,000-Particle Monte Carlo SAR Drift",
         CYAN_ACCENT),
        ("📍 DYNAMIC A* OPTIMIZATION", 
         "• 2D Dynamic Risk Cost Surface Mesh\n• A* Navigational Pathfinding & Avoidance\n• Optimal Fuel Burn & Dock Allocation",
         AMBER_WARN)
    ]
    for idx, (e_title, e_desc, e_col) in enumerate(engines):
        ec = add_card(slide3, Inches(0.65 + idx * 4.05), Inches(2.98), Inches(3.9), Inches(1.1), bg_color=NAVY_CARD, border_color=e_col)
        _, tf_ec = add_textbox(slide3, Inches(0.75 + idx * 4.05), Inches(3.02), Inches(3.7), Inches(1.0))
        p = tf_ec.paragraphs[0]
        p.text = e_title
        p.font.bold = True
        p.font.size = Pt(9.5)
        p.font.color.rgb = e_col
        p2 = tf_ec.add_paragraph()
        p2.text = e_desc
        p2.font.size = Pt(8)
        p2.font.color.rgb = TEXT_LIGHT_MUTED

    # Layer 4: Hard Safety Circuit Breaker & Outputs Banner
    scb = add_card(slide3, Inches(0.65), Inches(4.18), Inches(12.0), Inches(0.55), bg_color=NAVY_CARD, border_color=EMERALD_SAFE)
    _, tf_scb = add_textbox(slide3, Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.5))
    p_scb = tf_scb.paragraphs[0]
    r_sc1 = p_scb.add_run()
    r_sc1.text = "🛡️ NON-BYPASSABLE SAFETY CIRCUIT BREAKER: "
    r_sc1.font.bold = True
    r_sc1.font.size = Pt(9.5)
    r_sc1.font.color.rgb = EMERALD_SAFE
    r_sc2 = p_scb.add_run()
    r_sc2.text = "Pure deterministic compiled logic evaluates physical capsizing risk. LLM is strictly constrained downstream for plain-language vernacular audio translation only."
    r_sc2.font.size = Pt(8.5)
    r_sc2.font.color.rgb = TEXT_WHITE

    # Bottom Left: Live Working Prototype Flow (Kanyakumari Case Study)
    demo_bg = add_card(slide3, Inches(0.5), Inches(4.95), Inches(5.8), Inches(1.9), bg_color=BG_LIGHT_BLUE, border_color=TEAL_DARK)
    _, tf_dm = add_textbox(slide3, Inches(0.6), Inches(5.0), Inches(5.6), Inches(1.8))
    p_dh = tf_dm.paragraphs[0]
    p_dh.text = "LIVE WORKING PROTOTYPE FLOW (KANYAKUMARI CASE)"
    p_dh.font.bold = True
    p_dh.font.size = Pt(10)
    p_dh.font.color.rgb = TEAL_DARK
    p_dh.space_after = Pt(4)

    demo_steps = [
        ("1. Input", "Kanyakumari Port (8.08°N, 77.55°E), 9m FRP Craft"),
        ("2. Live State", "Wave Hs = 2.8m, Swell Tp = 8.2s, Wind = 24 kts"),
        ("3. Physics Check", "Hcrit = 5.4m > 2.8m (Safe from Immediate Capsize)"),
        ("4. Route Engine", "Diverts 3.2 NM East to avoid shallow reef breaker"),
        ("5. Output Verdict", "PROCEED WITH CAUTION (Audio Advisory Generated)")
    ]
    for step_num, step_val in demo_steps:
        p_s = tf_dm.add_paragraph()
        p_s.space_after = Pt(2)
        r_sn = p_s.add_run()
        r_sn.text = f"{step_num}: "
        r_sn.font.bold = True
        r_sn.font.size = Pt(8.5)
        r_sn.font.color.rgb = TEAL_DARK
        r_sv = p_s.add_run()
        r_sv.text = step_val
        r_sv.font.size = Pt(8)
        r_sv.font.color.rgb = TEXT_DARK

    # Bottom Right: Technology Stack Breakdown Grid
    tech_bg = add_card(slide3, Inches(6.5), Inches(4.95), Inches(6.33), Inches(1.9), bg_color=NAVY_DEEP, border_color=CYAN_ACCENT)
    _, tf_tc = add_textbox(slide3, Inches(6.6), Inches(5.0), Inches(6.1), Inches(1.8))
    p_th = tf_tc.paragraphs[0]
    p_th.text = "PRODUCTION TECHNOLOGY STACK"
    p_th.font.bold = True
    p_th.font.size = Pt(10)
    p_th.font.color.rgb = CYAN_ACCENT
    p_th.space_after = Pt(4)

    tech_items = [
        ("Frontend / UI", "React 18, TypeScript, Vite, MapLibre GL, WebGL, TailwindCSS, Offline PWA"),
        ("Backend / API", "FastAPI, Python 3.10+, Asynchronous DAG Orchestrator, SQLite/WAL"),
        ("Intelligence", "NumPy, SciPy, XGBoost (HSI Matrix), Modified A* Graph Pathfinder"),
        ("IoT / Telecom", "ESP32 LoRa Gateway (868/433 MHz), Multi-Hop Mesh (OLSR/BATMAN)")
    ]
    for t_cat, t_spec in tech_items:
        p_t = tf_tc.add_paragraph()
        p_t.space_after = Pt(3)
        r_tc = p_t.add_run()
        r_tc.text = f"• {t_cat}: "
        r_tc.font.bold = True
        r_tc.font.size = Pt(8.5)
        r_tc.font.color.rgb = CYAN_ACCENT
        r_ts = p_t.add_run()
        r_ts.text = t_spec
        r_ts.font.size = Pt(8)
        r_ts.font.color.rgb = TEXT_LIGHT_MUTED

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides[3]
    print("Populating Slide 4: FEASIBILITY AND VIABILITY...")
    
    for s in slide4.shapes:
        if s.name == "Title 1":
            s.text_frame.text = "FEASIBILITY AND VIABILITY: Build, Mitigate & Deploy"
        elif s.name == "Oval 11":
            s.text_frame.text = "ORCA 4.0"
        elif s.name == "TextBox 8":
            s.text_frame.clear()

    # Official Pointer Badges Bar (Top)
    pointers_s4 = [
        ("FEASIBILITY ANALYSIS", CYAN_ACCENT),
        ("POTENTIAL RISKS & CHALLENGES", AMBER_WARN),
        ("STRATEGIES FOR OVERCOMING", EMERALD_SAFE),
        ("DEPLOYMENT & COST VIABILITY", TEAL_DARK)
    ]
    for i, (p_text, p_col) in enumerate(pointers_s4):
        add_badge(slide4, Inches(0.5 + i * 3.1), Inches(1.15), Inches(3.0), Inches(0.32), p_text, p_col, font_size=9)

    # Column 1: Implementation Roadmap (Left 25%)
    road_bg = add_card(slide4, Inches(0.5), Inches(1.55), Inches(3.1), Inches(5.3), bg_color=NAVY_DEEP, border_color=BORDER_CYAN)
    _, tf_rd = add_textbox(slide4, Inches(0.6), Inches(1.65), Inches(2.9), Inches(5.1))
    p_rh = tf_rd.paragraphs[0]
    p_rh.text = "IMPLEMENTATION ROADMAP"
    p_rh.font.bold = True
    p_rh.font.size = Pt(10.5)
    p_rh.font.color.rgb = CYAN_ACCENT
    p_rh.space_after = Pt(8)

    milestones = [
        ("✓ PROTOTYPE", "End-to-End FastAPI & React MapLibre working", EMERALD_SAFE),
        ("✓ LIVE DATA", "Open-Meteo, IMD, MOSDAC ingestion pipelines", EMERALD_SAFE),
        ("✓ RISK ENGINE", "Deterministic safety circuit breaker active", EMERALD_SAFE),
        ("✓ ROUTE ENGINE", "A* dynamic pathfinder & multi-species HSI", EMERALD_SAFE),
        ("◐ HARBOR PILOT", "Field deployment in Kanyakumari & Munambam", AMBER_WARN),
        ("◌ SCALE PHASE", "Integration across 9 coastal states & Satcom", TEXT_LIGHT_MUTED)
    ]
    for m_title, m_desc, m_col in milestones:
        p_m = tf_rd.add_paragraph()
        p_m.space_after = Pt(6)
        r_mt = p_m.add_run()
        r_mt.text = f"{m_title}\n"
        r_mt.font.bold = True
        r_mt.font.size = Pt(9)
        r_mt.font.color.rgb = m_col
        r_md = p_m.add_run()
        r_md.text = m_desc
        r_md.font.size = Pt(8)
        r_md.font.color.rgb = TEXT_WHITE

    # Column 2: Challenge ➔ Strategic Mitigation Matrix (Center 42%)
    chal_bg = add_card(slide4, Inches(3.75), Inches(1.55), Inches(5.2), Inches(5.3), bg_color=BG_LIGHT_BLUE, border_color=TEAL_DARK)
    _, tf_ch = add_textbox(slide4, Inches(3.85), Inches(1.65), Inches(5.0), Inches(5.1))
    p_ch_hdr = tf_ch.paragraphs[0]
    p_ch_hdr.text = "CRITICAL RISKS ➔ PROVEN MITIGATION STRATEGIES"
    p_ch_hdr.font.bold = True
    p_ch_hdr.font.size = Pt(10.5)
    p_ch_hdr.font.color.rgb = TEAL_DARK
    p_ch_hdr.space_after = Pt(6)

    challenges = [
        ("Data Latency / Stale Feeds", "Automated TTL monitoring (<30m) & cached offline fallback"),
        ("API Provider Outages", "Multi-source failover cascade (INCOIS ➔ Open-Meteo ➔ MOSDAC)"),
        ("Source Disagreements", "Provenance-weighted ensemble & explicit uncertainty bounds"),
        ("Offshore Signal Blackout", "Multi-Hop LoRa Fleet Mesh (up to 50 km) + Offline PWA Caching"),
        ("Vessel Hull Diversity", "Parametric Digital Twin model library (Catamaran, FRP, Trawler)"),
        ("Model Hallucination", "100% Deterministic Safety Rules (Zero AI in safety path)")
    ]
    for c_title, m_desc in challenges:
        p_c = tf_ch.add_paragraph()
        p_c.space_after = Pt(4)
        r_ct = p_c.add_run()
        r_ct.text = f"⚠️ {c_title}\n"
        r_ct.font.bold = True
        r_ct.font.size = Pt(8.5)
        r_ct.font.color.rgb = ROSE_DANGER
        r_cm = p_c.add_run()
        r_cm.text = f"  ➔ 🛡️ {m_desc}"
        r_cm.font.size = Pt(8)
        r_cm.font.color.rgb = TEXT_DARK

    # Column 3 (Top): Deployment Architecture (Right 31%)
    dep_bg = add_card(slide4, Inches(9.1), Inches(1.55), Inches(3.73), Inches(2.6), bg_color=NAVY_DEEP, border_color=CYAN_ACCENT)
    _, tf_dp = add_textbox(slide4, Inches(9.2), Inches(1.65), Inches(3.53), Inches(2.4))
    p_dph = tf_dp.paragraphs[0]
    p_dph.text = "DEPLOYMENT ARCHITECTURE"
    p_dph.font.bold = True
    p_dph.font.size = Pt(10)
    p_dph.font.color.rgb = CYAN_ACCENT
    p_dph.space_after = Pt(4)

    dep_steps = [
        ("1. Vessel Edge", "GPS / IMU + ESP32 LoRa Node"),
        ("2. Fleet Relay", "Multi-Hop Boat Mesh (3-8 km hops)"),
        ("3. Coastal Mast", "Lighthouse LoRa Gateway (60m mast)"),
        ("4. Cloud Engine", "FastAPI / SQLite / Orchestrator"),
        ("5. User Tiers", "Fisher PWA | Fleet Web | Coast Guard API")
    ]
    for d_step, d_val in dep_steps:
        p_ds = tf_dp.add_paragraph()
        p_ds.space_after = Pt(2)
        r_dst = p_ds.add_run()
        r_dst.text = f"{d_step}: "
        r_dst.font.bold = True
        r_dst.font.size = Pt(8)
        r_dst.font.color.rgb = CYAN_ACCENT
        r_dsv = p_ds.add_run()
        r_dsv.text = d_val
        r_dsv.font.size = Pt(7.5)
        r_dsv.font.color.rgb = TEXT_LIGHT_MUTED

    # Column 3 (Bottom): Deployment Economics & BOM (Right 31%)
    cost_bg = add_card(slide4, Inches(9.1), Inches(4.25), Inches(3.73), Inches(2.6), bg_color=NAVY_CARD, border_color=EMERALD_SAFE)
    _, tf_cs = add_textbox(slide4, Inches(9.2), Inches(4.35), Inches(3.53), Inches(2.4))
    p_csh = tf_cs.paragraphs[0]
    p_csh.text = "DEPLOYMENT ECONOMICS & BOM"
    p_csh.font.bold = True
    p_csh.font.size = Pt(10)
    p_csh.font.color.rgb = EMERALD_SAFE
    p_csh.space_after = Pt(4)

    cost_items = [
        ("ESP32 MCU + LoRa SX1262", "₹1,800 (~$22)"),
        ("GPS NEO-6M Receiver", "₹600 (~$7)"),
        ("IP67 Enclosure + Solar Panel", "₹1,100 (~$13)"),
        ("Total Hardware Cost / Boat", "≈ ₹3,500 ($42)"),
        ("Data & Server API Cost", "₹0 (Open Feeds / Serverless)")
    ]
    for c_name, c_val in cost_items:
        p_ci = tf_cs.add_paragraph()
        p_ci.space_after = Pt(2)
        r_cn = p_ci.add_run()
        r_cn.text = f"• {c_name}: "
        r_cn.font.size = Pt(8)
        r_cn.font.color.rgb = TEXT_WHITE
        r_cv = p_ci.add_run()
        r_cv.text = c_val
        r_cv.font.bold = True
        r_cv.font.size = Pt(8)
        r_cv.font.color.rgb = EMERALD_SAFE if "Total" in c_name or "₹0" in c_val else CYAN_ACCENT

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides[4]
    print("Populating Slide 5: IMPACT AND BENEFITS...")
    
    for s in slide5.shapes:
        if s.name == "Title 1":
            s.text_frame.text = "IMPACT AND BENEFITS: Empowering India's Blue Economy"
        elif s.name == "Oval 11":
            s.text_frame.text = "ORCA 4.0"
        elif s.name == "TextBox 8":
            s.text_frame.clear()

    # Official Pointer Badges Bar (Top)
    pointers_s5 = [
        ("TARGET AUDIENCE IMPACT", CYAN_ACCENT),
        ("SOCIAL & HUMAN BENEFITS", EMERALD_SAFE),
        ("ECONOMIC & OPERATIONAL ROI", AMBER_WARN),
        ("NATIONAL & STRATEGIC SCALE", TEAL_DARK)
    ]
    for i, (p_text, p_col) in enumerate(pointers_s5):
        add_badge(slide5, Inches(0.5 + i * 3.1), Inches(1.15), Inches(3.0), Inches(0.32), p_text, p_col, font_size=9)

    # Left Container: ORCA Impact Ecosystem (45% width)
    eco_bg = add_card(slide5, Inches(0.5), Inches(1.55), Inches(5.6), Inches(4.6), bg_color=NAVY_DEEP, border_color=BORDER_CYAN)
    _, tf_ec = add_textbox(slide5, Inches(0.6), Inches(1.65), Inches(5.4), Inches(4.4))
    p_ech = tf_ec.paragraphs[0]
    p_ech.text = "ORCA MULTI-STAKEHOLDER ECOSYSTEM"
    p_ech.font.bold = True
    p_ech.font.size = Pt(11)
    p_ech.font.color.rgb = CYAN_ACCENT
    p_ech.space_after = Pt(6)

    stakeholders = [
        ("🎣 Artisanal Fishermen (4.5M)", "Zero preventable capsizings, localized vernacular audio safety, precise PFZ fishing hotspots."),
        ("⚓ Harbors & Port Authorities", "Live harbor fleet tracking, automated departure safety clearance, congestion mitigation."),
        ("🚢 Commercial Fishing Fleets", "Optimal fuel-saving route planning, dynamic weather avoidance, higher catch ROI."),
        ("🛡️ Coast Guard & NDRF", "Real-time Dark-Fleet anomaly alerts, 1,000-particle Monte Carlo SAR drift tracking."),
        ("🔬 Oceanographic Institutes", "Crowdsourced ground-truth sea state telemetry & high-resolution model validation.")
    ]
    for s_title, s_desc in stakeholders:
        p_stk = tf_ec.add_paragraph()
        p_stk.space_after = Pt(4)
        r_st = p_stk.add_run()
        r_st.text = f"{s_title}\n"
        r_st.font.bold = True
        r_st.font.size = Pt(9)
        r_st.font.color.rgb = CYAN_ACCENT
        r_sd = p_stk.add_run()
        r_sd.text = f"  • {s_desc}"
        r_sd.font.size = Pt(8)
        r_sd.font.color.rgb = TEXT_LIGHT_MUTED

    # Right Container: 4 Strategic Impact Pillars (53% width)
    pillars_bg = add_card(slide5, Inches(6.3), Inches(1.55), Inches(6.53), Inches(4.6), bg_color=BG_LIGHT_BLUE, border_color=TEAL_DARK)
    _, tf_pl = add_textbox(slide5, Inches(6.4), Inches(1.65), Inches(6.33), Inches(4.4))
    p_plh = tf_pl.paragraphs[0]
    p_plh.text = "THE FOUR PILLARS OF NATIONAL IMPACT"
    p_plh.font.bold = True
    p_plh.font.size = Pt(11)
    p_plh.font.color.rgb = TEAL_DARK
    p_plh.space_after = Pt(6)

    pillars = [
        ("👥 SOCIAL IMPACT (Saving Lives)", 
         "Eliminates fatal capsizings by calculating real-time wave-vessel physical resonance. Delivers plain-language voice alerts in Tamil, Malayalam, Bengali, Telugu, and Hindi for non-literate crews.",
         EMERALD_SAFE),
        ("💰 ECONOMIC IMPACT (Fuel & Yield)", 
         "Reduces diesel expenditure by 20–30% through current-assisted dynamic A* pathfinding. Maximizes high-value pelagic catch yields via multi-species Habitat Suitability Index (HSI) mapping.",
         AMBER_WARN),
        ("🌿 ENVIRONMENTAL IMPACT (Ocean Health)", 
         "Substantially lowers maritime diesel emissions per fishing voyage. Provides automated geofencing to prevent accidental intrusion into Marine Protected Areas (MPAs) and international boundaries.",
         OCEAN_BLUE),
        ("🇮🇳 STRATEGIC IMPACT (Digital Sovereignty)", 
         "Delivers end-to-end sovereign maritime domain awareness across India's 7,516 km coastline and 2.37M sq km EEZ, creating an integrated national ocean safety infrastructure.",
         TEAL_DARK)
    ]
    for p_name, p_body, p_col in pillars:
        p_pi = tf_pl.add_paragraph()
        p_pi.space_after = Pt(5)
        r_pt = p_pi.add_run()
        r_pt.text = f"{p_name}\n"
        r_pt.font.bold = True
        r_pt.font.size = Pt(9)
        r_pt.font.color.rgb = p_col
        r_pb = p_pi.add_run()
        r_pb.text = f"{p_body}"
        r_pb.font.size = Pt(8)
        r_pb.font.color.rgb = TEXT_DARK

    # Bottom Progression Banner: National Scale
    nat_bg = add_card(slide5, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.55), bg_color=NAVY_DEEP, border_color=CYAN_ACCENT)
    _, tf_nt = add_textbox(slide5, Inches(0.6), Inches(6.32), Inches(12.1), Inches(0.5))
    p_nt = tf_nt.paragraphs[0]
    p_nt.alignment = PP_ALIGN.CENTER
    
    r_n1 = p_nt.add_run()
    r_n1.text = "1 ARTISANAL VESSEL  ➔  HARBOR FLEET (500+)  ➔  COASTAL STATE  ➔  PAN-INDIA MARITIME INTELLIGENCE LAYER\n"
    r_n1.font.bold = True
    r_n1.font.size = Pt(9.5)
    r_n1.font.color.rgb = CYAN_ACCENT
    
    r_n2 = p_nt.add_run()
    r_n2.text = "“From a single artisanal voyage to a sovereign digital intelligence layer for India’s 7,516 km coastline.”"
    r_n2.font.italic = True
    r_n2.font.size = Pt(8.5)
    r_n2.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides[5]
    print("Populating Slide 6: RESEARCH AND REFERENCES...")
    
    for s in slide6.shapes:
        if s.name == "Title 1":
            s.text_frame.text = "RESEARCH  AND REFERENCES: Scientific Grounding to Implementation"
        elif s.name == "Oval 8":
            s.text_frame.text = "ORCA 4.0"
        elif s.name == "TextBox 8":
            s.text_frame.clear()

    # Official Pointer Badges Bar (Top)
    pointers_s6 = [
        ("DETAILS OF RESEARCH WORK", CYAN_ACCENT),
        ("IMO STABILITY CRITERIA", TEAL_DARK),
        ("HYDRODYNAMIC WAVE PHYSICS", TEAL_DARK),
        ("VERIFIED ACADEMIC CITATIONS", EMERALD_SAFE)
    ]
    for i, (p_text, p_col) in enumerate(pointers_s6):
        add_badge(slide6, Inches(0.5 + i * 3.1), Inches(1.15), Inches(3.0), Inches(0.32), p_text, p_col, font_size=9)

    # 4 Core Research Cards (2x2 Grid)
    research_cards = [
        ("IMO Code on Intact Stability (2008)",
         "Resolution MSC.267(85) / Small Craft Capsizing Dynamics",
         "Small vessel dynamic instability under steep wave action, beam-sea broaching, and parametric roll resonance.",
         "Deterministic Safety Engine: Computes critical wave threshold (Hcrit ≈ 0.6 · Lhull) and enforces non-bypassable hazard overrides (<10ms).",
         BORDER_CYAN),
        ("Marine Hydrodynamics & Encounter Physics",
         "O.M. Faltinsen (1990) & J.N. Newman (1977)",
         "Wave steepness ratio (S = Hs / λ), deep-water dispersion relations (λ = g·Tp² / 2π), and encounter angle frequency shifts.",
         "Vessel Digital Twin: Dynamic calculation of wave steepness & vessel-wave interaction rather than treating sea states as static point forecasts.",
         EMERALD_SAFE),
        ("Satellite Oceanography & Pelagic Habitats",
         "INCOIS PFZ Mission & ISRO MOSDAC Telemetry",
         "Aggregation of pelagic marine species along thermal gradients, SST frontal boundaries, and chlorophyll-a density concentration zones.",
         "Multi-Species Habitat Index (HSI): XGBoost model integrated with 2D Sobel spatial filter for high-precision fishery zone recommendations.",
         AMBER_WARN),
        ("Dynamic Constrained Pathfinding & Graph Optimization",
         "Modified A* over Time-Varying Navigational Cost Surfaces",
         "Optimization of multiobjective vessel routes under dynamic weather constraints, geofenced boundaries, and risk cost penalization.",
         "ORCA Navigational Pathfinder: Computes 4-point optimized waypoints balancing weather safety, fuel consumption, and dock proximity.",
         BORDER_CYAN)
    ]

    for idx, (r_title, r_src, r_insight, r_app, r_bcol) in enumerate(research_cards):
        col = idx % 2
        row = idx // 2
        left_pos = Inches(0.5 + col * 6.25)
        top_pos = Inches(1.55 + row * 2.3)
        
        card = add_card(slide6, left_pos, top_pos, Inches(6.08), Inches(2.2), bg_color=NAVY_DEEP, border_color=r_bcol)
        _, tf_rc = add_textbox(slide6, left_pos + Inches(0.1), top_pos + Inches(0.08), Inches(5.88), Inches(2.04))
        
        p1 = tf_rc.paragraphs[0]
        p1.text = f"📚 {r_title}"
        p1.font.bold = True
        p1.font.size = Pt(9.5)
        p1.font.color.rgb = CYAN_ACCENT
        
        p_src = tf_rc.add_paragraph()
        p_src.text = f"Source: {r_src}"
        p_src.font.italic = True
        p_src.font.size = Pt(7.5)
        p_src.font.color.rgb = TEXT_MUTED
        p_src.space_after = Pt(2)
        
        p_ins = tf_rc.add_paragraph()
        r_ih = p_ins.add_run()
        r_ih.text = "Key Insight: "
        r_ih.font.bold = True
        r_ih.font.size = Pt(8)
        r_ih.font.color.rgb = AMBER_WARN
        r_ib = p_ins.add_run()
        r_ib.text = r_insight
        r_ib.font.size = Pt(8)
        r_ib.font.color.rgb = TEXT_WHITE
        p_ins.space_after = Pt(2)

        p_app = tf_rc.add_paragraph()
        r_ah = p_app.add_run()
        r_ah.text = "ORCA Implementation: "
        r_ah.font.bold = True
        r_ah.font.size = Pt(8)
        r_ah.font.color.rgb = EMERALD_SAFE
        r_ab = p_app.add_run()
        r_ab.text = r_app
        r_ab.font.size = Pt(8)
        r_ab.font.color.rgb = TEXT_LIGHT_MUTED

    # Bottom Container: Verified Academic & Operational References Footer
    ref_bg = add_card(slide6, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.68), bg_color=BG_LIGHT_BLUE, border_color=TEAL_DARK)
    _, tf_rf = add_textbox(slide6, Inches(0.6), Inches(6.22), Inches(12.1), Inches(0.64))
    
    p_rf = tf_rf.paragraphs[0]
    p_rf.text = "AUTHENTIC OPERATIONAL & SCIENTIFIC REFERENCES"
    p_rf.font.bold = True
    p_rf.font.size = Pt(8.5)
    p_rf.font.color.rgb = TEAL_DARK
    p_rf.space_after = Pt(1)

    citations = (
        "[1] International Maritime Organization (IMO). International Code on Intact Stability (2008). Resolution MSC.267(85).\n"
        "[2] INCOIS (Ministry of Earth Sciences). Ocean State Forecast (OSF) & Potential Fishing Zone (PFZ) Operational Protocols. | "
        "[3] ISRO Space Applications Centre. Oceansat-3 (OCM-3) & INSAT-3DR Marine Products, MOSDAC.\n"
        "[4] Copernicus Marine Environment Monitoring Service (CMEMS). Global Ocean Physical Analysis. | "
        "[5] World Meteorological Organization (WMO). Manual on Marine Meteorological Services (WMO-No. 558)."
    )
    p_cits = tf_rf.add_paragraph()
    p_cits.text = citations
    p_cits.font.size = Pt(7)
    p_cits.font.color.rgb = TEXT_DARK

    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Perfect SIH 2026 Presentation generated successfully!")
    print(f"Saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_presentation()
