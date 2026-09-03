import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 widescreen slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_COLOR = RGBColor(2, 11, 20)          # #020b14 Dark Ocean
    CARD_BG = RGBColor(5, 25, 45)           # Deep Card Blue
    CYAN_PRIMARY = RGBColor(6, 182, 212)    # #06b6d4 Bright Cyan
    EMERALD_ACCENT = RGBColor(52, 211, 153) # #34d399 Emerald Safe
    AMBER_ACCENT = RGBColor(251, 191, 36)   # #fbbf24 Warning Amber
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)    # Slate 400

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, title_text, subtitle_text=""):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = "SIH26176 · ISRO / INCOIS BLUE ECONOMY HACKATHON"
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = CYAN_PRIMARY

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 1: Title Slide
    # ---------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide1)

    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.8))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "SMART INDIA HACKATHON 2026 · NATIONAL FINALS"
    p_badge.font.size = Pt(14)
    p_badge.font.bold = True
    p_badge.font.color.rgb = CYAN_PRIMARY

    p_title = tf1.add_paragraph()
    p_title.text = "ORCA 4.0"
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    p_sub = tf1.add_paragraph()
    p_sub.text = "Universal Autonomous Marine Operating System & Bio-Physical Intelligence Platform"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = CYAN_PRIMARY

    p_desc = tf1.add_paragraph()
    p_desc.text = "\nProblem Statement ID: SIH26176 | Ministry: ISRO / Department of Space & INCOIS\nEmpowering 4.5M Artisanal Fishermen with Edge Deterministic Safety & ISRO Satellite Telemetry"
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 2: Ground Reality & Critical Problems
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide2)
    add_header(slide2, "The Ground Reality: 4.5M Vulnerable Lives at Sea", "Existing marine systems fail due to generalized broadcasts, data silos, and zero offshore connectivity.")

    cards_data2 = [
        ("800+ Annual Fatalities", "Traditional SMS broadcasts ignore vessel hull length, draft, and wave steepness, leading to capsizing.", AMBER_ACCENT),
        ("₹45,000 Cr Fuel Lost", "Artisanal vessels burn 40% of their diesel searching blindly for schooling pelagic fish without spatial guidance.", CYAN_PRIMARY),
        ("Disconnected Data Silos", "ISRO satellite rasters, INCOIS PFZ bulletins, and Coast Guard radar operate in complete isolation.", TEXT_WHITE),
        ("Offshore Signal Blackout", "Beyond 12 nautical miles, 4G/5G mobile signals drop to 0%, rendering standard web apps completely dead.", EMERALD_ACCENT),
    ]

    for i, (title, desc, color) in enumerate(cards_data2):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(2.0 + row * 2.4)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.2), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 3: The ORCA 4.0 Solution
    # ---------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide3)
    add_header(slide3, "The Breakthrough Solution: Edge-Native Maritime OS", "ORCA 4.0 unifies cloud multi-agent reasoning with sub-100ms deterministic edge physics.")

    pillars = [
        ("1. Deterministic Safety (<100ms)", "Zero LLM hallucinations in safety. Mathematical capsizing thresholds: H_crit = 0.22*L + 0.05*B."),
        ("2. Multi-Species HSI Matrix", "Dynamic habitat suitability modeling for Pomfret, Mackerel, Tuna & Sardine via INCOIS OCM-3."),
        ("3. Multi-Objective Route Pareto", "Solves 3 trajectories simultaneously: Safest Detour, Lowest Fuel, and Maximum Net Catch Value."),
        ("4. 16-Byte LoRa Telemetry", "Bit-packed protocol transmitted over 868MHz LoRa & NavIC satellite receivers with HMAC-SHA256."),
        ("5. Regional Dialect Voice First", "Native spoken guidance in Marathi (Koli/Malvani), Hindi, Gujarati, Tamil & English."),
        ("6. Coast Guard OSINT Hub", "Sentinel-1 SAR Dark Fleet radar correlation, 1,000-particle Monte Carlo SAR drift, & audit ledgers.")
    ]

    for i, (title, desc) in enumerate(pillars):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(2.0 + row * 2.4)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CYAN_PRIMARY
        card.line.width = Pt(1)

        tb = slide3.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.4), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 4: Architecture & 16-Agent DAG
    # ---------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide4)
    add_header(slide4, "16-Agent Directed Acyclic Graph (DAG) Architecture", "Modular microservice pipeline connected via typed, asynchronous event bus with zero circular locks.")

    arch_points = [
        ("Data Ingestion Tier", "INCOIS ERDDAP REST client + Open-Meteo ocean telemetry + 2D Spatial DINEOF Satellite Cloud Gap Filler."),
        ("Spatio-Temporal World Model", "Assembles Vessel Digital Twin, dynamic ocean currents, wave steepness, and legal maritime geofences into unified state."),
        ("Edge Deterministic Safety Engine", "Sub-100ms circuit breaker. Evaluates wave-hull resonance and triggers immediate hard overrides if capsizing risk detected."),
        ("Multi-Objective Pareto Engine", "Calculates hydro-acoustic slip, brake specific fuel consumption (BSFC), and wholesale harbor auction price arbitrage."),
        ("Institutional Governance Tier", "Coast Guard SAR drift centroid calculation, Sentinel-1 radar anomaly matcher, and immutable override audit ledger.")
    ]

    for i, (title, desc) in enumerate(arch_points):
        y = Inches(2.0 + i * 0.95)
        bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CARD_BG
        bar.line.color.rgb = CYAN_PRIMARY
        bar.line.width = Pt(1)

        tb = slide4.shapes.add_textbox(Inches(1.0), y + Inches(0.05), Inches(11.3), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {title}: "
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = CYAN_PRIMARY
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 5: Space & Satellite Data Grounding
    # ---------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide5)
    add_header(slide5, "Deep Space Telemetry: ISRO & INCOIS Integration", "Direct integration with official Indian and European Space Agency earth observation missions.")

    sats = [
        ("Oceansat-3 (OCM-3)", "Ocean Color Monitor", "360m spatial resolution Chlorophyll-a layers for pelagic plankton bloom and bio-thermal front identification."),
        ("INSAT-3DR / MODIS", "Geostationary Thermal SST", "Continuous Sea Surface Temperature thermal gradient tracking to detect upwelling zones where fish school."),
        ("Sentinel-1 C-Band SAR", "Synthetic Aperture Radar", "Penetrates monsoonal cloud cover to detect Dark Fleet vessels (no-AIS) and offshore oil slicks."),
        ("2D Spatial DINEOF Engine", "Cloud Gap Reconstruction", "Gaussian covariance kernel reconstruction algorithm reconstructing 100% of cloudy satellite observation gaps.")
    ]

    for i, (name, role, desc) in enumerate(sats):
        x = Inches(0.8 + i * 2.95)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CYAN_PRIMARY
        card.line.width = Pt(1.5)

        tb = slide5.shapes.add_textbox(x + Inches(0.15), Inches(2.2), Inches(2.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p2 = tf.add_paragraph()
        p2.text = role
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = CYAN_PRIMARY
        p3 = tf.add_paragraph()
        p3.text = f"\n{desc}"
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 6: Multi-Objective Routing & Eco-Economics
    # ---------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide6)
    add_header(slide6, "Eco-Economic Optimization & Wholesale Arbitrage", "Maximizing net livelihood profit by uniting physics fuel models with real-time wholesale auction rates.")

    eco_cards = [
        ("Multi-Objective Cost Formula", "Cost = w1*Hazard + w2*Fuel(BSFC, Slip) + w3*Weather - w4*Catch_Value\nSolves 3 Pareto-optimal trajectories simultaneously.", CYAN_PRIMARY),
        ("Wholesale Harbor Arbitrage", "Ingests daily auction rates across 12+ coastal harbors (e.g. Mirkarwada, Sassoon Dock, Panaji) to identify highest net revenue port.", EMERALD_ACCENT),
        ("Hydrodynamic Digital Twin", "Computes propeller slip ratio, BSFC (240 g/hp-hr), and sea-state resistance for realistic liter-accurate diesel estimations.", TEXT_WHITE)
    ]

    for i, (title, desc, color) in enumerate(eco_cards):
        y = Inches(2.0 + i * 1.6)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = slide6.shapes.add_textbox(Inches(1.0), y + Inches(0.15), Inches(11.3), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 7: Zero-Connectivity LoRa Telemetry Protocol
    # ---------------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide7)
    add_header(slide7, "Zero-Connectivity Resilience: 16-Byte LoRa Protocol", "Sub-gigahertz bit-packed communication ensuring uninterrupted life-saving safety beyond cellular reach.")

    lora_bytes = [
        ("Byte 0", "Header / Packet Type (0x02 Alert, 0x01 Pos, 0x03 SOS)"),
        ("Bytes 1–3", "Latitude packed with millidegree spatial resolution"),
        ("Bytes 4–6", "Longitude packed with millidegree spatial resolution"),
        ("Byte 7", "Safety Risk Score (0–100 Seaworthiness Index)"),
        ("Bytes 8–11", "Recommended Waypoint / PFZ Target Grid Coordinates"),
        ("Bytes 12–15", "Cryptographic HMAC-SHA256 Integrity Signature")
    ]

    for i, (b_name, b_desc) in enumerate(lora_bytes):
        x = Inches(0.8 + (i % 3) * 3.9)
        y = Inches(2.0 + (i // 3) * 2.4)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CYAN_PRIMARY
        card.line.width = Pt(1.5)

        tb = slide7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = b_name
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = CYAN_PRIMARY
        pd = tf.add_paragraph()
        pd.text = f"\n{b_desc}"
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 8: National Security & Governance Hub
    # ---------------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide8)
    add_header(slide8, "Coast Guard SAR & Maritime Governance Hub", "Integrated Search & Rescue drift prediction and illicit Dark Fleet surveillance.")

    gov_cards = [
        ("1,000-Particle Monte Carlo SAR Drift", "Simulates Stokes wave drift, wind leeway (3%), and surface currents to generate search radius ellipses for missing craft.", CYAN_PRIMARY),
        ("Dark Fleet Sentinel-1 Radar Match", "Cross-correlates radar cross-sections against live AIS signals to flag unauthorized foreign trawlers.", AMBER_ACCENT),
        ("Immutable Authority Override Ledger", "Non-repudiable audit logging for Coast Guard zone closures, cyclone alerts, and naval restricted zones.", EMERALD_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(gov_cards):
        y = Inches(2.0 + i * 1.6)
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = slide8.shapes.add_textbox(Inches(1.0), y + Inches(0.15), Inches(11.3), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 9: Verification & Technology Benchmarks
    # ---------------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide9)
    add_header(slide9, "Engineering Rigor & Verification Benchmarks", "Production-grade codebase tested with comprehensive automated unit, integration, and physics suites.")

    metrics = [
        ("44 / 44", "PyTest Suites Passed", "100% test coverage across safety circuit breakers, geofences, and SAR physics.", EMERALD_ACCENT),
        ("< 100ms", "Deterministic Edge Verdict", "Sub-100 millisecond response time guaranteed for safety-critical evaluations.", CYAN_PRIMARY),
        ("0 Errors", "Vite PWA Build", "Clean production bundle with offline vector tile caching and service workers.", TEXT_WHITE),
        ("~50 Ports", "Global Ocean Coverage", "Worldwide coastal support across Arabian Sea, Bay of Bengal, Pacific, and Atlantic.", CYAN_PRIMARY)
    ]

    for i, (stat, label, desc, color) in enumerate(metrics):
        x = Inches(0.8 + i * 2.95)
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = slide9.shapes.add_textbox(x + Inches(0.15), Inches(2.2), Inches(2.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = stat
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p3 = tf.add_paragraph()
        p3.text = f"\n{desc}"
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 10: Impact, Blue Economy & Future Vision
    # ---------------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide10)
    add_header(slide10, "National Blue Economy Impact & Scalability Vision", "Transforming coastal livelihoods and setting a global benchmark in space-enabled marine operating systems.")

    vision_cards = [
        ("Economic Transformation", "Estimated ₹12,000–₹18,000 monthly fuel savings per vessel. ₹5,400+ Crore annual blue economy savings nationwide.", EMERALD_ACCENT),
        ("Zero Preventable Casualties", "Complete elimination of weather-induced capsizing in registered craft through hull-specific deterministic physics.", CYAN_PRIMARY),
        ("Pan-India Coastal Deployment", "Phased rollout with state fisheries departments across Maharashtra, Goa, Gujarat, Kerala, Tamil Nadu, and Odisha.", TEXT_WHITE),
        ("Global Export Potential", "Extensible architecture designed for artisanal coastal communities across South Asia, Africa, and Latin America.", CYAN_PRIMARY)
    ]

    for i, (title, desc, color) in enumerate(vision_cards):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(2.0 + row * 2.4)
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = slide10.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.2), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = TEXT_MUTED

    # Save presentation
    output_dir = "/Users/subham/code/orca"
    os.makedirs(output_dir, exist_ok=True)
    deck_path = os.path.join(output_dir, "ORCA_4.0_SIH_National_Finals_Deck.pptx")
    prs.save(deck_path)
    print(f"[SUCCESS] SIH PowerPoint presentation saved to: {deck_path}")

if __name__ == "__main__":
    create_deck()
