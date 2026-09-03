import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches

TEMPLATE_PATH = "/Users/subham/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx"
OUTPUT_PATH = "/Users/subham/Downloads/ORCA_4.0_SIH2026_Official_Deck.pptx"
ASSETS_DIR = "/Users/subham/code/orca/deck_assets"
os.makedirs(ASSETS_DIR, exist_ok=True)

# Color Palette
CLR_BG_DARK = "#0a192f"
CLR_CARD_DARK = "#0f2744"
CLR_CARD_LIGHT = "#f0f7ff"
CLR_CYAN = "#00d2ff"
CLR_BLUE = "#0077b6"
CLR_EMERALD = "#10b981"
CLR_AMBER = "#f59e0b"
CLR_ROSE = "#ef4444"
CLR_TEXT_WHITE = "#ffffff"
CLR_TEXT_MUTED = "#94a3b8"
CLR_TEXT_DARK = "#0f172a"
CLR_BORDER = "#1e3a5f"

def create_slide1_graphic():
    fig, ax = plt.subplots(figsize=(8.2, 7.2), dpi=300)
    fig.patch.set_facecolor(CLR_BG_DARK)
    ax.set_facecolor(CLR_BG_DARK)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')

    # Main Card Box
    card = patches.FancyBboxPatch((2, 2), 96, 96, boxstyle="round,pad=1.5,rounding_size=3",
                                  facecolor=CLR_CARD_DARK, edgecolor=CLR_CYAN, linewidth=2)
    ax.add_patch(card)

    # Subtitle Tag
    ax.text(50, 92, "SMART INDIA HACKATHON 2026 · PROJECT SHOWCASE", color=CLR_CYAN,
            fontsize=11, fontweight='bold', ha='center', va='center')

    # Title
    ax.text(50, 84, "ORCA 4.0", color=CLR_TEXT_WHITE, fontsize=32, fontweight='bold', ha='center', va='center')

    # Tagline
    ax.text(50, 77, "“FROM OCEAN DATA TO SAFER DECISIONS”", color=CLR_CYAN,
            fontsize=14, fontweight='bold', ha='center', va='center')

    # Sub-descriptor
    ax.text(50, 71.5, "AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform",
            color=CLR_TEXT_MUTED, fontsize=10.5, ha='center', va='center')

    # Flow Banner
    fb = patches.FancyBboxPatch((5, 60), 90, 8, boxstyle="round,pad=0.5,rounding_size=2",
                                facecolor="#062038", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(fb)
    ax.text(50, 64, "OCEAN ➔ WEATHER ➔ VESSEL ➔ ROUTE ➔ ORCA ➔ RISK ➔ DECISION",
            color="#38bdf8", fontsize=9.5, fontweight='bold', ha='center', va='center')

    # 3 Key Value Cards
    cards_data = [
        ("🛡️ DETERMINISTIC SAFETY SHIELD", "Hard physics capsizing check (Hcrit = 0.6·L) with zero AI hallucination risk", CLR_EMERALD, 44),
        ("🚢 PARAMETRIC VESSEL TWIN", "Dynamic wave encounter physics tailored to vessel hull length, draft & beam", CLR_CYAN, 28),
        ("📻 MULTI-HOP LORA FLEET MESH", "Overcomes 12 NM 4G blackout for zero-cost offshore telemetry up to 50 km", CLR_AMBER, 12)
    ]

    for title, desc, col, y_pos in cards_data:
        box = patches.FancyBboxPatch((5, y_pos), 90, 13, boxstyle="round,pad=0.8,rounding_size=2",
                                     facecolor="#061c33", edgecolor=col, linewidth=1.5)
        ax.add_patch(box)
        ax.text(8, y_pos + 8.5, title, color=col, fontsize=11, fontweight='bold', va='center')
        ax.text(8, y_pos + 3.8, desc, color=CLR_TEXT_WHITE, fontsize=9, va='center')

    img_path = os.path.join(ASSETS_DIR, "slide1_hero.png")
    plt.tight_layout()
    plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight', pad_inches=0.05)
    plt.close()
    return img_path

def create_slide2_graphic():
    fig, ax = plt.subplots(figsize=(15.5, 6.8), dpi=300)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')

    # Top Thesis Box
    th_box = patches.FancyBboxPatch((1, 88), 98, 10.5, boxstyle="round,pad=0.5,rounding_size=2",
                                    facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(th_box)
    ax.text(3, 94.5, "CORE PARADIGM SHIFT:  FORECAST ≠ OPERATIONAL RISK", color=CLR_AMBER, fontsize=12, fontweight='bold', va='center')
    ax.text(3, 90.5, "Existing systems report 'What is happening at sea?'. ORCA computes 'What does it mean for THIS vessel, on THIS route, at THIS time?'",
            color=CLR_TEXT_WHITE, fontsize=10, va='center')

    # Left Container: Data Transformation (58% width)
    left_bg = patches.FancyBboxPatch((1, 12), 58, 73, boxstyle="round,pad=0.8,rounding_size=2",
                                     facecolor="#0a192f", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(left_bg)
    ax.text(3, 81.5, "VISUAL TRANSFORMATION: FROM DATA CHAOS TO OPERATIONAL DECISION",
            color=CLR_CYAN, fontsize=11, fontweight='bold', va='center')

    # 8 Input Nodes (2 columns)
    inputs_col1 = [
        ("🌊 WAVES", "Hs, Swell Tp, Period"),
        ("💨 WIND", "Vel, Gusts, Direction"),
        ("🌀 CYCLONES", "IMD Tracks & Alerts"),
        ("🌊 CURRENTS", "Surface Vectors (u,v)")
    ]
    inputs_col2 = [
        ("🛰️ SATELLITES", "SST & Chlorophyll-a"),
        ("🚢 VESSEL", "Length, Draft, Beam"),
        ("📡 AIS FEEDS", "MMSI Coordinates"),
        ("🗺️ ROUTE", "Waypoints & Port Docks")
    ]

    for idx, (title, sub) in enumerate(inputs_col1):
        y = 66 - idx * 14.5
        box = patches.FancyBboxPatch((3, y), 13, 11.5, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#061c33", edgecolor="#38bdf8", linewidth=1)
        ax.add_patch(box)
        ax.text(9.5, y + 7.5, title, color=CLR_CYAN, fontsize=8.5, fontweight='bold', ha='center', va='center')
        ax.text(9.5, y + 3.2, sub, color="#cbd5e1", fontsize=7, ha='center', va='center')

    for idx, (title, sub) in enumerate(inputs_col2):
        y = 66 - idx * 14.5
        box = patches.FancyBboxPatch((17.5, y), 13, 11.5, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#061c33", edgecolor="#38bdf8", linewidth=1)
        ax.add_patch(box)
        ax.text(24, y + 7.5, title, color=CLR_CYAN, fontsize=8.5, fontweight='bold', ha='center', va='center')
        ax.text(24, y + 3.2, sub, color="#cbd5e1", fontsize=7, ha='center', va='center')

    # Convergence Arrows
    ax.annotate("", xy=(33.5, 45), xytext=(31, 45),
                arrowprops=dict(arrowstyle="->", color=CLR_CYAN, lw=2.5))

    # Central Engine
    eng_box = patches.FancyBboxPatch((34, 32), 23.5, 26, boxstyle="round,pad=0.6,rounding_size=2",
                                     facecolor="#0284c7", edgecolor=CLR_CYAN, linewidth=2)
    ax.add_patch(eng_box)
    ax.text(45.7, 52, "⚡ ORCA 4.0 ENGINE", color=CLR_TEXT_WHITE, fontsize=11, fontweight='bold', ha='center', va='center')
    ax.text(45.7, 45, "• 0.083° Grid Normalization\n• Parametric Vessel Twin\n• Multi-Species Habitat Matrix",
            color=CLR_TEXT_WHITE, fontsize=8.5, ha='center', va='center')

    # Output decisions
    out_cards = [
        ("🛡️ DETERMINISTIC RISK", "Hcrit capsizing check (<10ms)", CLR_EMERALD, 66),
        ("🗺️ ROUTE RISK FIELD", "Dynamic 2D risk surface A*", CLR_CYAN, 45),
        ("📢 ACTIONABLE VERDICT", "Clear Go/No-Go + Audio Voice", CLR_AMBER, 24)
    ]
    for title, sub, col, y in out_cards:
        # Mini arrow
        ax.annotate("", xy=(34, y + 5), xytext=(31, 45),
                    arrowprops=dict(arrowstyle="->", color=col, lw=1.5, alpha=0.7))

    # Right Container: Competitive Matrix (40% width)
    right_bg = patches.FancyBboxPatch((61, 12), 38, 73, boxstyle="round,pad=0.8,rounding_size=2",
                                      facecolor="#f0f7ff", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(right_bg)
    ax.text(63, 81.5, "COMPETITIVE DIFFERENTIATION MATRIX", color="#0d47a1", fontsize=11, fontweight='bold', va='center')

    comp_rows = [
        ("Capability", "Existing Systems", "ORCA 4.0"),
        ("Weather Forecast", "✅ Regional Only", "✅ High-Res (0.083°)"),
        ("Cyclone Warnings", "✅ Generic SMS", "✅ Dynamic Geofenced"),
        ("Vessel Hull Context", "❌ Ignored", "✅ Vessel Twin (Lwl, Draft)"),
        ("Wave Physics Math", "❌ None", "✅ Hcrit & Resonance Math"),
        ("Safe Route Pathfinder", "❌ Static Path", "✅ Dynamic A* Avoidance"),
        ("Deterministic Safety", "❌ Unverified", "✅ Non-bypassable Rule"),
        ("Offshore Reach", "❌ Dead >12 NM", "✅ LoRa Mesh (50 km)")
    ]

    for idx, (c1, c2, c3) in enumerate(comp_rows):
        y = 73 - idx * 7.6
        if idx == 0:
            ax.text(63, y, c1, color="#0d47a1", fontsize=9, fontweight='bold', va='center')
            ax.text(77, y, c2, color="#0d47a1", fontsize=9, fontweight='bold', va='center')
            ax.text(89, y, c3, color="#0d47a1", fontsize=9, fontweight='bold', va='center')
            line = patches.ConnectionPatch((63, y - 2.5), (97, y - 2.5), "data", "data", color="#0284c7", lw=1)
            ax.add_patch(line)
        else:
            ax.text(63, y, c1, color=CLR_TEXT_DARK, fontsize=8, va='center')
            c2_col = CLR_ROSE if "❌" in c2 else "#64748b"
            ax.text(77, y, c2, color=c2_col, fontsize=7.5, va='center')
            ax.text(89, y, c3, color=CLR_EMERALD, fontsize=8, fontweight='bold', va='center')

    # Bottom Innovation Strip
    btm_box = patches.FancyBboxPatch((1, 1.5), 98, 8.5, boxstyle="round,pad=0.5,rounding_size=2",
                                     facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(btm_box)
    ax.text(50, 5.7, "① MULTI-SOURCE HARMONIZATION  ➔  ② VESSEL-AWARE PHYSICS TWIN  ➔  ③ ROUTE RISK OPTIMIZATION  ➔  ④ EXPLAINABLE SAFETY SHIELD",
            color=CLR_CYAN, fontsize=9.5, fontweight='bold', ha='center', va='center')

    img_path = os.path.join(ASSETS_DIR, "slide2_diagram.png")
    plt.tight_layout()
    plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight', pad_inches=0.05)
    plt.close()
    return img_path

def create_slide3_graphic():
    fig, ax = plt.subplots(figsize=(15.5, 6.8), dpi=300)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')

    # Top Architecture Background Box
    top_bg = patches.FancyBboxPatch((1, 40), 98, 58, boxstyle="round,pad=0.8,rounding_size=2",
                                    facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(top_bg)

    # 6 Ingestion Nodes
    sources = [
        ("INCOIS ERDDAP", "Wave & Current"),
        ("IMD ALERTS", "Cyclone Tracks"),
        ("COPERNICUS", "0.083° Marine Grids"),
        ("ISRO MOSDAC", "SST & Oceansat-3"),
        ("AIS STREAMS", "MMSI Positions"),
        ("NMEA SENSORS", "Onboard GPS/IMU")
    ]
    for idx, (name, desc) in enumerate(sources):
        x = 2.5 + idx * 16
        box = patches.FancyBboxPatch((x, 82), 14.8, 13, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#061c33", edgecolor=CLR_CYAN, linewidth=1)
        ax.add_patch(box)
        ax.text(x + 7.4, 90, f"📡 {name}", color=CLR_CYAN, fontsize=8, fontweight='bold', ha='center', va='center')
        ax.text(x + 7.4, 85.5, desc, color="#cbd5e1", fontsize=7, ha='center', va='center')

    # Normalization Funnel Bar
    funnel = patches.FancyBboxPatch((2.5, 71.5), 95, 8, boxstyle="round,pad=0.4,rounding_size=1.5",
                                    facecolor="#0284c7", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(funnel)
    ax.text(50, 75.5, "⚡ DATA ENGINE & DUAL DIGITAL STATE: Spatial H3 Indexing + Environmental State ⟷ Parametric Vessel Twin (Lwl, Draft, Beam)",
            color=CLR_TEXT_WHITE, fontsize=9, fontweight='bold', ha='center', va='center')

    # Tripartite Engines (3 Cards)
    engines = [
        ("⚙️ DETERMINISTIC PHYSICS (<10ms)", 
         "• Critical Wave Height (Hcrit = 0.6·Lhull)\n• Wave Steepness (S = Hs / λ) & Period\n• Beam-Sea Encounter Angle & Roll Resonance",
         CLR_EMERALD, 2.5),
        ("🧠 PREDICTIVE ML & ANALYTICS", 
         "• Multi-Species Habitat Index (HSI XGBoost)\n• 2D Sobel Filter SST Thermal Fronts\n• 1,000-Particle Monte Carlo SAR Drift",
         CLR_CYAN, 34.5),
        ("📍 DYNAMIC A* OPTIMIZATION", 
         "• 2D Dynamic Risk Cost Surface Mesh\n• A* Navigational Pathfinding & Avoidance\n• Optimal Fuel Burn & Dock Allocation",
         CLR_AMBER, 66.5)
    ]
    for title, desc, col, x in engines:
        box = patches.FancyBboxPatch((x, 50.5), 31, 18.5, boxstyle="round,pad=0.5,rounding_size=1.5",
                                     facecolor="#061c33", edgecolor=col, linewidth=1.5)
        ax.add_patch(box)
        ax.text(x + 1.5, 65, title, color=col, fontsize=8.5, fontweight='bold', va='center')
        ax.text(x + 1.5, 56.5, desc, color=CLR_TEXT_WHITE, fontsize=7.5, va='center')

    # Non-Bypassable Circuit Breaker Shield
    cb = patches.FancyBboxPatch((2.5, 41.5), 95, 7.5, boxstyle="round,pad=0.4,rounding_size=1.5",
                                facecolor="#061c33", edgecolor=CLR_EMERALD, linewidth=1.5)
    ax.add_patch(cb)
    ax.text(50, 45.2, "🛡️ NON-BYPASSABLE SAFETY CIRCUIT BREAKER: Pure deterministic compiled logic evaluates physical capsizing risk (<10ms). Zero LLM hallucination risk.",
            color=CLR_EMERALD, fontsize=8.5, fontweight='bold', ha='center', va='center')

    # Bottom Left: Working Prototype Simulation (Kanyakumari Case)
    btm_left = patches.FancyBboxPatch((1, 1.5), 47.5, 36.5, boxstyle="round,pad=0.8,rounding_size=2",
                                      facecolor="#f0f7ff", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(btm_left)
    ax.text(3, 34.5, "LIVE WORKING PROTOTYPE FLOW (KANYAKUMARI CASE)", color="#0d47a1", fontsize=10, fontweight='bold', va='center')
    
    k_steps = [
        ("1. Input", "Kanyakumari Port (8.08°N, 77.55°E), 9m FRP Craft"),
        ("2. Live State", "Wave Hs = 2.8m, Swell Tp = 8.2s, Wind = 24 kts"),
        ("3. Physics Check", "Hcrit = 5.4m > 2.8m (Safe from Immediate Capsize)"),
        ("4. Route Engine", "Diverts 3.2 NM East to avoid shallow reef breaker"),
        ("5. Output Verdict", "PROCEED WITH CAUTION (Audio Advisory Generated)")
    ]
    for idx, (sn, sv) in enumerate(k_steps):
        y = 28 - idx * 5.8
        ax.text(3, y, f"{sn}:", color="#0d47a1", fontsize=8, fontweight='bold', va='center')
        ax.text(14, y, sv, color=CLR_TEXT_DARK, fontsize=7.5, va='center')

    # Bottom Right: Production Tech Stack
    btm_right = patches.FancyBboxPatch((51.5, 1.5), 47.5, 36.5, boxstyle="round,pad=0.8,rounding_size=2",
                                       facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(btm_right)
    ax.text(53.5, 34.5, "PRODUCTION TECHNOLOGY STACK", color=CLR_CYAN, fontsize=10, fontweight='bold', va='center')

    t_stack = [
        ("Frontend / UI", "React 18, TypeScript, Vite, MapLibre GL, WebGL, TailwindCSS, PWA"),
        ("Backend / API", "FastAPI, Python 3.10+, Async DAG Orchestrator, SQLite/WAL"),
        ("Intelligence", "NumPy, SciPy, XGBoost (HSI Matrix), Modified A* Pathfinder"),
        ("IoT / Telecom", "ESP32 LoRa Gateway (868/433 MHz), Multi-Hop Mesh (OLSR/BATMAN)")
    ]
    for idx, (cat, spec) in enumerate(t_stack):
        y = 28 - idx * 6.5
        ax.text(53.5, y, f"• {cat}:", color=CLR_CYAN, fontsize=8, fontweight='bold', va='center')
        ax.text(67, y, spec, color=CLR_TEXT_WHITE, fontsize=7.5, va='center')

    img_path = os.path.join(ASSETS_DIR, "slide3_architecture.png")
    plt.tight_layout()
    plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight', pad_inches=0.05)
    plt.close()
    return img_path

def create_slide4_graphic():
    fig, ax = plt.subplots(figsize=(15.5, 6.8), dpi=300)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')

    # Column 1: Implementation Roadmap (24% width)
    c1 = patches.FancyBboxPatch((1, 2), 24, 96, boxstyle="round,pad=0.8,rounding_size=2",
                                facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(c1)
    ax.text(13, 93.5, "IMPLEMENTATION ROADMAP", color=CLR_CYAN, fontsize=10.5, fontweight='bold', ha='center', va='center')

    milestones = [
        ("✓ PROTOTYPE", "FastAPI + React UI working", CLR_EMERALD),
        ("✓ LIVE DATA", "Open-Meteo & IMD ingestion", CLR_EMERALD),
        ("✓ RISK ENGINE", "Deterministic safety breaker", CLR_EMERALD),
        ("✓ ROUTE ENGINE", "A* dynamic pathfinder & HSI", CLR_EMERALD),
        ("◐ HARBOR PILOT", "Kanyakumari & Munambam", CLR_AMBER),
        ("◌ SCALE PHASE", "9 Coastal States & Satcom", CLR_TEXT_MUTED)
    ]
    for idx, (title, desc, col) in enumerate(milestones):
        y = 82 - idx * 13.5
        box = patches.FancyBboxPatch((2.5, y - 4), 21, 10.5, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#061c33", edgecolor=col, linewidth=1)
        ax.add_patch(box)
        ax.text(4, y + 2.5, title, color=col, fontsize=8.5, fontweight='bold', va='center')
        ax.text(4, y - 1.5, desc, color=CLR_TEXT_WHITE, fontsize=7.5, va='center')

    # Column 2: Challenges ➔ Mitigations (44% width)
    c2 = patches.FancyBboxPatch((27, 2), 43, 96, boxstyle="round,pad=0.8,rounding_size=2",
                                facecolor="#f0f7ff", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(c2)
    ax.text(48.5, 93.5, "CRITICAL RISKS ➔ PROVEN MITIGATION STRATEGIES", color="#0d47a1", fontsize=10.5, fontweight='bold', ha='center', va='center')

    challenges = [
        ("Data Latency / Stale Feeds", "Automated TTL monitoring (<30m) & cached offline fallback"),
        ("API Provider Outages", "Multi-source failover cascade (INCOIS ➔ Open-Meteo ➔ MOSDAC)"),
        ("Source Disagreements", "Provenance-weighted ensemble & explicit uncertainty bounds"),
        ("Offshore Signal Blackout", "Multi-Hop LoRa Fleet Mesh (up to 50 km) + Offline PWA Caching"),
        ("Vessel Hull Diversity", "Parametric Digital Twin model library (Catamaran, FRP, Trawler)"),
        ("Model Hallucination", "100% Deterministic Safety Rules (Zero AI in safety path)")
    ]
    for idx, (ct, md) in enumerate(challenges):
        y = 83 - idx * 13.5
        box = patches.FancyBboxPatch((28.5, y - 4.5), 40, 11, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#ffffff", edgecolor="#0284c7", linewidth=1)
        ax.add_patch(box)
        ax.text(30, y + 2.5, f"⚠️ {ct}", color=CLR_ROSE, fontsize=8.5, fontweight='bold', va='center')
        ax.text(30, y - 1.8, f"➔ 🛡️ {md}", color=CLR_TEXT_DARK, fontsize=7.5, va='center')

    # Column 3 Top: Deployment Architecture (28% width)
    c3_top = patches.FancyBboxPatch((72, 51), 27, 47, boxstyle="round,pad=0.8,rounding_size=2",
                                    facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(c3_top)
    ax.text(85.5, 93.5, "DEPLOYMENT ARCHITECTURE", color=CLR_CYAN, fontsize=10, fontweight='bold', ha='center', va='center')

    dep_steps = [
        ("1. Vessel Edge", "GPS / IMU + ESP32 LoRa Node"),
        ("2. Fleet Relay", "Multi-Hop Boat Mesh (3-8 km hops)"),
        ("3. Coastal Mast", "Lighthouse LoRa Gateway (60m mast)"),
        ("4. Cloud Engine", "FastAPI / SQLite / Orchestrator"),
        ("5. User Tiers", "Fisher PWA | Fleet Web | Coast Guard")
    ]
    for idx, (ds, dv) in enumerate(dep_steps):
        y = 84 - idx * 7.2
        ax.text(74, y, f"{ds}:", color=CLR_CYAN, fontsize=8, fontweight='bold', va='center')
        ax.text(83.5, y, dv, color=CLR_TEXT_WHITE, fontsize=7.2, va='center')

    # Column 3 Bottom: Deployment Economics & BOM
    c3_btm = patches.FancyBboxPatch((72, 2), 27, 47, boxstyle="round,pad=0.8,rounding_size=2",
                                    facecolor="#061c33", edgecolor=CLR_EMERALD, linewidth=1.5)
    ax.add_patch(c3_btm)
    ax.text(85.5, 44.5, "DEPLOYMENT ECONOMICS & BOM", color=CLR_EMERALD, fontsize=10, fontweight='bold', ha='center', va='center')

    cost_items = [
        ("ESP32 MCU + LoRa SX1262", "₹1,800 (~$22)"),
        ("GPS NEO-6M Receiver", "₹600 (~$7)"),
        ("IP67 Enclosure + Solar", "₹1,100 (~$13)"),
        ("Total Hardware Cost / Boat", "≈ ₹3,500 ($42)"),
        ("Data & Server API Cost", "₹0 (Open Feeds / Serverless)")
    ]
    for idx, (cn, cv) in enumerate(cost_items):
        y = 35 - idx * 7.2
        ax.text(74, y, f"• {cn}:", color=CLR_TEXT_WHITE, fontsize=7.5, va='center')
        col = CLR_EMERALD if "Total" in cn or "₹0" in cv else CLR_CYAN
        ax.text(90, y, cv, color=col, fontsize=7.5, fontweight='bold', ha='center', va='center')

    img_path = os.path.join(ASSETS_DIR, "slide4_feasibility.png")
    plt.tight_layout()
    plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight', pad_inches=0.05)
    plt.close()
    return img_path

def create_slide5_graphic():
    fig, ax = plt.subplots(figsize=(15.5, 6.8), dpi=300)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')

    # Left Container: ORCA Impact Ecosystem (46% width)
    left_bg = patches.FancyBboxPatch((1, 12), 46, 86, boxstyle="round,pad=0.8,rounding_size=2",
                                     facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(left_bg)
    ax.text(24, 93.5, "ORCA MULTI-STAKEHOLDER ECOSYSTEM", color=CLR_CYAN, fontsize=11, fontweight='bold', ha='center', va='center')

    stakeholders = [
        ("🎣 Artisanal Fishermen (4.5M Lives)", "Zero preventable capsizings, localized vernacular audio safety alerts, precision PFZ fishing zones."),
        ("⚓ Harbors & Port Authorities", "Real-time harbor fleet tracking, automated departure safety clearances, harbor congestion mitigation."),
        ("🚢 Commercial Fishing Fleets", "Optimal fuel-saving route planning, dynamic weather avoidance, higher catch yield ROI."),
        ("🛡️ Coast Guard & NDRF", "Real-time Dark-Fleet anomaly alerts, 1,000-particle Monte Carlo SAR drift tracking."),
        ("🔬 Oceanographic Institutes", "Crowdsourced ground-truth sea state telemetry & high-resolution model validation.")
    ]
    for idx, (st, sd) in enumerate(stakeholders):
        y = 82 - idx * 15
        box = patches.FancyBboxPatch((2.5, y - 5), 43, 12.5, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#061c33", edgecolor="#0284c7", linewidth=1)
        ax.add_patch(box)
        ax.text(4, y + 3.5, st, color=CLR_CYAN, fontsize=8.5, fontweight='bold', va='center')
        ax.text(4, y - 1.5, sd, color=CLR_TEXT_WHITE, fontsize=7.5, va='center')

    # Right Container: 4 Pillars of National Impact (51% width)
    right_bg = patches.FancyBboxPatch((49, 12), 50, 86, boxstyle="round,pad=0.8,rounding_size=2",
                                      facecolor="#f0f7ff", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(right_bg)
    ax.text(74, 93.5, "THE FOUR PILLARS OF NATIONAL IMPACT", color="#0d47a1", fontsize=11, fontweight='bold', ha='center', va='center')

    pillars = [
        ("👥 SOCIAL IMPACT (Saving Lives)", 
         "Eliminates fatal capsizings by calculating real-time wave-vessel physical resonance. Delivers plain-language voice alerts in Tamil, Malayalam, Bengali, Telugu, and Hindi for non-literate crews.",
         CLR_EMERALD),
        ("💰 ECONOMIC IMPACT (Fuel & Yield)", 
         "Reduces diesel expenditure by 20–30% through current-assisted dynamic A* pathfinding. Maximizes high-value pelagic catch yields via multi-species Habitat Suitability Index (HSI) mapping.",
         CLR_AMBER),
        ("🌿 ENVIRONMENTAL IMPACT (Ocean Health)", 
         "Substantially lowers maritime diesel emissions per fishing voyage. Provides automated geofencing to prevent accidental intrusion into Marine Protected Areas (MPAs) and international boundaries.",
         CLR_BLUE),
        ("🇮🇳 STRATEGIC IMPACT (Digital Sovereignty)", 
         "Delivers end-to-end sovereign maritime domain awareness across India's 7,516 km coastline and 2.37M sq km EEZ, creating an integrated national ocean safety infrastructure.",
         "#0d47a1")
    ]
    for idx, (p_name, p_body, p_col) in enumerate(pillars):
        y = 82 - idx * 17.5
        box = patches.FancyBboxPatch((50.5, y - 6.5), 47, 15, boxstyle="round,pad=0.4,rounding_size=1.5",
                                     facecolor="#ffffff", edgecolor=p_col, linewidth=1.5)
        ax.add_patch(box)
        ax.text(52, y + 4.5, p_name, color=p_col, fontsize=9, fontweight='bold', va='center')
        ax.text(52, y - 2, p_body, color=CLR_TEXT_DARK, fontsize=7.5, va='center')

    # Bottom Progression Banner: National Scale
    btm_box = patches.FancyBboxPatch((1, 1.5), 98, 8.5, boxstyle="round,pad=0.5,rounding_size=2",
                                     facecolor="#0a192f", edgecolor=CLR_CYAN, linewidth=1.5)
    ax.add_patch(btm_box)
    ax.text(50, 5.7, "1 ARTISANAL VESSEL  ➔  HARBOR FLEET (500+)  ➔  COASTAL STATE  ➔  PAN-INDIA MARITIME DIGITAL TWIN LAYER",
            color=CLR_CYAN, fontsize=9.5, fontweight='bold', ha='center', va='center')

    img_path = os.path.join(ASSETS_DIR, "slide5_impact.png")
    plt.tight_layout()
    plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight', pad_inches=0.05)
    plt.close()
    return img_path

def create_slide6_graphic():
    fig, ax = plt.subplots(figsize=(15.5, 6.8), dpi=300)
    fig.patch.set_facecolor("#ffffff")
    ax.set_facecolor("#ffffff")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis('off')

    # 4 Core Research Cards (2x2 Grid)
    research_cards = [
        ("IMO Code on Intact Stability (2008)",
         "Resolution MSC.267(85) / Small Craft Capsizing Dynamics",
         "Small vessel dynamic instability under steep wave action, beam-sea broaching, and parametric roll resonance.",
         "Deterministic Safety Engine: Computes critical wave threshold (Hcrit ≈ 0.6 · Lhull) and enforces non-bypassable hazard overrides (<10ms).",
         CLR_CYAN),
        ("Marine Hydrodynamics & Encounter Physics",
         "O.M. Faltinsen (1990) & J.N. Newman (1977)",
         "Wave steepness ratio (S = Hs / λ), deep-water dispersion relations (λ = g·Tp² / 2π), and encounter angle frequency shifts.",
         "Vessel Digital Twin: Dynamic calculation of wave steepness & vessel-wave interaction rather than treating sea states as static point forecasts.",
         CLR_EMERALD),
        ("Satellite Oceanography & Pelagic Habitats",
         "INCOIS PFZ Mission & ISRO MOSDAC Telemetry",
         "Aggregation of pelagic marine species along thermal gradients, SST frontal boundaries, and chlorophyll-a density concentration zones.",
         "Multi-Species Habitat Index (HSI): XGBoost model integrated with 2D Sobel spatial filter for high-precision fishery zone recommendations.",
         CLR_AMBER),
        ("Dynamic Constrained Pathfinding & Graph Optimization",
         "Modified A* over Time-Varying Navigational Cost Surfaces",
         "Optimization of multiobjective vessel routes under dynamic weather constraints, geofenced boundaries, and risk cost penalization.",
         "ORCA Navigational Pathfinder: Computes 4-point optimized waypoints balancing weather safety, fuel consumption, and dock proximity.",
         "#0284c7")
    ]

    for idx, (title, src, insight, app, bcol) in enumerate(research_cards):
        col = idx % 2
        row = idx // 2
        x = 1 + col * 49.5
        y = 54 if row == 0 else 14
        
        box = patches.FancyBboxPatch((x, y), 48.5, 38, boxstyle="round,pad=0.8,rounding_size=2",
                                     facecolor="#0a192f", edgecolor=bcol, linewidth=1.5)
        ax.add_patch(box)
        
        ax.text(x + 2, y + 32, f"📚 {title}", color=CLR_CYAN, fontsize=9.5, fontweight='bold', va='center')
        ax.text(x + 2, y + 27, f"Source: {src}", color="#94a3b8", fontsize=7.5, style='italic', va='center')
        
        ax.text(x + 2, y + 20, "Key Insight:", color=CLR_AMBER, fontsize=8, fontweight='bold', va='center')
        ax.text(x + 2, y + 14, insight, color=CLR_TEXT_WHITE, fontsize=7.5, va='center')
        
        ax.text(x + 2, y + 7.5, "ORCA Implementation:", color=CLR_EMERALD, fontsize=8, fontweight='bold', va='center')
        ax.text(x + 2, y + 2.5, app, color="#cbd5e1", fontsize=7.5, va='center')

    # Bottom Citations Container
    ref_bg = patches.FancyBboxPatch((1, 1.5), 98, 10.5, boxstyle="round,pad=0.5,rounding_size=2",
                                    facecolor="#f0f7ff", edgecolor="#0284c7", linewidth=1.5)
    ax.add_patch(ref_bg)
    ax.text(3, 8.5, "AUTHENTIC OPERATIONAL & SCIENTIFIC REFERENCES", color="#0d47a1", fontsize=8.5, fontweight='bold', va='center')

    cits = (
        "[1] International Maritime Organization (IMO). Intact Stability Code (2008). Res. MSC.267(85). | "
        "[2] INCOIS. Ocean State Forecast & Potential Fishing Zone Operational Guidelines.\n"
        "[3] ISRO Space Applications Centre. Oceansat-3 & INSAT-3DR Products, MOSDAC. | "
        "[4] Copernicus Marine (CMEMS). Global Ocean Analysis. | "
        "[5] WMO-No. 558: Marine Meteorological Services."
    )
    ax.text(3, 4.5, cits, color=CLR_TEXT_DARK, fontsize=7, va='center')

    img_path = os.path.join(ASSETS_DIR, "slide6_references.png")
    plt.tight_layout()
    plt.savefig(img_path, facecolor=fig.get_facecolor(), edgecolor='none', bbox_inches='tight', pad_inches=0.05)
    plt.close()
    return img_path

def build_presentation():
    print("Generating high-resolution diagram graphics...")
    img_s1 = create_slide1_graphic()
    img_s2 = create_slide2_graphic()
    img_s3 = create_slide3_graphic()
    img_s4 = create_slide4_graphic()
    img_s5 = create_slide5_graphic()
    img_s6 = create_slide6_graphic()

    print(f"Loading template: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)

    # 1. DELETE SLIDE 7 (Instruction Slide)
    if len(prs.slides) >= 7:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]
        print("Deleted instruction slide. Active slides: 6")

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides[0]
    print("Formatting Slide 1...")

    # Form Fields in TextBox 9
    for shape in slide1.shapes:
        if shape.name == "TextBox 9":
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            
            fields = [
                ("Problem Statement ID:", " SIH26176"),
                ("Problem Statement Title:", " AI-Assisted Maritime Decision Intelligence & Bio-Physical Safety Platform for Artisanal Fishermen & Coastal Fleets"),
                ("Theme:", " Space Technology / Blue Economy / Disaster Management"),
                ("PS Category:", " Software (with Edge Hardware compatibility)"),
                ("Team ID:", " [Team ID]"),
                ("Team Name:", " ORCA 4.0")
            ]
            
            for idx, (label, val) in enumerate(fields):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.space_after = Pt(8)
                
                r1 = p.add_run()
                r1.text = label
                r1.font.bold = True
                r1.font.size = Pt(13)
                r1.font.color.rgb = RGBColor(13, 71, 161)
                
                r2 = p.add_run()
                r2.text = val
                r2.font.bold = False
                r2.font.size = Pt(13)
                r2.font.color.rgb = RGBColor(15, 23, 42)

    # Insert Slide 1 Hero Image
    slide1.shapes.add_picture(img_s1, Inches(7.0), Inches(1.3), Inches(5.9), Inches(5.6))

    # =========================================================================
    # SLIDES 2 TO 6 CONFIGURATION
    # =========================================================================
    slide_configs = [
        (1, "IDEA TITLE: ORCA 4.0 — Maritime Decision Intelligence",
         [("PROPOSED SOLUTION", RGBColor(6, 182, 212)),
          ("DETAILED EXPLANATION", RGBColor(13, 71, 161)),
          ("HOW IT ADDRESSES PROBLEM", RGBColor(13, 71, 161)),
          ("INNOVATION & UNIQUENESS", RGBColor(16, 185, 129))],
         img_s2),
        (2, "TECHNICAL APPROACH: End-to-End Architecture",
         [("TECHNOLOGIES TO BE USED", RGBColor(6, 182, 212)),
          ("METHODOLOGY & PROCESS", RGBColor(13, 71, 161)),
          ("SYSTEM FLOWCHARTS", RGBColor(13, 71, 161)),
          ("PROTOTYPE SIMULATION", RGBColor(16, 185, 129))],
         img_s3),
        (3, "FEASIBILITY AND VIABILITY: Build, Mitigate & Deploy",
         [("FEASIBILITY ANALYSIS", RGBColor(6, 182, 212)),
          ("POTENTIAL RISKS & CHALLENGES", RGBColor(245, 158, 11)),
          ("STRATEGIES FOR OVERCOMING", RGBColor(16, 185, 129)),
          ("DEPLOYMENT & ECONOMICS", RGBColor(13, 71, 161))],
         img_s4),
        (4, "IMPACT AND BENEFITS: Empowering India's Blue Economy",
         [("TARGET AUDIENCE IMPACT", RGBColor(6, 182, 212)),
          ("SOCIAL & HUMAN BENEFITS", RGBColor(16, 185, 129)),
          ("ECONOMIC & OPERATIONAL ROI", RGBColor(245, 158, 11)),
          ("NATIONAL STRATEGIC SCALE", RGBColor(13, 71, 161))],
         img_s5),
        (5, "RESEARCH  AND REFERENCES: Science to Implementation",
         [("DETAILS OF RESEARCH WORK", RGBColor(6, 182, 212)),
          ("IMO STABILITY CRITERIA", RGBColor(13, 71, 161)),
          ("HYDRODYNAMIC WAVE PHYSICS", RGBColor(13, 71, 161)),
          ("VERIFIED ACADEMIC CITATIONS", RGBColor(16, 185, 129))],
         img_s6)
    ]

    for slide_idx, title_text, badges, img_path in slide_configs:
        slide = prs.slides[slide_idx]
        print(f"Formatting Slide {slide_idx + 1}: {title_text[:30]}...")

        # Update and clean placeholders
        for s in slide.shapes:
            if s.name == "Title 1":
                # Position title cleanly so it DOES NOT collide with Oval on left or Logo on right
                s.left = Inches(1.85)
                s.top = Inches(0.15)
                s.width = Inches(8.8)
                s.height = Inches(0.85)
                tf = s.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.name = "Arial"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(13, 71, 161)
            elif "Oval" in s.name:
                s.text_frame.text = "ORCA 4.0"
                if len(s.text_frame.paragraphs) > 0:
                    s.text_frame.paragraphs[0].font.size = Pt(11)
                    s.text_frame.paragraphs[0].font.bold = True
            elif "TextBox" in s.name:
                s.text_frame.clear()

        # Add Official Idea Pointer Badges Across Top
        for b_idx, (b_text, b_col) in enumerate(badges):
            badge_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                Inches(0.6 + b_idx * 3.05), Inches(1.05),
                                                Inches(2.95), Inches(0.32))
            badge_shape.fill.solid()
            badge_shape.fill.fore_color.rgb = b_col
            badge_shape.line.fill.background()
            
            tf_b = badge_shape.text_frame
            tf_b.word_wrap = True
            tf_b.margin_left = Inches(0.05)
            tf_b.margin_right = Inches(0.05)
            tf_b.margin_top = Inches(0.02)
            tf_b.margin_bottom = Inches(0.02)
            p_b = tf_b.paragraphs[0]
            p_b.text = b_text
            p_b.font.name = "Arial"
            p_b.font.size = Pt(9)
            p_b.font.bold = True
            p_b.font.color.rgb = RGBColor(255, 255, 255)
            p_b.alignment = PP_ALIGN.CENTER

        # Add High-Resolution Diagram Graphic into Content Area
        slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.45), Inches(12.13), Inches(5.35))

    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Elite SIH 2026 Presentation generated successfully!")
    print(f"Saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    build_presentation()
